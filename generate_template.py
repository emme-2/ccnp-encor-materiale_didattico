#!/usr/bin/env python3
"""
generate_template.py — Master Template Generator
CCNP ENCOR 350-401 · Materiale Didattico
Genera TEMPLATE/master_template.pptx con 10 layout slide.

Uso:  python generate_template.py
Deps: python-pptx >= 1.0, lxml
"""

import os
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree

# ─────────────────────────────────────────────────────────────────────────────
# Palette
# ─────────────────────────────────────────────────────────────────────────────
ARANCIO = RGBColor(0xEF, 0x6A, 0x01)   # #EF6A01  primario, accent, titoli
NERO    = RGBColor(0x00, 0x00, 0x00)   # #000000  testo, struttura
AVORIO  = RGBColor(0xED, 0xEB, 0xDC)   # #EDEBDC  sfondi neutri
GIALLO  = RGBColor(0xFF, 0xB6, 0x00)   # #FFB600  highlight, exam tips
BIANCO  = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF

# Colori derivati (non nella palette ufficiale — uso interno per elementi blended)
DARK_BG   = RGBColor(0x1A, 0x1A, 0x1A)   # quasi-nero per code box
PALE_ORA  = RGBColor(0xFB, 0xE4, 0xCC)   # arancio 15% su avorio — key concept box
DARK_YEL  = RGBColor(0x4A, 0x3D, 0x00)   # giallo 30% su dark — code highlight

# ─────────────────────────────────────────────────────────────────────────────
# Dimensioni slide 16:9 — 33.87 × 19.05 cm
# ─────────────────────────────────────────────────────────────────────────────
W = Cm(33.87)
H = Cm(19.05)

# ─────────────────────────────────────────────────────────────────────────────
# Font
# ─────────────────────────────────────────────────────────────────────────────
F_SANS = "Montserrat"
F_MONO = "Courier New"

# ─────────────────────────────────────────────────────────────────────────────
# Misure ricorrenti
# ─────────────────────────────────────────────────────────────────────────────
PAD_L    = Cm(1.06)   # left padding ≈ 40px a 96dpi
HDR_60   = Cm(1.59)   # header 60px
HDR_55   = Cm(1.46)   # header 55px
STRIPE_T = Cm(0.25)   # top stripe cover (8px)
STRIPE_B = Cm(0.15)   # bottom stripe cover (4px)

# ─────────────────────────────────────────────────────────────────────────────
# Helper — slide vuota
# ─────────────────────────────────────────────────────────────────────────────

def new_slide(prs):
    """Aggiunge una slide con layout Blank."""
    blank = prs.slide_layouts[6]   # index 6 = "Blank" nel tema default
    return prs.slides.add_slide(blank)


# ─────────────────────────────────────────────────────────────────────────────
# Helpers — rettangoli
# ─────────────────────────────────────────────────────────────────────────────

def rect(slide, x, y, w, h, fill=None, border=None, border_w=Pt(1)):
    """
    Aggiunge un rettangolo.
    fill=None  → sfondo trasparente
    border=None → nessun bordo
    """
    shp = slide.shapes.add_shape(1, int(x), int(y), int(w), int(h))
    # Fill
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    # Bordo
    if border is not None:
        shp.line.color.rgb = border
        shp.line.width = int(border_w)
    else:
        _remove_border(shp)
    return shp


def _remove_border(shp):
    """Rimuove completamente il bordo da un shape via XML."""
    spPr = shp._element.spPr
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        spPr.remove(ln)
    ns = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    spPr.append(parse_xml(f'<a:ln {ns}><a:noFill/></a:ln>'))


def rect_dashed(slide, x, y, w, h, color=ARANCIO, lw=Pt(1.5)):
    """Rettangolo con bordo tratteggiato e fill trasparente."""
    shp = slide.shapes.add_shape(1, int(x), int(y), int(w), int(h))
    shp.fill.background()
    shp.line.color.rgb = color
    shp.line.width = int(lw)
    shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return shp


# ─────────────────────────────────────────────────────────────────────────────
# Helpers — testo
# ─────────────────────────────────────────────────────────────────────────────

def txt(slide, x, y, w, h, text,
        font=F_SANS, size=18, color=NERO,
        bold=False, italic=False,
        align=PP_ALIGN.LEFT, wrap=True):
    """Aggiunge un textbox con singolo paragrafo."""
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tb


def hdr(slide, title, fill=NERO, color=BIANCO, height=HDR_55,
        size=22, pad=PAD_L):
    """
    Barra header a tutta larghezza con testo verticalmente centrato.
    Ritorna lo shape.
    """
    from pptx.enum.text import MSO_ANCHOR
    shp = rect(slide, 0, 0, W, height, fill=fill)
    tf = shp.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = 0
    tf.margin_left = int(pad)
    tf.margin_right = int(pad)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = F_SANS
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = True
    return shp


def code_box(slide, x, y, w, h, lines, fill=DARK_BG):
    """
    Box di codice monospace con sfondo scuro e bordo ARANCIO.
    lines: lista di stringhe, una per riga.
    Restituisce (shape, line_h_cm) dove line_h_cm è l'altezza stimata di una riga.
    """
    from pptx.enum.text import MSO_ANCHOR
    shp = rect(slide, x, y, w, h, fill=fill, border=ARANCIO, border_w=Pt(1))
    tf = shp.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_top = int(Cm(0.3))
    tf.margin_bottom = int(Cm(0.3))
    tf.margin_left = int(Cm(0.4))
    tf.margin_right = int(Cm(0.4))
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.name = F_MONO
        run.font.size = Pt(15)
        run.font.color.rgb = BIANCO
        run.font.bold = False
    return shp


def style_cell(cell, text, bg, fg, size=15, bold=False, align=PP_ALIGN.LEFT):
    """Stile uniforme per celle di tabella."""
    from pptx.enum.text import MSO_ANCHOR
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg
    tf = cell.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = int(Pt(6))
    tf.margin_left = tf.margin_right = int(Pt(10))
    p = tf.paragraphs[0]
    p.alignment = align
    for r in p.runs:
        p._p.remove(r._r)
    run = p.add_run()
    run.text = text
    run.font.name = F_SANS
    run.font.size = Pt(size)
    run.font.color.rgb = fg
    run.font.bold = bold


# ─────────────────────────────────────────────────────────────────────────────
# Layout 01 — Cover Module
# Sfondo NERO · stripe ARANCIO top/bottom · titolo ARANCIO centrato
# ─────────────────────────────────────────────────────────────────────────────

def layout_01_cover(prs):
    slide = new_slide(prs)

    # Sfondo
    rect(slide, 0, 0, W, H, fill=NERO)

    # Stripe top
    rect(slide, 0, 0, W, STRIPE_T, fill=ARANCIO)

    # Stripe bottom
    rect(slide, 0, H - STRIPE_B, W, STRIPE_B, fill=ARANCIO)

    # Logo placeholder — top right
    txt(slide, W - Cm(5.5), Cm(0.5), Cm(5), Cm(0.9),
        "LOGO", size=14, color=AVORIO, bold=True, align=PP_ALIGN.RIGHT)

    # Linea decorativa sotto logo
    rect(slide, W - Cm(5.5), Cm(1.5), Cm(5), Cm(0.05), fill=ARANCIO)

    # Titolo modulo — centrato verticalmente
    txt(slide, Cm(2.5), Cm(6.2), W - Cm(5), Cm(3.2),
        "MOD-XX — Titolo del Modulo",
        size=40, color=ARANCIO, bold=True, align=PP_ALIGN.CENTER)

    # Separatore thin
    sep_w = Cm(15)
    rect(slide, (W - sep_w) / 2, Cm(9.8), sep_w, Cm(0.06), fill=ARANCIO)

    # Sottotitolo
    txt(slide, Cm(2.5), Cm(10.2), W - Cm(5), Cm(1.8),
        "Area Tematica  ·  Ore  ·  Codici Syllabus",
        size=18, color=AVORIO, align=PP_ALIGN.CENTER)

    return slide


# ─────────────────────────────────────────────────────────────────────────────
# Layout 02 — Agenda
# Sfondo AVORIO · header NERO · lista numerata con numeri ARANCIO
# ─────────────────────────────────────────────────────────────────────────────

def layout_02_agenda(prs):
    slide = new_slide(prs)

    rect(slide, 0, 0, W, H, fill=AVORIO)
    hdr(slide, "AGENDA", fill=NERO, color=BIANCO, height=HDR_60, size=24)

    items = [
        ("01", "Argomento uno — Introduzione e panoramica del modulo"),
        ("02", "Argomento due — Concetti fondamentali e teoria"),
        ("03", "Argomento tre — Configurazione pratica passo-passo"),
        ("04", "Argomento quattro — Verifica, output e troubleshooting"),
        ("05", "Argomento cinque — Exam Tips e riepilogo finale"),
    ]

    num_w    = Cm(2.2)
    lbl_x    = PAD_L + num_w + Cm(0.3)
    lbl_w    = W - lbl_x - Cm(1.5)
    row_h    = Cm(2.6)
    start_y  = HDR_60 + Cm(0.9)

    for i, (num, label) in enumerate(items):
        y = start_y + i * row_h

        # Numero ARANCIO Bold
        txt(slide, PAD_L, y, num_w, row_h,
            num, size=22, color=ARANCIO, bold=True, align=PP_ALIGN.LEFT)

        # Separatore verticale sottile
        rect(slide, PAD_L + num_w + Cm(0.1), y + Cm(0.3),
             Cm(0.05), row_h - Cm(0.6), fill=ARANCIO)

        # Label NERO Regular
        txt(slide, lbl_x, y + Cm(0.3), lbl_w, row_h - Cm(0.3),
            label, size=20, color=NERO, align=PP_ALIGN.LEFT)

    return slide


# ─────────────────────────────────────────────────────────────────────────────
# Layout 03 — Section Header
# Sfondo ARANCIO pieno · titolo BIANCO centrato · separatore + sottotitolo
# ─────────────────────────────────────────────────────────────────────────────

def layout_03_section(prs):
    slide = new_slide(prs)

    rect(slide, 0, 0, W, H, fill=ARANCIO)

    # Titolo — zona verticale centrale-alta
    title_y = H / 2 - Cm(4.5)
    txt(slide, Cm(2.5), title_y, W - Cm(5), Cm(3.8),
        "Titolo della Sezione",
        size=48, color=BIANCO, bold=True, align=PP_ALIGN.CENTER)

    # Separatore bianco centrato
    sep_w = Cm(12)
    sep_y = title_y + Cm(4.0)
    rect(slide, (W - sep_w) / 2, sep_y, sep_w, Cm(0.08), fill=BIANCO)

    # Sottotitolo
    txt(slide, Cm(3), sep_y + Cm(0.35), W - Cm(6), Cm(2.0),
        "Sottotitolo descrittivo — ambito e obiettivi della sezione",
        size=22, color=BIANCO, align=PP_ALIGN.CENTER)

    return slide


# ─────────────────────────────────────────────────────────────────────────────
# Layout 04 — Teoria Concetto
# Sfondo AVORIO · header NERO · bordo sinistro ARANCIO · key-concept box bottom
# ─────────────────────────────────────────────────────────────────────────────

def layout_04_teoria(prs):
    slide = new_slide(prs)

    rect(slide, 0, 0, W, H, fill=AVORIO)
    hdr(slide, "Titolo Concetto Teorico", fill=NERO, color=BIANCO,
        height=HDR_55, size=22)

    # Metriche area corpo
    box_h      = Cm(2.5)            # altezza key-concept box in fondo
    body_top   = HDR_55 + Cm(0.4)
    body_bot   = H - box_h - Cm(0.5)
    body_h     = body_bot - body_top
    accent_w   = Cm(0.35)
    body_left  = Cm(0.25) + accent_w + Cm(0.3)
    body_w     = W - body_left - Cm(1.0)

    # Bordo sinistro ARANCIO
    rect(slide, Cm(0.25), body_top, accent_w, body_h, fill=ARANCIO)

    # Corpo testo — max 6 righe
    body_text = (
        "Il testo del concetto teorico viene inserito qui.\n"
        "Massimo 5-6 righe per mantenere la leggibilità ottimale.\n\n"
        "• Punto chiave uno — definizione e contesto del concetto\n"
        "• Punto chiave due — applicazione pratica in ambiente IOS\n"
        "• Punto chiave tre — differenza rispetto all'approccio legacy\n"
        "• Punto chiave quattro — implicazioni per l'esame CCNP ENCOR"
    )
    txt(slide, body_left, body_top + Cm(0.2), body_w, body_h - Cm(0.4),
        body_text, size=18, color=NERO, wrap=True)

    # Key-concept box — fondo arancio tenue (PALE_ORA = ~15% ARANCIO su AVORIO)
    box_y = H - box_h - Cm(0.3)
    rect(slide, Cm(0.5), box_y, W - Cm(1.0), box_h, fill=PALE_ORA)
    txt(slide, Cm(1.0), box_y + Cm(0.4), W - Cm(2.0), box_h - Cm(0.4),
        "💡  Concetto chiave: sintesi del punto essenziale da ricordare per l'esame.",
        size=16, color=ARANCIO, bold=True, wrap=True)

    return slide


# ─────────────────────────────────────────────────────────────────────────────
# Layout 05 — Diagramma / Topologia
# Sfondo BIANCO · header NERO · area tratteggiata ARANCIO · didascalia bottom
# ─────────────────────────────────────────────────────────────────────────────

def layout_05_diagramma(prs):
    slide = new_slide(prs)

    rect(slide, 0, 0, W, H, fill=BIANCO)
    hdr(slide, "Diagramma / Topologia", fill=NERO, color=BIANCO,
        height=HDR_55, size=22)

    caption_h = Cm(0.9)
    margin    = Cm(0.8)
    diag_top  = HDR_55 + Cm(0.4)
    diag_h    = H - diag_top - caption_h - Cm(0.6)

    # Area diagramma — bordo tratteggiato ARANCIO
    rect_dashed(slide, margin, diag_top, W - 2 * margin, diag_h)

    # Placeholder testuale centrato nell'area
    txt(slide,
        margin + Cm(1), diag_top + diag_h / 2 - Cm(0.7),
        W - 2 * margin - Cm(2), Cm(1.4),
        "[ Diagramma / Topologia ]",
        size=16, color=ARANCIO, align=PP_ALIGN.CENTER)

    # Didascalia
    cap_y = H - caption_h - Cm(0.2)
    rect(slide, 0, cap_y - Cm(0.1), W, Cm(0.05), fill=AVORIO)
    txt(slide, PAD_L, cap_y, W - 2 * PAD_L, caption_h,
        "Figura X — Descrizione della topologia o del diagramma",
        size=14, color=NERO, italic=True)

    return slide


# ─────────────────────────────────────────────────────────────────────────────
# Layout 06 — Config / Comando
# Sfondo NERO · header ARANCIO · code box con highlight riga · badge device
# ─────────────────────────────────────────────────────────────────────────────

def layout_06_config(prs):
    slide = new_slide(prs)

    rect(slide, 0, 0, W, H, fill=NERO)
    hdr(slide, "Configurazione / Comando",
        fill=ARANCIO, color=NERO, height=HDR_55, size=22)

    m        = Cm(1.0)
    code_top = HDR_55 + Cm(0.55)
    code_h   = H - code_top - Cm(0.6)
    code_w   = W - 2 * m

    # Righe di codice — riga 4 (index 3) sarà evidenziata
    lines = [
        "Router# show ip ospf neighbor",
        "",
        "Neighbor ID    Pri   State       Dead    Address        Interface",
        "2.2.2.2          1   FULL/DR   00:00:38  10.0.12.2  Ethernet0/0.12",
        "3.3.3.3          1   FULL/BDR  00:00:31  10.0.13.3  Ethernet0/1.13",
        "",
        "Router# show ip route ospf",
        "O    10.3.3.3/32 [110/2] via 10.0.13.3, 00:12:41, Ethernet0/1.13",
    ]
    code_box(slide, m, code_top, code_w, code_h, lines)

    # Highlight DARK_YEL su riga 4 (index 3, 0-based)
    # Stima: font 15pt → ~0.63cm/riga; margin_top 0.3cm
    line_pitch = Cm(0.635)
    margin_top = Cm(0.3)
    hl_idx     = 3
    hl_y       = code_top + margin_top + hl_idx * line_pitch
    rect(slide,
         m + Cm(0.1), hl_y,
         code_w - Cm(0.2), line_pitch,
         fill=DARK_YEL)

    # Badge device ARANCIO — top-right del code box
    badge_w = Cm(1.8)
    badge_h = Cm(0.65)
    badge_x = m + code_w - badge_w - Cm(0.15)
    badge_y = code_top + Cm(0.15)
    rect(slide, badge_x, badge_y, badge_w, badge_h, fill=ARANCIO)
    txt(slide, badge_x, badge_y, badge_w, badge_h,
        "R1", size=12, color=NERO, bold=True, align=PP_ALIGN.CENTER)

    return slide


# ─────────────────────────────────────────────────────────────────────────────
# Layout 07 — Output Verifica
# Sfondo NERO · header ARANCIO · label GIALLO · output con riga ARANCIO + freccia
# ─────────────────────────────────────────────────────────────────────────────

def layout_07_verifica(prs):
    slide = new_slide(prs)

    rect(slide, 0, 0, W, H, fill=NERO)
    hdr(slide, "Output di Verifica",
        fill=ARANCIO, color=NERO, height=HDR_55, size=22)

    # Label "VERIFICA ATTESA"
    lbl_y = HDR_55 + Cm(0.3)
    txt(slide, PAD_L, lbl_y, Cm(14), Cm(0.7),
        "▼  VERIFICA ATTESA",
        size=14, color=GIALLO, bold=True)

    # Code box
    m        = Cm(1.0)
    code_top = lbl_y + Cm(0.8)
    code_h   = H - code_top - Cm(0.6)
    code_w   = W - 2 * m

    lines = [
        "HUB# show crypto isakmp sa",
        "",
        "IPv4 Crypto ISAKMP SA",
        "dst             src             state     conn-id  status",
        "198.51.100.254  192.0.2.254     QM_IDLE   1001     ACTIVE",
        "203.0.113.254   192.0.2.254     QM_IDLE   1002     ACTIVE",
        "",
        "HUB# show crypto ikev2 sa",
        "",
        "Tunnel-id  Local              Remote             Status",
        "1          192.0.2.254/4500   198.51.100.254/4500  READY",
        "2          192.0.2.254/4500   203.0.113.254/4500   READY",
    ]
    code_box(slide, m, code_top, code_w, code_h, lines)

    # Riga attesa (index 3 = header tabella isakmp) evidenziata in ARANCIO
    line_pitch = Cm(0.635)
    margin_top = Cm(0.3)
    hl_idx     = 3
    hl_y       = code_top + margin_top + hl_idx * line_pitch

    # Overlay testuale ARANCIO sulla riga evidenziata
    txt(slide,
        m + Cm(0.4), hl_y,
        code_w - Cm(6), line_pitch + Cm(0.1),
        "dst             src             state     conn-id  status",
        font=F_MONO, size=15, color=ARANCIO, bold=False)

    # Freccia "← atteso"
    txt(slide,
        W - Cm(6.5), hl_y,
        Cm(6), line_pitch + Cm(0.1),
        "← atteso",
        size=14, color=ARANCIO, bold=True, align=PP_ALIGN.RIGHT)

    return slide


# ─────────────────────────────────────────────────────────────────────────────
# Layout 08 — Troubleshooting
# Sfondo AVORIO · header NERO · tabella 2 col: SINTOMO | CAUSA+FIX
# ─────────────────────────────────────────────────────────────────────────────

def layout_08_trouble(prs):
    slide = new_slide(prs)

    rect(slide, 0, 0, W, H, fill=AVORIO)
    hdr(slide, "Troubleshooting Guide",
        fill=NERO, color=BIANCO, height=HDR_55, size=22)

    # Tabella
    t_m    = Cm(0.8)
    t_top  = HDR_55 + Cm(0.45)
    t_w    = W - 2 * t_m
    t_h    = H - t_top - Cm(0.4)
    rows   = 4   # 1 header + 3 data
    cols   = 2

    tbl_shp = slide.shapes.add_table(rows, cols, int(t_m), int(t_top),
                                      int(t_w), int(t_h))
    tbl = tbl_shp.table
    tbl.columns[0].width = int(t_w * 0.36)
    tbl.columns[1].width = int(t_w * 0.64)

    # Header
    style_cell(tbl.cell(0, 0), "SINTOMO",
               NERO, BIANCO, size=16, bold=True, align=PP_ALIGN.CENTER)
    style_cell(tbl.cell(0, 1), "CAUSA + FIX",
               ARANCIO, NERO, size=16, bold=True, align=PP_ALIGN.CENTER)

    # Dati
    data = [
        ("show crypto ikev2 sa vuoto dopo ping",
         "Proposal non corrispondenti — verificare encryption / integrity / group su entrambi i peer."),
        ("Phase 1 up (QM_IDLE) ma pkts encaps = 0",
         "Transform-set non corrispondente (Phase 2 fallisce) — show crypto ipsec transform-set."),
        ("Tunnel UP ma traffico VRF non raggiunge destinazione",
         "Route statiche VRF mancanti — verificare ip route vrf CUST-A ... TunnelXXX."),
    ]
    row_bgs = [BIANCO, AVORIO, BIANCO]

    for i, (symp, fix) in enumerate(data):
        bg = row_bgs[i]
        style_cell(tbl.cell(i + 1, 0), symp, bg, NERO, size=14)
        style_cell(tbl.cell(i + 1, 1), fix,  bg, NERO, size=14)

    return slide


# ─────────────────────────────────────────────────────────────────────────────
# Layout 09 — Exam Tips
# Sfondo GIALLO · header NERO · icona · 4 bullet · box Q&A
# ─────────────────────────────────────────────────────────────────────────────

def layout_09_exam(prs):
    slide = new_slide(prs)

    rect(slide, 0, 0, W, H, fill=GIALLO)
    hdr(slide, "Exam Tips",
        fill=NERO, color=BIANCO, height=HDR_55, size=22)

    # Icona 📋
    icon_y = HDR_55 + Cm(0.4)
    txt(slide, PAD_L, icon_y, Cm(2.4), Cm(2.4),
        "📋", size=36, color=NERO, align=PP_ALIGN.CENTER)

    # Bullet tips (a destra dell'icona)
    tips = [
        "▶  ESP usa protocollo IP 50 — non TCP/UDP. Ricordare SPI (identifica SA) e ICV (integrità).",
        "▶  IKEv2: IKE_SA_INIT + IKE_AUTH → Phase 1 (QM_IDLE). CREATE_CHILD_SA → Phase 2.",
        "▶  tunnel protection ipsec profile: obbligatorio per DMVPN — crypto map non è compatibile con mGRE.",
        "▶  no split-horizon + no next-hop-self su HUB DMVPN: entrambi critici per Phase 2/3.",
    ]
    tip_x   = PAD_L + Cm(2.8)
    tip_w   = W - tip_x - Cm(1.2)
    tip_h   = Cm(1.8)
    tip_y0  = HDR_55 + Cm(0.45)

    for i, tip in enumerate(tips):
        txt(slide, tip_x, tip_y0 + i * tip_h, tip_w, tip_h,
            tip, size=17, color=NERO, wrap=True)

    # Box Q&A
    qa_top = tip_y0 + len(tips) * tip_h + Cm(0.6)
    qa_h   = H - qa_top - Cm(0.5)
    rect(slide, PAD_L, qa_top, W - 2 * PAD_L, qa_h,
         fill=BIANCO, border=NERO, border_w=Pt(1))

    qa_text = (
        "Q:  Qual è il vantaggio principale di tunnel protection ipsec profile rispetto a crypto map?\n\n"
        "A:  Compatibile con mGRE (DMVPN) e non richiede ACL per selezionare il traffico. "
        "Crypto map è legacy e non scalabile in topologie hub-and-spoke dinamiche."
    )
    txt(slide,
        PAD_L + Cm(0.4), qa_top + Cm(0.3),
        W - 2 * PAD_L - Cm(0.8), qa_h - Cm(0.5),
        qa_text, size=16, color=NERO, wrap=True)

    return slide


# ─────────────────────────────────────────────────────────────────────────────
# Layout 10 — Summary / Takeaway
# Sfondo ARANCIO · titolo BIANCO · 3 box NERO · footer
# ─────────────────────────────────────────────────────────────────────────────

def layout_10_summary(prs):
    slide = new_slide(prs)

    rect(slide, 0, 0, W, H, fill=ARANCIO)

    # Titolo
    txt(slide, Cm(1), Cm(1.0), W - Cm(2), Cm(2.2),
        "TAKEAWAY",
        size=36, color=BIANCO, bold=True, align=PP_ALIGN.CENTER)

    # Separatore
    sep_w = Cm(20)
    rect(slide, (W - sep_w) / 2, Cm(3.1), sep_w, Cm(0.06), fill=BIANCO)

    # 3 box concetti
    n      = 3
    mx     = Cm(1.3)
    gap    = Cm(0.8)
    bw     = (W - 2 * mx - (n - 1) * gap) / n
    bt     = Cm(3.5)
    bh     = H - bt - Cm(1.8)

    labels  = ["Concetto 1", "Concetto 2", "Concetto 3"]
    numbers = ["1", "2", "3"]
    bodies  = [
        "Primo punto chiave del modulo — concetto fondamentale.",
        "Secondo punto chiave del modulo — applicazione pratica.",
        "Terzo punto chiave del modulo — implicazione per l'esame.",
    ]

    for i in range(n):
        bx = mx + i * (bw + gap)

        # Box NERO
        rect(slide, bx, bt, bw, bh, fill=NERO)

        # Numero ARANCIO
        txt(slide, bx, bt + Cm(0.5), bw, Cm(1.8),
            numbers[i], size=28, color=ARANCIO, bold=True,
            align=PP_ALIGN.CENTER)

        # Linea separatrice interna
        rect(slide, bx + Cm(1.0), bt + Cm(2.2), bw - Cm(2.0), Cm(0.06),
             fill=ARANCIO)

        # Label
        txt(slide, bx + Cm(0.3), bt + Cm(2.5), bw - Cm(0.6), Cm(1.3),
            labels[i], size=18, color=BIANCO, bold=True,
            align=PP_ALIGN.CENTER)

        # Body
        txt(slide, bx + Cm(0.4), bt + Cm(3.9), bw - Cm(0.8), bh - Cm(4.2),
            bodies[i], size=15, color=AVORIO, wrap=True,
            align=PP_ALIGN.CENTER)

    # Footer
    txt(slide, Cm(0), H - Cm(1.2), W, Cm(1.0),
        "CCNP ENCOR 350-401  —  Materiale Didattico",
        size=12, color=BIANCO, align=PP_ALIGN.CENTER)

    return slide


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────

def main():
    os.makedirs("TEMPLATE", exist_ok=True)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    layouts = [
        ("01 — Cover Module",           layout_01_cover),
        ("02 — Agenda",                  layout_02_agenda),
        ("03 — Section Header",          layout_03_section),
        ("04 — Teoria Concetto",         layout_04_teoria),
        ("05 — Diagramma / Topologia",   layout_05_diagramma),
        ("06 — Config / Comando",        layout_06_config),
        ("07 — Output Verifica",         layout_07_verifica),
        ("08 — Troubleshooting",         layout_08_trouble),
        ("09 — Exam Tips",               layout_09_exam),
        ("10 — Summary / Takeaway",      layout_10_summary),
    ]

    print("Generazione master_template.pptx...")
    for name, fn in layouts:
        fn(prs)
        print(f"  [OK] Layout {name}")

    out = "TEMPLATE/master_template.pptx"
    prs.save(out)
    size_kb = os.path.getsize(out) // 1024
    print(f"\n[DONE] Salvato: {out}  ({size_kb} KB)")
    print(f"       Slide generate: {len(prs.slides)}")
    print(f"       Dimensioni: {prs.slide_width.cm:.2f} x {prs.slide_height.cm:.2f} cm (16:9)")


if __name__ == "__main__":
    main()
