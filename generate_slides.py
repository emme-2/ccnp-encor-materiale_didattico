#!/usr/bin/env python3
"""
generate_slides.py — Genera i deck .pptx per tutti i 35 moduli CCNP ENCOR 350-401.
Uso: python generate_slides.py [MOD-XX]  oppure senza argomenti per tutti.
Deps: python-pptx >= 1.0, lxml
"""

import os
import sys
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# ─── Palette ──────────────────────────────────────────────────────────────────
ARANCIO  = RGBColor(0xEF, 0x6A, 0x01)
NERO     = RGBColor(0x00, 0x00, 0x00)
AVORIO   = RGBColor(0xED, 0xEB, 0xDC)
GIALLO   = RGBColor(0xFF, 0xB6, 0x00)
BIANCO   = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG  = RGBColor(0x1A, 0x1A, 0x1A)
PALE_ORA = RGBColor(0xFB, 0xE4, 0xCC)
DARK_YEL = RGBColor(0x4A, 0x3D, 0x00)

# ─── Dimensioni ───────────────────────────────────────────────────────────────
W = Cm(33.87); H = Cm(19.05)
F_SANS = "Montserrat"; F_MONO = "Courier New"
PAD_L = Cm(1.06); HDR_55 = Cm(1.46); HDR_60 = Cm(1.59)

# ─── Helpers ──────────────────────────────────────────────────────────────────

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _remove_border(shp):
    spPr = shp._element.spPr
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        spPr.remove(ln)
    ns = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    spPr.append(parse_xml(f'<a:ln {ns}><a:noFill/></a:ln>'))

def rect(slide, x, y, w, h, fill=None, border=None, border_w=Pt(1)):
    shp = slide.shapes.add_shape(1, int(x), int(y), int(w), int(h))
    if fill is not None:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if border is not None:
        shp.line.color.rgb = border; shp.line.width = int(border_w)
    else:
        _remove_border(shp)
    return shp

def rect_dashed(slide, x, y, w, h, color=ARANCIO, lw=Pt(1.5)):
    shp = slide.shapes.add_shape(1, int(x), int(y), int(w), int(h))
    shp.fill.background()
    shp.line.color.rgb = color; shp.line.width = int(lw)
    shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return shp

def txt(slide, x, y, w, h, text, font=F_SANS, size=18, color=NERO,
        bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = font; run.font.size = Pt(size)
    run.font.color.rgb = color; run.font.bold = bold; run.font.italic = italic
    return tb

def hdr(slide, title, fill=NERO, color=BIANCO, height=HDR_55, size=22, pad=PAD_L):
    shp = rect(slide, 0, 0, W, height, fill=fill)
    tf = shp.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = 0
    tf.margin_left = int(pad); tf.margin_right = int(pad)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = title
    run.font.name = F_SANS; run.font.size = Pt(size)
    run.font.color.rgb = color; run.font.bold = True
    return shp

def code_box(slide, x, y, w, h, lines):
    shp = rect(slide, x, y, w, h, fill=DARK_BG, border=ARANCIO, border_w=Pt(1))
    tf = shp.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_top = int(Cm(0.3)); tf.margin_bottom = int(Cm(0.3))
    tf.margin_left = int(Cm(0.4)); tf.margin_right = int(Cm(0.4))
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = line
        run.font.name = F_MONO; run.font.size = Pt(13)
        run.font.color.rgb = BIANCO
    return shp

def style_cell(cell, text, bg, fg, size=13, bold=False, align=PP_ALIGN.LEFT):
    cell.fill.solid(); cell.fill.fore_color.rgb = bg
    tf = cell.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = int(Pt(5))
    tf.margin_left = tf.margin_right = int(Pt(8))
    p = tf.paragraphs[0]; p.alignment = align
    for r in p.runs: p._p.remove(r._r)
    run = p.add_run(); run.text = text
    run.font.name = F_SANS; run.font.size = Pt(size)
    run.font.color.rgb = fg; run.font.bold = bold

# ─── Layout builders ──────────────────────────────────────────────────────────

def slide_cover(prs, mod_id, title, area, hours, codes):
    slide = new_slide(prs)
    rect(slide, 0, 0, W, H, fill=NERO)
    rect(slide, 0, 0, W, Cm(0.25), fill=ARANCIO)
    rect(slide, 0, H - Cm(0.15), W, Cm(0.15), fill=ARANCIO)
    txt(slide, W - Cm(5.5), Cm(0.5), Cm(5), Cm(0.9), "MAGNETICO",
        size=12, color=AVORIO, bold=True, align=PP_ALIGN.RIGHT)
    txt(slide, Cm(2.5), Cm(5.5), W - Cm(5), Cm(3.8),
        f"{mod_id} — {title}",
        size=38, color=ARANCIO, bold=True, align=PP_ALIGN.CENTER)
    sep_w = Cm(15)
    rect(slide, (W - sep_w) / 2, Cm(9.7), sep_w, Cm(0.06), fill=ARANCIO)
    txt(slide, Cm(2.5), Cm(10.1), W - Cm(5), Cm(1.8),
        f"{area}  ·  {hours}  ·  {codes}",
        size=17, color=AVORIO, align=PP_ALIGN.CENTER)
    txt(slide, Cm(0), H - Cm(1.5), W, Cm(1.0),
        "CCNP ENCOR 350-401  —  Materiale Didattico",
        size=11, color=AVORIO, align=PP_ALIGN.CENTER)

def slide_agenda(prs, items):
    slide = new_slide(prs)
    rect(slide, 0, 0, W, H, fill=AVORIO)
    hdr(slide, "AGENDA", fill=NERO, color=BIANCO, height=HDR_60, size=24)
    num_w = Cm(2.2); lbl_x = PAD_L + num_w + Cm(0.3)
    lbl_w = W - lbl_x - Cm(1.5); row_h = Cm(2.4)
    start_y = HDR_60 + Cm(0.8)
    for i, label in enumerate(items[:6]):
        y = start_y + i * row_h
        txt(slide, PAD_L, y, num_w, row_h,
            f"{i+1:02d}", size=22, color=ARANCIO, bold=True)
        rect(slide, PAD_L + num_w + Cm(0.1), y + Cm(0.3),
             Cm(0.05), row_h - Cm(0.6), fill=ARANCIO)
        txt(slide, lbl_x, y + Cm(0.2), lbl_w, row_h - Cm(0.2),
            label, size=19, color=NERO)

def slide_section(prs, title, subtitle=""):
    slide = new_slide(prs)
    rect(slide, 0, 0, W, H, fill=ARANCIO)
    title_y = H / 2 - Cm(4.0)
    txt(slide, Cm(2.5), title_y, W - Cm(5), Cm(3.5),
        title, size=44, color=BIANCO, bold=True, align=PP_ALIGN.CENTER)
    sep_w = Cm(12)
    sep_y = title_y + Cm(3.8)
    rect(slide, (W - sep_w) / 2, sep_y, sep_w, Cm(0.08), fill=BIANCO)
    if subtitle:
        txt(slide, Cm(3), sep_y + Cm(0.35), W - Cm(6), Cm(2.0),
            subtitle, size=20, color=BIANCO, align=PP_ALIGN.CENTER)

def slide_teoria(prs, title, bullet_points, key_concept=""):
    """bullet_points: lista di stringhe (max 6)"""
    slide = new_slide(prs)
    rect(slide, 0, 0, W, H, fill=AVORIO)
    hdr(slide, title, fill=NERO, color=BIANCO, height=HDR_55, size=20)
    box_h = Cm(2.2)
    body_top = HDR_55 + Cm(0.35)
    body_bot = H - box_h - Cm(0.45)
    body_h = body_bot - body_top
    accent_w = Cm(0.35)
    body_left = Cm(0.25) + accent_w + Cm(0.3)
    body_w = W - body_left - Cm(1.0)
    rect(slide, Cm(0.25), body_top, accent_w, body_h, fill=ARANCIO)
    body_text = "\n".join(f"• {p}" if not p.startswith("•") else p
                          for p in bullet_points)
    txt(slide, body_left, body_top + Cm(0.15), body_w, body_h - Cm(0.3),
        body_text, size=17, color=NERO, wrap=True)
    box_y = H - box_h - Cm(0.25)
    rect(slide, Cm(0.5), box_y, W - Cm(1.0), box_h, fill=PALE_ORA)
    kc = key_concept or bullet_points[-1] if bullet_points else ""
    txt(slide, Cm(1.0), box_y + Cm(0.3), W - Cm(2.0), box_h - Cm(0.35),
        f"Concetto chiave:  {kc}", size=15, color=ARANCIO, bold=True, wrap=True)

def slide_diagramma(prs, title, caption, nodes, links=None):
    """Disegna un diagramma semplice con rettangoli e testo."""
    slide = new_slide(prs)
    rect(slide, 0, 0, W, H, fill=BIANCO)
    hdr(slide, title, fill=NERO, color=BIANCO, height=HDR_55, size=20)
    caption_h = Cm(0.85)
    margin = Cm(0.7)
    diag_top = HDR_55 + Cm(0.35)
    diag_h = H - diag_top - caption_h - Cm(0.5)
    # Disegna area diagramma con bordo tratteggiato
    rect_dashed(slide, margin, diag_top, W - 2 * margin, diag_h)
    # Disegna nodi: nodes = [(label, x_frac, y_frac), ...]
    for label, xf, yf in nodes:
        nx = margin + (W - 2 * margin) * xf - Cm(2.0)
        ny = diag_top + diag_h * yf - Cm(0.5)
        rect(slide, nx, ny, Cm(4.0), Cm(1.0), fill=NERO)
        txt(slide, nx, ny, Cm(4.0), Cm(1.0),
            label, size=11, color=BIANCO, bold=True, align=PP_ALIGN.CENTER)
    cap_y = H - caption_h - Cm(0.15)
    txt(slide, PAD_L, cap_y, W - 2 * PAD_L, caption_h,
        caption, size=13, color=NERO, italic=True)

def slide_config(prs, title, lines, device="R1", highlight_idx=None):
    slide = new_slide(prs)
    rect(slide, 0, 0, W, H, fill=NERO)
    hdr(slide, title, fill=ARANCIO, color=NERO, height=HDR_55, size=20)
    m = Cm(1.0); code_top = HDR_55 + Cm(0.5)
    code_h = H - code_top - Cm(0.5); code_w = W - 2 * m
    code_box(slide, m, code_top, code_w, code_h, lines)
    if highlight_idx is not None:
        lp = Cm(0.6); mt = Cm(0.3)
        hl_y = code_top + mt + highlight_idx * lp
        rect(slide, m + Cm(0.1), hl_y, code_w - Cm(0.2), lp, fill=DARK_YEL)
    bw = Cm(2.0); bh = Cm(0.6)
    bx = m + code_w - bw - Cm(0.15); by = code_top + Cm(0.15)
    rect(slide, bx, by, bw, bh, fill=ARANCIO)
    txt(slide, bx, by, bw, bh, device, size=11, color=NERO, bold=True,
        align=PP_ALIGN.CENTER)

def slide_verifica(prs, title, lines, atteso_idx=None):
    slide = new_slide(prs)
    rect(slide, 0, 0, W, H, fill=NERO)
    hdr(slide, title, fill=ARANCIO, color=NERO, height=HDR_55, size=20)
    lbl_y = HDR_55 + Cm(0.25)
    txt(slide, PAD_L, lbl_y, Cm(14), Cm(0.65),
        "▼  VERIFICA ATTESA", size=13, color=GIALLO, bold=True)
    m = Cm(1.0); code_top = lbl_y + Cm(0.75)
    code_h = H - code_top - Cm(0.45); code_w = W - 2 * m
    code_box(slide, m, code_top, code_w, code_h, lines)
    if atteso_idx is not None:
        lp = Cm(0.6); mt = Cm(0.3)
        hl_y = code_top + mt + atteso_idx * lp
        txt(slide, W - Cm(6.5), hl_y, Cm(6), lp + Cm(0.1),
            "← atteso", size=13, color=ARANCIO, bold=True, align=PP_ALIGN.RIGHT)

def slide_troubleshooting(prs, rows):
    """rows: [(sintomo, causa_fix), ...]  max 4 righe"""
    slide = new_slide(prs)
    rect(slide, 0, 0, W, H, fill=AVORIO)
    hdr(slide, "Troubleshooting Guide",
        fill=NERO, color=BIANCO, height=HDR_55, size=20)
    t_m = Cm(0.7); t_top = HDR_55 + Cm(0.4)
    t_w = W - 2 * t_m; t_h = H - t_top - Cm(0.35)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, 2,
                                  int(t_m), int(t_top),
                                  int(t_w), int(t_h)).table
    tbl.columns[0].width = int(t_w * 0.35)
    tbl.columns[1].width = int(t_w * 0.65)
    style_cell(tbl.cell(0, 0), "SINTOMO", NERO, BIANCO, size=14, bold=True,
               align=PP_ALIGN.CENTER)
    style_cell(tbl.cell(0, 1), "CAUSA + FIX", ARANCIO, NERO, size=14, bold=True,
               align=PP_ALIGN.CENTER)
    bgs = [BIANCO, AVORIO, BIANCO, AVORIO]
    for i, (s, f) in enumerate(rows):
        bg = bgs[i % 4]
        style_cell(tbl.cell(i + 1, 0), s, bg, NERO, size=13)
        style_cell(tbl.cell(i + 1, 1), f, bg, NERO, size=13)

def slide_exam_tips(prs, tips, q_a=None):
    """tips: lista 3-5 stringhe; q_a: [(domanda, risposta), ...]"""
    slide = new_slide(prs)
    rect(slide, 0, 0, W, H, fill=GIALLO)
    hdr(slide, "Exam Tips", fill=NERO, color=BIANCO, height=HDR_55, size=22)
    icon_y = HDR_55 + Cm(0.35)
    txt(slide, PAD_L, icon_y, Cm(2.2), Cm(2.2), "📋",
        size=32, color=NERO, align=PP_ALIGN.CENTER)
    tip_x = PAD_L + Cm(2.6); tip_w = W - tip_x - Cm(1.0)
    tip_h = Cm(1.65); tip_y0 = HDR_55 + Cm(0.4)
    for i, tip in enumerate(tips[:4]):
        txt(slide, tip_x, tip_y0 + i * tip_h, tip_w, tip_h,
            f"▶  {tip}", size=16, color=NERO, wrap=True)
    qa_top = tip_y0 + min(len(tips), 4) * tip_h + Cm(0.45)
    qa_h = H - qa_top - Cm(0.4)
    if qa_h > Cm(1.5):
        rect(slide, PAD_L, qa_top, W - 2 * PAD_L, qa_h,
             fill=BIANCO, border=NERO, border_w=Pt(1))
        if q_a:
            qa_text = "\n\n".join(f"Q:  {q}\nA:  {a}" for q, a in q_a[:2])
        else:
            qa_text = ""
        if qa_text:
            txt(slide, PAD_L + Cm(0.4), qa_top + Cm(0.25),
                W - 2 * PAD_L - Cm(0.8), qa_h - Cm(0.4),
                qa_text, size=14, color=NERO, wrap=True)

def slide_summary(prs, labels, bodies):
    """labels: 3 label; bodies: 3 testi"""
    slide = new_slide(prs)
    rect(slide, 0, 0, W, H, fill=ARANCIO)
    txt(slide, Cm(1), Cm(0.9), W - Cm(2), Cm(2.0),
        "TAKEAWAY", size=34, color=BIANCO, bold=True, align=PP_ALIGN.CENTER)
    sep_w = Cm(20)
    rect(slide, (W - sep_w) / 2, Cm(2.9), sep_w, Cm(0.06), fill=BIANCO)
    n = 3; mx = Cm(1.2); gap = Cm(0.7)
    bw = (W - 2 * mx - (n - 1) * gap) / n
    bt = Cm(3.3); bh = H - bt - Cm(1.6)
    for i in range(n):
        bx = mx + i * (bw + gap)
        rect(slide, bx, bt, bw, bh, fill=NERO)
        txt(slide, bx, bt + Cm(0.4), bw, Cm(1.6),
            str(i + 1), size=26, color=ARANCIO, bold=True,
            align=PP_ALIGN.CENTER)
        rect(slide, bx + Cm(0.9), bt + Cm(2.0), bw - Cm(1.8), Cm(0.06),
             fill=ARANCIO)
        txt(slide, bx + Cm(0.25), bt + Cm(2.3), bw - Cm(0.5), Cm(1.2),
            labels[i], size=16, color=BIANCO, bold=True,
            align=PP_ALIGN.CENTER)
        txt(slide, bx + Cm(0.35), bt + Cm(3.6), bw - Cm(0.7), bh - Cm(3.8),
            bodies[i], size=14, color=AVORIO, wrap=True,
            align=PP_ALIGN.CENTER)
    txt(slide, Cm(0), H - Cm(1.1), W, Cm(0.9),
        "CCNP ENCOR 350-401  —  Materiale Didattico",
        size=11, color=BIANCO, align=PP_ALIGN.CENTER)

# ─── Generatore deck ──────────────────────────────────────────────────────────

def make_deck(module_id, content):
    """
    content = {
      'title': str, 'area': str, 'hours': str, 'codes': str,
      'agenda': [str, ...],                 # titoli agenda (concetti)
      'topology': {'title', 'caption', 'nodes': [(label,xf,yf)]},
      'sections': [                          # blocchi teorici
        {'section': str, 'subtitle': str,   # slide Section Header
         'slides': [                         # slide Teoria o Tabella
           {'type': 'teoria', 'title': str, 'points': [str], 'key': str},
           {'type': 'config', 'title': str, 'lines': [str], 'device': str, 'hl': int|None},
           {'type': 'verifica', 'title': str, 'lines': [str], 'hl': int|None},
         ]
        }, ...
      ],
      'config_section': {'title': str, 'slides': [...]},  # opz
      'trouble': [(sintomo, fix), ...],
      'exam_tips': [str, ...],
      'exam_qa': [(q, a), ...],
      'summary': {'labels': [str,str,str], 'bodies': [str,str,str]}
    }
    """
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H

    # 1 — Cover
    slide_cover(prs, module_id, content['title'],
                content['area'], content['hours'], content['codes'])

    # 2 — Agenda
    slide_agenda(prs, content['agenda'])

    # 3 — Topologia
    topo = content.get('topology')
    if topo:
        slide_diagramma(prs, topo['title'], topo['caption'], topo['nodes'])

    # 4 — Sezioni teoriche
    for sec in content.get('sections', []):
        slide_section(prs, sec['section'], sec.get('subtitle', ''))
        for sl in sec.get('slides', []):
            if sl['type'] == 'teoria':
                slide_teoria(prs, sl['title'], sl['points'], sl.get('key', ''))
            elif sl['type'] == 'config':
                slide_config(prs, sl['title'], sl['lines'],
                             sl.get('device', 'R1'), sl.get('hl'))
            elif sl['type'] == 'verifica':
                slide_verifica(prs, sl['title'], sl['lines'], sl.get('hl'))

    # 5 — Sezione config/verifica se separata
    cfg_sec = content.get('config_section')
    if cfg_sec:
        slide_section(prs, "CONFIGURAZIONE & VERIFICA")
        for sl in cfg_sec.get('slides', []):
            if sl['type'] == 'config':
                slide_config(prs, sl['title'], sl['lines'],
                             sl.get('device', 'R1'), sl.get('hl'))
            elif sl['type'] == 'verifica':
                slide_verifica(prs, sl['title'], sl['lines'], sl.get('hl'))

    # 6 — Troubleshooting
    if content.get('trouble'):
        slide_troubleshooting(prs, content['trouble'])

    # 7 — Exam Tips
    slide_exam_tips(prs, content['exam_tips'], content.get('exam_qa'))

    # 8 — Summary
    s = content['summary']
    slide_summary(prs, s['labels'], s['bodies'])

    return prs


# ══════════════════════════════════════════════════════════════════════════════
# CONTENUTO MODULI
# ══════════════════════════════════════════════════════════════════════════════

MODULES = {}

# ─── MOD-01 ───────────────────────────────────────────────────────────────────
MODULES['MOD-01'] = {
    'title': 'OSPFv2 Fondamenta',
    'area': 'AREA 1 — OSPF', 'hours': '2h', 'codes': '3.2.a · 3.2.b',
    'agenda': [
        'Sub-interface 802.1Q: configurazione e best practice',
        'OSPF: processo, Router-ID e selezione',
        'Network statement vs ip ospf area',
        'DR/BDR Election: regole e controllo',
        'Network type point-to-point e costo OSPF',
        'Troubleshooting adiacenze: stati e cause',
    ],
    'topology': {
        'title': 'Topologia MOD-01 — OSPFv2 Multi-Area',
        'caption': 'R3-R6: Area 0 Backbone (broadcast + ring P2P) · R1/R7: Area 15/99 · R2: Area 25 · ABR: R5',
        'nodes': [
            ('R1\nABR 15/99', 0.1, 0.4),
            ('R3\nDROTHER', 0.3, 0.25),
            ('R4\nDR prio 255', 0.45, 0.4),
            ('R5\nABR 0/15/25', 0.6, 0.4),
            ('R6\nBDR prio 100', 0.75, 0.25),
            ('R2\nArea 25', 0.9, 0.55),
            ('R7\nArea 99', 0.1, 0.7),
        ],
    },
    'sections': [
        {
            'section': 'Sub-interface 802.1Q',
            'subtitle': 'Multiplexing logico su interfaccia fisica',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Sub-interface 802.1Q — Come funziona',
                    'points': [
                        'Una sub-interface (virtuale) suddivide logicamente una porta fisica con tag VLAN 802.1Q',
                        'Sintassi IOS: interface eth0/0.<vlan>  →  encapsulation dot1Q <vlan>',
                        "L'interfaccia fisica padre deve avere no ip address + no shutdown (nessuna configurazione IP)",
                        'Il numero sub-interface coincide per convenzione con il numero VLAN (leggibilità)',
                        'Ogni sub-interface ha il proprio indirizzo IP e partecipa independentemente ai protocolli',
                        'In ambiente IOU: tutta la connettività inter-router avviene via sub-interface (no link fisici diretti)',
                    ],
                    'key': 'La fisica è uno, il logico è molti — tag VLAN separa il traffico di più reti su un solo cavo.',
                },
                {
                    'type': 'config',
                    'title': 'Sub-interface 802.1Q — Snippet',
                    'lines': [
                        'interface Ethernet0/0',
                        ' no ip address',
                        ' no shutdown          ! padre deve essere UP',
                        '!',
                        'interface Ethernet0/0.34',
                        ' encapsulation dot1Q 34     ! tag VLAN 34',
                        ' ip address 10.0.34.1 255.255.255.252',
                        ' description P2P_R3-R4_Area0',
                        ' no shutdown',
                        ' ip ospf 100 area 0         ! abilita direttamente OSPF',
                    ],
                    'device': 'R3', 'hl': 5,
                },
            ],
        },
        {
            'section': 'OSPF: Processo e Router-ID',
            'subtitle': 'Identità univoca nel dominio di routing',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'OSPF Router-ID — Selezione e Best Practice',
                    'points': [
                        'Il Router-ID identifica univocamente il router nella LSDB OSPF',
                        'Ordine di selezione: 1) router-id configurato  2) IP loopback più alto  3) IP fisico più alto',
                        'Best practice: configurare sempre esplicitamente (router-id x.x.x.x)',
                        'Senza configurazione esplicita: cambio IP = cambio Router-ID = instabilità LSDB',
                        'Convenzione laboratorio: R1 → 1.1.1.1,  R2 → 2.2.2.2, ... R7 → 7.7.7.7',
                        'Network statement vs ip ospf area: due metodi equivalenti per abilitare OSPF su un\'interfaccia',
                    ],
                    'key': 'router-id esplicito = comportamento deterministico. Senza: rischio instabilità al riavvio.',
                },
                {
                    'type': 'teoria',
                    'title': 'Network Statement vs ip ospf area',
                    'points': [
                        '• network 10.0.0.0 0.0.0.7 area 0 → abilita tutte le interfacce con IP nel range',
                        '• ip ospf 100 area 0 (su interfaccia) → esplicito, per-interfaccia, più leggibile',
                        'passive-interface: impedisce l\'invio di Hello ma mantiene il prefisso nella LSDB',
                        'Tecnica sicura: passive-interface default + no passive-interface su porte OSPF attive',
                        'ip ospf network point-to-point: elimina DR/BDR su link con un solo neighbor (sub-if P2P)',
                        'Broadcast (default Ethernet): DR/BDR election — usare solo su segmenti multi-router',
                    ],
                    'key': 'ip ospf area per-interfaccia è più preciso e meno soggetto a errori di wildcard.',
                },
            ],
        },
        {
            'section': 'DR/BDR Election',
            'subtitle': 'Riduzione adiacenze su segmenti broadcast',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'DR/BDR — Perché esiste e Come funziona',
                    'points': [
                        'Problema: N router su broadcast → N*(N-1)/2 adiacenze (45 con 10 router)',
                        'Soluzione: DR centralizza — tutti inviano aggiornamenti al DR (224.0.0.6)',
                        'DR li riflette a tutti (224.0.0.5) — riduce a N-1 adiacenze reali',
                        'Regole election: 1) priority più alta  2) Router-ID più alto (pareggio)',
                        'NON preemptive: il DR caduto che torna diventa DROTHER anche con priority alta',
                        'Forza rielelezione: clear ip ospf process (interrompe tutte le adiacenze)',
                    ],
                    'key': 'Election non preemptive: modificare priority non basta — serve clear ip ospf process.',
                },
                {
                    'type': 'verifica',
                    'title': 'Verifica DR/BDR Election',
                    'lines': [
                        'R4# show ip ospf interface ethernet 0/0.3456',
                        '  Network Type BROADCAST, Cost: 10',
                        '  Designated Router (ID) 4.4.4.4, Interface address 10.0.0.4',
                        '  Backup Designated router (ID) 6.6.6.6, Interface address 10.0.0.6',
                        '',
                        'R4# show ip ospf neighbor',
                        'Neighbor ID  Pri  State      Dead Time  Address    Interface',
                        '3.3.3.3        0  FULL/DROTHER  00:00:38  10.0.0.3  Et0/0.3456',
                        '5.5.5.5        0  FULL/DROTHER  00:00:36  10.0.0.5  Et0/0.3456',
                        '6.6.6.6      100  FULL/BDR      00:00:37  10.0.0.6  Et0/0.3456',
                    ],
                    'hl': 7,
                },
            ],
        },
        {
            'section': 'Troubleshooting Adiacenze OSPF',
            'subtitle': 'Macchina a stati e cause di blocco',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Stati Adiacenza OSPF — Dove si blocca e Perché',
                    'points': [
                        'DOWN → INIT: Hello ricevuto ma router locale non è nella lista neighbor (link unidirezionale)',
                        'INIT → 2-WAY: Hello bidirezionale OK (DROTHER-DROTHER fermano qui — normale)',
                        'EXSTART: negoziazione master/slave per DB exchange → MTU mismatch blocca qui',
                        'EXCHANGE: scambio DBD → MTU mismatch o database corrotto',
                        'FULL: adiacenza completa e LSDB sincronizzata — stato desiderato',
                        'Triade errori più comuni: timer mismatch · area mismatch · authentication mismatch',
                    ],
                    'key': 'EXSTART/EXCHANGE = quasi sempre MTU mismatch. Rimedio: ip ospf mtu-ignore.',
                },
            ],
        },
    ],
    'config_section': {
        'slides': [
            {
                'type': 'config',
                'title': 'OSPF — Comandi Chiave',
                'lines': [
                    'router ospf 100',
                    ' router-id 3.3.3.3            ! sempre esplicito',
                    ' network 10.0.0.0 0.0.0.7 area 0',
                    ' passive-interface default     ! blocca Hello su tutte',
                    ' no passive-interface Et0/0.3456  ! riabilita OSPF attivo',
                    '!',
                    'interface Ethernet0/0.34',
                    ' ip ospf 100 area 0            ! metodo diretto',
                    ' ip ospf network point-to-point ! niente DR/BDR',
                    ' ip ospf cost 1000             ! sfavorisce il percorso',
                    ' ip ospf priority 0            ! escluso dalla election',
                ],
                'device': 'R3', 'hl': 4,
            },
        ],
    },
    'trouble': [
        ('Adiacenza bloccata in EXSTART', 'MTU mismatch — ip ospf mtu-ignore su entrambi i lati'),
        ('Neighbor in INIT, non avanza', 'Autenticazione MD5 su un solo lato — verificare show ip ospf interface'),
        ('DR non è il router desiderato', 'Election non preemptive — eseguire clear ip ospf process dopo aver modificato priority'),
        ('Sub-interface non passa traffico', 'encapsulation dot1Q mancante o VLAN non permessa sul trunk switch'),
    ],
    'exam_tips': [
        'Router-ID: configurare sempre esplicitamente. Senza: cambio IP = instabilità OSPF',
        'ip ospf network point-to-point su link P2P elimina DR/BDR e accelera convergenza',
        'Election DR/BDR NON è preemptive: priority alta non basta, serve clear ip ospf process',
        'EXSTART/EXCHANGE bloccato = MTU mismatch (quasi sempre)',
    ],
    'exam_qa': [
        ('Un router OSPF resta in EXSTART. Causa più probabile?',
         'MTU mismatch — impedisce lo scambio dei DBD. Fix: ip ospf mtu-ignore'),
        ('Come escludere un router dalla election senza disabilitare OSPF?',
         'ip ospf priority 0 sull\'interfaccia — partecipa all\'area ma non vince mai DR/BDR'),
    ],
    'summary': {
        'labels': ['Router-ID Esplicito', 'Network Type P2P', 'DR/BDR Election'],
        'bodies': [
            'Sempre configurare router-id x.x.x.x: comportamento deterministico, nessuna sorpresa al riavvio.',
            'ip ospf network point-to-point su sub-if con un solo neighbor: elimina DR/BDR, convergenza più rapida.',
            'Non preemptive: modificare priority richiede clear ip ospf process per forzare la rielelezione.',
        ],
    },
}

# ─── MOD-02 ───────────────────────────────────────────────────────────────────
MODULES['MOD-02'] = {
    'title': 'OSPFv2 Aree & Summarization',
    'area': 'AREA 1 — OSPF', 'hours': '2h', 'codes': '3.2.a · 3.2.b',
    'agenda': [
        'Architettura OSPF multi-area: ABR e backbone',
        'Summarization inter-area con area range',
        'Stub Area e Totally-Stub Area',
        'NSSA e redistribuzione esterna',
        'Virtual Link: quando e come usarlo',
        'ASBR e summary-address per rotte esterne',
    ],
    'topology': {
        'title': 'Topologia MOD-02 — Aree e ABR',
        'caption': 'R5: ABR Area0/15/25 · R1: ABR Area15/99 · Area 0: backbone con R3-R4-R5-R6 · summarization su R5',
        'nodes': [
            ('R7\nArea 99', 0.08, 0.3),
            ('R1\nABR 15/99', 0.22, 0.3),
            ('R5\nABR 0/15/25', 0.5, 0.5),
            ('R3/R4/R6\nArea 0', 0.65, 0.35),
            ('R2\nArea 25', 0.85, 0.5),
        ],
    },
    'sections': [
        {
            'section': 'OSPF Multi-Area',
            'subtitle': 'Scalabilità e riduzione della LSDB',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Architettura Multi-Area OSPF',
                    'points': [
                        'Area 0 (backbone): obbligatoria — tutte le altre aree devono connettersi ad area 0',
                        'ABR (Area Border Router): connette due o più aree, mantiene LSDB separate per ciascuna',
                        'LSA Type 1 (Router LSA): restano confinati nell\'area di origine',
                        'LSA Type 3 (Summary LSA): generati dall\'ABR per annunciare prefissi inter-area',
                        'Vantaggio: fallimento in area X non causa recalcolo SPF nelle altre aree',
                        'ASBR: redistribuisce rotte esterne (LSA Type 5) verso tutto il dominio OSPF',
                    ],
                    'key': 'Ogni area ha la propria LSDB — ABR traduce Type 1/2 in Type 3 verso le altre aree.',
                },
                {
                    'type': 'teoria',
                    'title': 'Stub, Totally-Stub e NSSA',
                    'points': [
                        'Stub area: blocca LSA Type 5 (rotte esterne) — ABR inietta default route (Type 3)',
                        'Totally-Stub: blocca anche LSA Type 3 (inter-area) — solo default route rimane',
                        'NSSA (Not-So-Stubby Area): come stub ma permette redistribuzione locale (LSA Type 7)',
                        'LSA Type 7 → Type 5: l\'ABR converte quando esce dall\'NSSA verso il backbone',
                        'Tutti i router in un\'area stub/NSSA devono avere la stessa configurazione (area X stub)',
                        'Totally-NSSA (area X nssa no-summary): massima riduzione LSDB con redistribuzione locale',
                    ],
                    'key': 'Totally-Stub = minima LSDB possibile: solo prefissi interni + default route dall\'ABR.',
                },
            ],
        },
        {
            'section': 'Summarization e Virtual Link',
            'subtitle': 'Ottimizzazione e continuità backbone',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Summarization Inter-Area e ASBR',
                    'points': [
                        'area range (su ABR): aggrega prefissi inter-area — riduce LSDB nelle aree adiacenti',
                        'Sintassi: area 15 range 10.15.0.0 255.255.252.0  (su ABR che confina con Area 15)',
                        'summary-address (su ASBR): aggrega rotte esterne (Type 5) prima dell\'annuncio',
                        'Vantaggio: meno LSA = convergenza più rapida + meno memoria nei router interni',
                        'Attenzione: la summary route è attiva solo se almeno un prefisso subordinato esiste',
                        'Discard route: IOS installa automaticamente una rotta Null0 per evitare loop',
                    ],
                    'key': 'area range sull\'ABR: N prefissi dell\'area diventano 1 summary verso il backbone.',
                },
                {
                    'type': 'teoria',
                    'title': 'Virtual Link — Quando e Come',
                    'points': [
                        'Problema: una nuova area non può connettersi direttamente ad Area 0',
                        'Virtual Link: estende logicamente Area 0 attraverso una transit area non-stub',
                        'Sintassi: area <transit-area> virtual-link <router-id-remoto>  su entrambi gli ABR',
                        'Limitazioni: la transit area NON può essere stub/totally-stub/NSSA',
                        'Il Virtual Link è un\'interfaccia OSPF logica — vive in Area 0, transita altrove',
                        'Uso reale limitato: preferire ridisegno topologico se possibile',
                    ],
                    'key': 'Virtual Link come eccezione, non come design normale — transit area non può essere stub.',
                },
            ],
        },
    ],
    'config_section': {
        'slides': [
            {
                'type': 'config',
                'title': 'ABR Summarization + Stub + Virtual Link',
                'lines': [
                    'router ospf 100',
                    ' area 15 range 10.15.0.0 255.255.252.0  ! summarize Area 15',
                    ' area 25 range 10.25.0.0 255.255.252.0  ! summarize Area 25',
                    ' area 25 stub                            ! totally-stub: aggiungere no-summary',
                    ' area 99 virtual-link 1.1.1.1            ! VL via Area 15 (R5→R1)',
                    '!',
                    '! Su ASBR per rotte redistribuite:',
                    ' summary-address 192.168.100.0 255.255.252.0',
                ],
                'device': 'R5', 'hl': 1,
            },
            {
                'type': 'verifica',
                'title': 'Verifica Summarization e Stub',
                'lines': [
                    'R5# show ip route ospf',
                    'O IA 10.15.0.0/22 [110/11] via 10.1.15.2   ! summary Area 15',
                    'O*IA 0.0.0.0/0   [110/1]  via 10.0.0.4     ! default in stub',
                    '',
                    'R5# show ip ospf border-routers',
                    'OSPF Process 100 internal Routing Table',
                    'Codes: i - Intra-area route, I - Inter-area route',
                    'i 1.1.1.1 [11] via 10.1.15.2, Et0/0.51, ABR, Area 15',
                ],
                'hl': 1,
            },
        ],
    },
    'trouble': [
        ('Virtual Link non si forma', 'Transit area configurata come stub — VL richiede transit area non-stub'),
        ('Summary non appare nella routing table', 'Nessun prefisso subordinato attivo — verificare che almeno una rotta dell\'area esista'),
        ('Default route non ricevuta in stub area', 'Non tutti i router dell\'area hanno "area X stub" — configurazione deve essere uniforme'),
        ('LSA Type 5 presenti in stub area', 'Un router dell\'area non ha "area X stub" — blocca propagazione'),
    ],
    'exam_tips': [
        'Tutti i router in una stub area devono avere "area X stub" — configurazione non-uniforme = adiacenza down',
        'area range va configurato sull\'ABR, lato area da aggregare — non nel backbone',
        'Virtual Link: transit area non può essere stub/NSSA — altrimenti non funziona',
        'LSA Type 7 (NSSA) viene convertito in Type 5 dall\'ABR verso il backbone',
    ],
    'exam_qa': [
        ('Differenza tra Stub e Totally-Stub?',
         'Stub blocca Type 5 (esterne). Totally-Stub blocca anche Type 3 (inter-area) — solo default route rimane.'),
        ('Perché la transit area di un Virtual Link non può essere stub?',
         'Il VL appartiene ad Area 0: richiederebbe LSA di Area 0 attraverso la transit area, impossibile in stub.'),
    ],
    'summary': {
        'labels': ['Gerarchia Aree', 'Summarization ABR', 'Stub/Totally-Stub'],
        'bodies': [
            'Area 0 obbligatoria. ABR traduce LSA Type 1/2 in Type 3. Virtual Link solo come eccezione.',
            'area range sull\'ABR aggrega N prefissi in 1 summary — riduce LSDB e accelera convergenza.',
            'Stub blocca Type 5; Totally-Stub blocca anche Type 3. NSSA permette redistribuzione locale.',
        ],
    },
}

# ─── MOD-03 ───────────────────────────────────────────────────────────────────
MODULES['MOD-03'] = {
    'title': 'OSPFv3 Dual-Stack',
    'area': 'AREA 1 — OSPF', 'hours': '1.5h', 'codes': '3.2.b',
    'agenda': [
        'OSPFv2 vs OSPFv3: differenze architetturali',
        'OSPFv3 per IPv6: processo e configurazione',
        'Address-Family: dual-stack su singolo processo',
        'Link-Local come sorgente Hello OSPFv3',
        'Verifica e troubleshooting OSPFv3',
    ],
    'topology': {
        'title': 'Topologia MOD-03 — OSPFv3 Dual-Stack',
        'caption': 'Stessa topologia MOD-02 con indirizzi IPv6. Area 0: 2001:db8:0::/64. R5 ABR dual-stack.',
        'nodes': [
            ('R1\nArea 15\nIPv6', 0.15, 0.35),
            ('R5\nABR 0/15/25\nDual-stack', 0.5, 0.5),
            ('R3/R4/R6\nArea 0\n2001:db8:0::/64', 0.7, 0.3),
            ('R2\nArea 25\nIPv6', 0.85, 0.6),
        ],
    },
    'sections': [
        {
            'section': 'OSPFv3 vs OSPFv2',
            'subtitle': 'Evoluzione per il supporto IPv6 e dual-stack',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Differenze Chiave OSPFv2 → OSPFv3',
                    'points': [
                        'OSPFv2: solo IPv4. OSPFv3: IPv6 native, poi esteso a dual-stack con address-family',
                        'Transport: OSPFv3 usa link-local IPv6 (fe80::) come sorgente dei pacchetti Hello',
                        'Router-ID: rimane un valore 32-bit anche in OSPFv3 (non è un indirizzo IPv6)',
                        'Autenticazione: OSPFv2 usa area/interface auth; OSPFv3 usa IPSec (AH/ESP)',
                        'Configurazione: ipv6 router ospf <pid> oppure router ospfv3 <pid> address-family',
                        'Link-local come next-hop: OSPFv3 usa il link-local del neighbor come next-hop nelle rotte',
                    ],
                    'key': 'OSPFv3 usa link-local IPv6 per i Hello — il Router-ID rimane 32-bit come in OSPFv2.',
                },
                {
                    'type': 'teoria',
                    'title': 'Configurazione OSPFv3 — Metodi',
                    'points': [
                        'Metodo 1 (classic): ipv6 router ospf <pid> + ipv6 ospf <pid> area X sull\'interfaccia',
                        'Metodo 2 (address-family): router ospfv3 <pid> + address-family ipv6 unicast per IPv6',
                        'Dual-stack: router ospfv3 include sia address-family ipv4 che ipv6 — un solo processo',
                        'Su ogni interfaccia: ospfv3 <pid> ipv4 area X  e  ospfv3 <pid> ipv6 area X',
                        'ipv6 unicast-routing: obbligatorio per abilitare il forwarding IPv6 sul router',
                        'passive-interface funziona anche in OSPFv3 — stessa logica OSPFv2',
                    ],
                    'key': 'address-family nel router ospfv3: un solo processo gestisce sia IPv4 che IPv6 — approccio moderno.',
                },
            ],
        },
        {
            'section': 'Verifica e Differenze Operative',
            'subtitle': 'Comandi show e interpretazione output',
            'slides': [
                {
                    'type': 'verifica',
                    'title': 'Verifica OSPFv3 Dual-Stack',
                    'lines': [
                        'R5# show ospfv3 neighbor',
                        'OSPFv3 100 address-family ipv6 (router-id 5.5.5.5)',
                        'Neighbor ID  Pri  State   Dead  Interface-ID  Interface',
                        '3.3.3.3        1  FULL/DROTHER  00:00:38  6  Et0/0.3456',
                        '4.4.4.4        1  FULL/DR       00:00:36  7  Et0/0.3456',
                        '',
                        'R5# show ipv6 route ospf',
                        'OI 2001:db8:15::/64 [110/11] via FE80::1, Et0/0.51',
                    ],
                    'hl': 3,
                },
            ],
        },
    ],
    'config_section': {
        'slides': [
            {
                'type': 'config',
                'title': 'OSPFv3 Address-Family — Snippet',
                'lines': [
                    'ipv6 unicast-routing',
                    '!',
                    'router ospfv3 100',
                    ' router-id 5.5.5.5',
                    ' address-family ipv4 unicast',
                    '  area 0 range 10.0.0.0 255.255.248.0',
                    ' address-family ipv6 unicast',
                    '  area 0 range 2001:db8::/48',
                    '!',
                    'interface Ethernet0/0.3456',
                    ' ospfv3 100 ipv4 area 0',
                    ' ospfv3 100 ipv6 area 0',
                ],
                'device': 'R5', 'hl': 4,
            },
        ],
    },
    'trouble': [
        ('Adiacenza OSPFv3 non si forma', 'ipv6 unicast-routing non abilitato — prerequisito assoluto per OSPFv3'),
        ('Router-ID non configurato', 'Nessuna interfaccia IPv4 attiva — configurare router-id manualmente in ospfv3'),
        ('Next-hop nelle rotte IPv6 è link-local', 'Comportamento normale in OSPFv3 — il link-local fe80:: è usato come next-hop'),
        ('show ospfv3 neighbor vuoto', 'Interfaccia non abilitata in OSPFv3 — verificare ospfv3 <pid> area su ogni interfaccia'),
    ],
    'exam_tips': [
        'OSPFv3 usa link-local come sorgente Hello — il Router-ID rimane 32-bit (come OSPFv2)',
        'ipv6 unicast-routing: obbligatorio prima di configurare qualsiasi protocollo IPv6',
        'address-family ipv4/ipv6 in ospfv3: un processo gestisce entrambi i protocolli',
        'Autenticazione OSPFv3 usa IPSec (non MD5 come OSPFv2)',
    ],
    'exam_qa': [
        ('OSPFv3 usa quale indirizzo come sorgente dei pacchetti Hello?',
         'Il link-local IPv6 (fe80::) dell\'interfaccia — non l\'indirizzo globale.'),
        ('Come mantenere la compatibilità dual-stack con un solo processo OSPF?',
         'router ospfv3 con address-family ipv4 unicast e address-family ipv6 unicast.'),
    ],
    'summary': {
        'labels': ['Link-Local Source', 'Router-ID 32-bit', 'Address-Family'],
        'bodies': [
            'OSPFv3 usa fe80:: come sorgente Hello — indirizzo link-local, non globale.',
            'Il Router-ID rimane 32-bit anche in OSPFv3 — deve essere configurato se mancano interfacce IPv4.',
            'ospfv3 con address-family: un processo gestisce IPv4 e IPv6 — approccio moderno e pulito.',
        ],
    },
}

# ─── MOD-04 ───────────────────────────────────────────────────────────────────
MODULES['MOD-04'] = {
    'title': 'OSPF Troubleshooting',
    'area': 'AREA 1 — OSPF', 'hours': '4h', 'codes': '1.10.a-d',
    'agenda': [
        'Metodologia troubleshooting OSPF strutturata',
        'Scenario A: adiacenze mancanti (timer, auth, area)',
        'Scenario B: rotte assenti nella routing table',
        'Scenario C: External Type E1 vs E2 e metriche',
        'Scenario D: OSPFv3 dual-stack diagnostica',
        'Tool: debug ip ospf e interpretazione output',
    ],
    'topology': {
        'title': 'Topologia MOD-04 — Troubleshooting Lab',
        'caption': 'R3: ABR + ASBR (redistribuisce static) · R4: ABR · R5: Area 2 Stub · auth MD5 su R1-R4',
        'nodes': [
            ('R1\nArea 0', 0.12, 0.3),
            ('R2\nArea 0', 0.35, 0.3),
            ('R3\nABR+ASBR', 0.58, 0.5),
            ('R4\nABR 0/1/2', 0.75, 0.3),
            ('R5\nArea 2 Stub', 0.9, 0.5),
        ],
    },
    'sections': [
        {
            'section': 'Metodologia Troubleshooting',
            'subtitle': 'Approccio sistematico a problemi OSPF',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Framework Troubleshooting OSPF',
                    'points': [
                        'Livello 1 — Fisica: interfaccia UP/UP? show interfaces / show ip interface brief',
                        'Livello 2 — Adiacenza: show ip ospf neighbor → stato e Dead timer',
                        'Livello 3 — LSDB: show ip ospf database → LSA presenti e completi',
                        'Livello 4 — Routing Table: show ip route ospf → prefissi installati nel RIB',
                        'Livello 5 — Forwarding: ping / traceroute → traffico effettivo',
                        'Debug: debug ip ospf adj | hello → usare con cautela, poi undebug all',
                    ],
                    'key': 'Seguire la pila dal basso: fisica → adiacenza → LSDB → RIB → forwarding.',
                },
                {
                    'type': 'teoria',
                    'title': 'E1 vs E2 — Metriche Rotte Esterne',
                    'points': [
                        'E1 (External Type 1): metrica = costo OSPF interno + seed metric esterna',
                        'E2 (External Type 2, default): metrica = solo seed metric esterna (costo interno ignorato)',
                        'E2 preferisce il percorso con seed metric più bassa — indifferente alla topologia interna',
                        'E1 preferisce il percorso totale più corto — tiene conto della topologia OSPF',
                        'Redistribuzione: redistribute static metric 20 metric-type 1 subt type come E1',
                        'In pratica: E1 per rotte dove conta la topologia interna; E2 per rotte "flat"',
                    ],
                    'key': 'E2 default: metrica statica. E1: metrica cumulativa. Usare E1 quando la topologia interna conta.',
                },
            ],
        },
    ],
    'config_section': {
        'slides': [
            {
                'type': 'config',
                'title': 'Debug e Diagnostica OSPF',
                'lines': [
                    '! --- Sequenza diagnostica standard ---',
                    'show ip ospf neighbor          ! stato adiacenze',
                    'show ip ospf neighbor detail   ! MTU, auth, timers',
                    'show ip ospf database          ! LSDB completa',
                    'show ip ospf database external ! solo LSA Type 5',
                    'show ip route ospf             ! RIB OSPF',
                    '!',
                    'debug ip ospf adj              ! messaggi adiacenza',
                    'debug ip ospf hello            ! pacchetti Hello',
                    'undebug all                    ! stop debug',
                ],
                'device': 'R1', 'hl': 1,
            },
        ],
    },
    'trouble': [
        ('Area 0 non si vede da area stub', 'LSA Type 5 bloccati in stub area — usare NSSA se serve redistribuzione locale'),
        ('E2 route preferisce percorso sbagliato', 'E2 ignora costo interno — passare a metric-type 1 per routing topology-aware'),
        ('MD5 auth: adiacenza bloccata in INIT', 'Key mismatch o key-id diverso — verificare ip ospf authentication message-digest e key id'),
        ('OSPFv3 prefix non appare nel RIB', 'ipv6 unicast-routing non abilitato o interfaccia non abilitata in area'),
    ],
    'exam_tips': [
        'E2 (default): metrica solo esterna. E1: metrica esterna + interna. E1 più corretto per topologie complesse',
        'show ip ospf neighbor detail: mostra MTU, dead timer, auth type — tutto ciò che causa blocco',
        'Stub area: area X stub su TUTTI i router dell\'area — anche su un solo router = adiacenza down',
        'debug ip ospf adj: vedere perché due router non formano adiacenza in tempo reale',
    ],
    'exam_qa': [
        ('Differenza tra E1 e E2 nelle rotte OSPF esterne?',
         'E2 (default): metrica = solo seed metric esterna. E1: metrica = seed + costo OSPF interno accumulato.'),
        ('Quale comando mostra l\'autenticazione configurata su un\'interfaccia OSPF?',
         'show ip ospf interface <int> — mostra auth type, key-id e stato.'),
    ],
    'summary': {
        'labels': ['Framework L1→L5', 'E1 vs E2', 'Debug Selettivo'],
        'bodies': [
            'Fisica → adiacenza → LSDB → RIB → forwarding. Non saltare livelli.',
            'E2 metrica piatta (default); E1 cumulativa. Usare E1 quando la topologia interna conta.',
            'debug ip ospf adj/hello + undebug all. show ip ospf neighbor detail = primo posto dove guardare.',
        ],
    },
}

# ─── MOD-05 ───────────────────────────────────────────────────────────────────
MODULES['MOD-05'] = {
    'title': 'BGP Fondamenta',
    'area': 'AREA 2 — BGP', 'hours': '2h', 'codes': '3.2.c · 1.11.a · 1.11.b',
    'agenda': [
        'BGP: Autonomous System e sessioni eBGP/iBGP',
        'iBGP full-mesh: update-source e next-hop-self',
        'BGP network statement e Origin attribute',
        'Processo di selezione best-path BGP',
        'Prefix-list e route-map come filtri',
        'Troubleshooting peering BGP',
    ],
    'topology': {
        'title': 'Topologia MOD-05 — BGP Multi-AS',
        'caption': 'AS 65001 (ISP): R1-R2-R3 iBGP full-mesh. AS 65000 (Customer): R4-R5-R6. eBGP: R1↔R4 e R3↔R5.',
        'nodes': [
            ('R1\nAS 65001\nBorder', 0.15, 0.3),
            ('R2\nAS 65001\nInternal', 0.35, 0.5),
            ('R3\nAS 65001\nBorder', 0.55, 0.3),
            ('R4\nAS 65000\nBorder', 0.15, 0.7),
            ('R5\nAS 65000\nBorder', 0.55, 0.7),
            ('R6\nAS 65000\nInternal', 0.85, 0.5),
        ],
    },
    'sections': [
        {
            'section': 'BGP Fondamenta',
            'subtitle': 'Il protocollo di routing di Internet',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'BGP — Autonomous System e Sessioni',
                    'points': [
                        'BGP: Border Gateway Protocol — routing inter-AS (EGP), path vector protocol',
                        'AS (Autonomous System): insieme di reti sotto un\'unica amministrazione, numero 16/32-bit',
                        'eBGP: sessione tra router in AS diversi — tipicamente link diretto, TTL=1',
                        'iBGP: sessione tra router nello stesso AS — tipicamente via loopback, no TTL limit',
                        'TCP port 179: BGP usa TCP come transport — no hello periodici, keepalive ogni 60s',
                        'BGP non scopre la topologia: riceve prefissi con path attributes da peer',
                    ],
                    'key': 'eBGP = tra AS diversi. iBGP = stesso AS via loopback. BGP è path vector, non distance vector.',
                },
                {
                    'type': 'teoria',
                    'title': 'iBGP Full-Mesh: update-source e next-hop-self',
                    'points': [
                        'iBGP full-mesh: ogni router deve avere sessione con tutti gli altri (N*(N-1)/2)',
                        'update-source Loopback0: usa loopback come sorgente TCP — più stabile di un indirizzo fisico',
                        'Problema next-hop iBGP: il next-hop di una rotta eBGP non cambia in iBGP',
                        'next-hop-self: forza il router a diventare next-hop per i prefissi annunciati agli iBGP peer',
                        'IBGP split-horizon: un router iBGP non ri-annuncia prefissi appresi da iBGP ad altri iBGP',
                        'Soluzione full-mesh: Route Reflector (MOD-07) o Confederation per grandi AS',
                    ],
                    'key': 'Sempre: update-source Lo0 + next-hop-self su ogni iBGP neighbor che non ha visibilità diretta.',
                },
                {
                    'type': 'teoria',
                    'title': 'BGP Best-Path Selection — Regole',
                    'points': [
                        '1. Weight (Cisco proprietario): più alto vince — locale al router, non propagato',
                        '2. Local Preference: più alto vince — propagato in iBGP, controlla traffico uscente',
                        '3. Locally originated: network/redistribute/aggregate preferito su appreso',
                        '4. AS-Path length: più corto vince',
                        '5. Origin: IGP (i) < EGP (e) < Incomplete (?)',
                        '6. MED: più basso vince — suggerisce al peer l\'exit point preferito (entrante)',
                    ],
                    'key': 'Mnemonico: "We Love Oranges As Oranges Mean Pure Refreshment" — Weight, LP, Locally originated...',
                },
            ],
        },
    ],
    'config_section': {
        'slides': [
            {
                'type': 'config',
                'title': 'BGP iBGP Full-Mesh — Configurazione Base',
                'lines': [
                    'router bgp 65000',
                    ' bgp router-id 4.4.4.4',
                    ' bgp log-neighbor-changes',
                    ' neighbor 5.5.5.5 remote-as 65000      ! iBGP',
                    ' neighbor 5.5.5.5 update-source Lo0    ! stabilità',
                    ' neighbor 5.5.5.5 next-hop-self        ! fix next-hop',
                    ' neighbor 172.16.14.1 remote-as 65001  ! eBGP',
                    ' !',
                    ' address-family ipv4',
                    '  network 4.4.4.4 mask 255.255.255.255 ! annuncia loopback',
                    '  neighbor 5.5.5.5 activate',
                    '  neighbor 172.16.14.1 activate',
                ],
                'device': 'R4', 'hl': 4,
            },
            {
                'type': 'verifica',
                'title': 'Verifica BGP Peering',
                'lines': [
                    'R4# show bgp summary',
                    'Neighbor      V  AS    MsgRcvd  MsgSent  Up/Down   State/PfxRcd',
                    '5.5.5.5       4  65000     150     148  01:22:14         3',
                    '172.16.14.1   4  65001      88      90  00:45:30         5',
                    '',
                    'R4# show bgp ipv4 unicast 100.0.0.1',
                    'BGP routing table entry for 100.0.0.0/8',
                    '  100.0.0.0/8, from 172.16.14.1, Origin IGP, metric 0, localpref 100',
                    '  Path: 65001',
                ],
                'hl': 2,
            },
        ],
    },
    'trouble': [
        ('Sessione BGP in ACTIVE', 'TCP non raggiunge il peer — verificare routing verso update-source, ACL, TTL eBGP'),
        ('Prefisso non appare nel BGP table', 'network statement: il prefisso deve esistere nel RIB locale con la maschera esatta'),
        ('Next-hop non raggiungibile in iBGP', 'Manca next-hop-self sul border router — il next-hop eBGP non è noto internamente'),
        ('Prefisso appresto da iBGP non ri-annunciato', 'iBGP split-horizon — serve full-mesh o Route Reflector'),
    ],
    'exam_tips': [
        'update-source Loopback0 + next-hop-self: coppia inseparabile per iBGP corretto',
        'BGP usa TCP 179 — se la sessione non sale, prima verificare connettività TCP',
        'network statement BGP: il prefisso DEVE esistere nel RIB con maschera esatta',
        'Local Preference (default 100): più alto = preferito. Controlla traffico USCENTE dall\'AS.',
    ],
    'exam_qa': [
        ('Perché next-hop-self è necessario in iBGP?',
         'Il next-hop eBGP non cambia in iBGP: i router interni non sanno raggiungere il next-hop esterno.'),
        ('Differenza tra Weight e Local Preference?',
         'Weight: Cisco only, locale al router, non propagato. Local-Pref: propagato in iBGP, condiviso nell\'AS.'),
    ],
    'summary': {
        'labels': ['eBGP vs iBGP', 'update-source + NH-self', 'Best-Path'],
        'bodies': [
            'eBGP: AS diversi, TTL=1. iBGP: stesso AS, via loopback, no TTL limit, split-horizon.',
            'iBGP via loopback richiede: update-source Lo0 (stabilità) + next-hop-self (next-hop accessibile).',
            'Weight → LP → Local → AS-Path → Origin → MED. LP controlla uscita; AS-Path controlla entrata.',
        ],
    },
}

# ─── MOD-06 ───────────────────────────────────────────────────────────────────
MODULES['MOD-06'] = {
    'title': 'BGP Traffic Engineering',
    'area': 'AREA 2 — BGP', 'hours': '2h', 'codes': '1.11.c · 1.11.d · 1.11.e',
    'agenda': [
        'Default route da ISP: network vs default-originate',
        'Local Preference per controllare traffico uscente',
        'AS-Path Prepend per controllare traffico entrante',
        'BGP Community: comunicazione policy inter-AS',
        'MED: suggerire entry-point preferito al peer',
    ],
    'topology': {
        'title': 'Topologia MOD-06 — BGP Traffic Engineering',
        'caption': 'AS 65000 (Customer) dual-homed: R4↔R1 primario, R5↔R3 secondario. TE con LP e AS-Path Prepend.',
        'nodes': [
            ('R1\nAS 65001\nPrimario', 0.2, 0.2),
            ('R3\nAS 65001\nSecondario', 0.7, 0.2),
            ('R4\nAS 65000\nBorder Prim', 0.2, 0.7),
            ('R5\nAS 65000\nBorder Sec', 0.7, 0.7),
            ('R6\nAS 65000\nInternal', 0.5, 0.5),
        ],
    },
    'sections': [
        {
            'section': 'BGP Traffic Engineering',
            'subtitle': 'Controllare il flusso del traffico in e out',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Local Preference — Traffico Uscente',
                    'points': [
                        'Local Preference: attributo well-known discretionary, propagato solo in iBGP',
                        'Valore default: 100. Più alto = preferito per il traffico uscente dall\'AS',
                        'Usa route-map per impostare LP: neighbor X route-map SET-LP in',
                        'Scenario: link primario LP=200, link secondario LP=50 → tutto esce dal primario',
                        'Failover automatico: se LP=200 cade, LP=50 diventa best path',
                        'LP influenza solo il traffico uscente dall\'AS (OUTBOUND)',
                    ],
                    'key': 'Local Preference: alto = preferito per uscita. Configurato INBOUND sul border router.',
                },
                {
                    'type': 'teoria',
                    'title': 'AS-Path Prepend — Traffico Entrante',
                    'points': [
                        'AS-Path: lista di AS attraversati — più corto = preferito dalla selezione BGP standard',
                        'AS-Path Prepend: aggiunge ripetizioni del proprio AS nel path annunciato',
                        'Effetto: il path artificialmente più lungo viene sfavorito dai peer',
                        'Sintassi route-map: set as-path prepend <AS> <AS> <AS> (ripetizioni)',
                        'Configurato OUTBOUND sul border router che vuole sfavorire quel link',
                        'Influenza il traffico ENTRANTE nell\'AS (peer vede path più lungo)',
                    ],
                    'key': 'AS-Path Prepend: rendere un path più lungo per dirottare il traffico entrante sul link preferito.',
                },
                {
                    'type': 'teoria',
                    'title': 'BGP Community — Policy Inter-AS',
                    'points': [
                        'Community: tag 32-bit (AA:NN) applicato ai prefissi BGP — comunicazione di policy',
                        'Community well-known: NO_EXPORT (non uscire dall\'AS), NO_ADVERTISE (non annunciare)',
                        'Community custom: accordo bilaterale tra AS per comunicare preferenze',
                        'Uso tipico: ISP offre community 65001:100=preferenza alta, 65001:50=bassa',
                        'Il Customer tagga i prefissi con la community ISP per influenzare routing',
                        'send-community: obbligatorio nella neighbor config (default: community non trasmessa)',
                    ],
                    'key': 'Community = etichetta su un prefisso BGP per comunicare policy tra AS — send-community obbligatorio.',
                },
            ],
        },
    ],
    'config_section': {
        'slides': [
            {
                'type': 'config',
                'title': 'Local Preference + AS-Path Prepend',
                'lines': [
                    '! Local Pref sul link primario (R4)',
                    'route-map SET-LP-HIGH permit 10',
                    ' set local-preference 200',
                    'router bgp 65000',
                    ' neighbor 172.16.14.1 route-map SET-LP-HIGH in',
                    '!',
                    '! AS-Path Prepend sul link secondario (R5)',
                    'route-map PREPEND-SECONDARY permit 10',
                    ' set as-path prepend 65000 65000',
                    'router bgp 65000',
                    ' neighbor 172.16.35.1 route-map PREPEND-SECONDARY out',
                ],
                'device': 'R4/R5', 'hl': 2,
            },
        ],
    },
    'trouble': [
        ('LP non cambia best path', 'LP cambia ma serve clear ip bgp soft — oppure route-map non matchata'),
        ('AS-Path Prepend non ha effetto', 'send-community mancante o route-map non applicata outbound'),
        ('Community non propagata al peer', 'Manca neighbor X send-community — default: community non inviata'),
        ('Default route non ricevuta', 'neighbor X default-originate sul provider, oppure network 0.0.0.0 nel RIB'),
    ],
    'exam_tips': [
        'Local Preference (alto = preferito) controlla traffico USCENTE — configurato inbound sul border',
        'AS-Path Prepend (path lungo = sfavorito) controlla traffico ENTRANTE — configurato outbound',
        'Community: send-community deve essere esplicito nella neighbor config',
        'MED: più basso = preferito — suggerisce al peer il punto di ingresso preferito (solo tra stesso AS)',
    ],
    'exam_qa': [
        ('Come far uscire tutto il traffico dal link primario?',
         'Local Preference alta (es. 200) sul link primario con route-map in inbound.'),
        ('Come far entrare il traffico dal link primario?',
         'AS-Path Prepend sul link secondario (outbound) — path più lungo = sfavorito dal peer.'),
    ],
    'summary': {
        'labels': ['Local Preference', 'AS-Path Prepend', 'BGP Community'],
        'bodies': [
            'LP alto = uscita preferita. Configurato inbound sul border. Propagato in iBGP nell\'AS.',
            'Prepend = path lungo = sfavorito. Configurato outbound. Influenza traffico entrante.',
            'Community = etichetta policy. send-community obbligatorio. NO_EXPORT, NO_ADVERTISE built-in.',
        ],
    },
}

# ─── MOD-07 ───────────────────────────────────────────────────────────────────
MODULES['MOD-07'] = {
    'title': 'BGP Route Reflector & IPv6 BGP',
    'area': 'AREA 2 — BGP', 'hours': '1.5h', 'codes': '1.11.d',
    'agenda': [
        'Problema del full-mesh iBGP: scalabilità',
        'Route Reflector: cluster e RR client',
        'BGP Confederation: alternativa al RR',
        'MP-BGP: address-family per IPv6',
        'Exam Tips e riepilogo scaling BGP',
    ],
    'topology': {
        'title': 'Topologia MOD-07 — Route Reflector',
        'caption': 'RR centrale connesso a 4 RR-client. Senza RR: 10 sessioni iBGP. Con RR: 4 sessioni.',
        'nodes': [
            ('RR\nRoute Reflector', 0.5, 0.5),
            ('Client1', 0.15, 0.2),
            ('Client2', 0.85, 0.2),
            ('Client3', 0.15, 0.8),
            ('Client4', 0.85, 0.8),
        ],
    },
    'sections': [
        {
            'section': 'Route Reflector',
            'subtitle': 'Eliminare il full-mesh iBGP',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Il Problema del Full-Mesh iBGP',
                    'points': [
                        'iBGP split-horizon: un router non ri-annuncia prefix iBGP ad altri peer iBGP',
                        'Soluzione originale: full-mesh — ogni router ha sessione con tutti (N*(N-1)/2)',
                        'Con 100 router: 4950 sessioni BGP — non scalabile in grandi AS',
                        'Route Reflector (RR): può ri-annunciare prefix iBGP ai client — viola split-horizon in modo controllato',
                        'Confederation: divide l\'AS in sotto-AS (sub-AS) — eBGP interno tra sub-AS',
                        'RR è lo standard moderno; Confederation più raro, usata in grosse telco',
                    ],
                    'key': 'RR risolve il full-mesh iBGP: il RR ri-annuncia prefix ai client — da N*(N-1)/2 a N sessioni.',
                },
                {
                    'type': 'teoria',
                    'title': 'Route Reflector — Cluster e Regole',
                    'points': [
                        'RR cluster: uno o più RR + i loro client — identificato da cluster-id',
                        'Client → RR: RR riflette a tutti gli altri client e ai non-client iBGP',
                        'Non-client → RR: RR riflette solo ai client (non ad altri non-client)',
                        'Cluster-list: attributo aggiunto dal RR per prevenire loop tra RR multipli',
                        'Originator-ID: RR aggiunge l\'ID del router originante per prevenire loop',
                        'Configurazione: neighbor X route-reflector-client sull\'RR',
                    ],
                    'key': 'RR aggiunge cluster-list e originator-id per prevenire loop — trasparente ai client.',
                },
                {
                    'type': 'teoria',
                    'title': 'MP-BGP — Multi-Protocol BGP per IPv6',
                    'points': [
                        'MP-BGP (RFC 4760): estende BGP per portare NLRI di famiglie diverse da IPv4',
                        'address-family ipv6 unicast: abilita il trasporto di prefissi IPv6 in BGP',
                        'I session BGP stessi (TCP 179) — solo i NLRI cambiano address-family',
                        'Neighbor per IPv6: usa indirizzo IPv6 o loopback IPv4 con activate in af ipv6',
                        'next-hop per IPv6: il next-hop nella af ipv6 è un indirizzo IPv6',
                        'VPNv4, VPNv6, L2VPN: tutte le varianti MPLS usano MP-BGP',
                    ],
                    'key': 'MP-BGP = una sessione TCP porta NLRI di più famiglie. address-family seleziona il tipo.',
                },
            ],
        },
    ],
    'config_section': {
        'slides': [
            {
                'type': 'config',
                'title': 'Route Reflector — Configurazione',
                'lines': [
                    'router bgp 65000',
                    ' bgp router-id 10.0.0.1',
                    ' ! Sessioni verso i client',
                    ' neighbor 1.1.1.1 remote-as 65000',
                    ' neighbor 1.1.1.1 update-source Lo0',
                    ' neighbor 1.1.1.1 route-reflector-client   ! questo è client',
                    ' neighbor 2.2.2.2 remote-as 65000',
                    ' neighbor 2.2.2.2 update-source Lo0',
                    ' neighbor 2.2.2.2 route-reflector-client',
                    ' ! Non-client: sessione normale (no rr-client)',
                    ' neighbor 3.3.3.3 remote-as 65000',
                ],
                'device': 'RR', 'hl': 5,
            },
        ],
    },
    'trouble': [
        ('Loop BGP con RR multipli', 'Manca cluster-id configurato uguale su RR ridondanti dello stesso cluster'),
        ('Client non riceve prefissi da altri client', 'RR non ha route-reflector-client configurato su tutti i client'),
        ('MP-BGP IPv6: no prefix nella table', 'neighbor non activato in address-family ipv6 unicast'),
        ('Sessione RR non si forma', 'update-source mancante — RR deve raggiungere il loopback del client'),
    ],
    'exam_tips': [
        'RR evita il full-mesh: route-reflector-client sul RR — il client non sa di essere reflectato',
        'cluster-id su RR ridondanti evita loop — stessa cluster-id su tutti gli RR dello stesso cluster',
        'MP-BGP: address-family ipv6 unicast + neighbor X activate — stessa sessione TCP, diversi NLRI',
        'Confederation: sub-AS con eBGP interno — alternativa più complessa al RR',
    ],
    'exam_qa': [
        ('Perché il Route Reflector non rompe la regola split-horizon?',
         'RR aggiunge cluster-list e originator-id — i loop vengono rilevati e i prefix scartati.'),
        ('Differenza tra RR e Confederation?',
         'RR: singolo AS, client trasparenti, più semplice. Confederation: divide in sub-AS con eBGP interno.'),
    ],
    'summary': {
        'labels': ['Full-mesh → RR', 'Cluster + Loop Prev.', 'MP-BGP'],
        'bodies': [
            'RR risolve scalabilità iBGP: da N*(N-1)/2 a N sessioni. route-reflector-client sul RR.',
            'cluster-list + originator-id prevengono loop. RR ridondanti: stesso cluster-id.',
            'MP-BGP porta IPv6, VPNv4, L2VPN in BGP. address-family + neighbor activate.',
        ],
    },
}

# ─── MOD-08 ───────────────────────────────────────────────────────────────────
MODULES['MOD-08'] = {
    'title': 'Redistribuzione BGP↔OSPF & Prefix Filtering',
    'area': 'AREA 3 — ROUTE MANIPULATION', 'hours': '2h', 'codes': '1.3 · 1.4 · 1.5 · 3.2.d',
    'agenda': [
        'Prefix-list: filtrare prefissi per lunghezza',
        'Route-map: strumento universale di manipolazione',
        'Redistribuzione OSPF→BGP e BGP→OSPF',
        'Loop di redistribuzione e route tagging',
        'Route-map in redistribuzione vs BGP neighbor policy',
    ],
    'topology': {
        'title': 'Topologia MOD-08 — Redistribuzione Multi-AS',
        'caption': 'CORE: punto di redistribuzione BGP↔OSPF. WAN-A/WAN-B: dual-homed. LAN-A/LAN-B: solo OSPF.',
        'nodes': [
            ('ISP-A\nAS 100', 0.08, 0.3),
            ('WAN-A\nAS 65001', 0.25, 0.3),
            ('CORE\nAS 65000', 0.5, 0.5),
            ('WAN-B\nAS 65002', 0.75, 0.3),
            ('ISP-B\nAS 200', 0.92, 0.3),
            ('LAN-A\nOSPF', 0.35, 0.75),
            ('LAN-B\nOSPF', 0.65, 0.75),
        ],
    },
    'sections': [
        {
            'section': 'Prefix-List e Route-Map',
            'subtitle': 'Strumenti fondamentali di policy routing',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Prefix-List — Filtrare per Rete e Lunghezza',
                    'points': [
                        'Prefix-list: filtra prefissi IP per network e/o lunghezza maschera',
                        'Sintassi: ip prefix-list NAME [seq N] permit|deny <prefix/len> [ge X] [le Y]',
                        'ge (greater-equal): maschera ≥ ge. le (less-equal): maschera ≤ le',
                        'Esempio: 10.0.0.0/8 ge 24 le 32 → tutti gli host dentro 10.0.0.0/8 con /24-/32',
                        'Più efficiente di una ACL per il filtraggio di routing (valutazione O(1) con PATRICIA tree)',
                        'Implicit deny alla fine — come ACL: aggiungere permit any se necessario',
                    ],
                    'key': 'prefix-list ge/le: filtrare range di maschere in un unico statement. Più preciso delle ACL.',
                },
                {
                    'type': 'teoria',
                    'title': 'Route-Map — Strumento Universale',
                    'points': [
                        'Route-map: serie ordinata di clausole permit/deny con condizioni match e azioni set',
                        'match: ip address (ACL/prefix-list), ip next-hop, tag, metric, route-type...',
                        'set: metric, local-preference, as-path prepend, community, tag, weight...',
                        'Uso: redistribuzione (seleziona cosa redistribuire), BGP (policy neighbor), PBR',
                        'Clausola senza match → matcha tutto (catch-all)',
                        'Implicit deny finale: i prefix non matchati non vengono redistribuiti',
                    ],
                    'key': 'Route-map = IF (match) THEN (set). Implicit deny finale — aggiungere permit 999 se necessario.',
                },
            ],
        },
        {
            'section': 'Redistribuzione e Loop Prevention',
            'subtitle': 'Collegare BGP e OSPF senza loop',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Redistribuzione Bidirezionale BGP↔OSPF',
                    'points': [
                        'Redistribuzione OSPF→BGP: redistribute ospf <pid> [route-map FILTER] in router bgp',
                        'Redistribuzione BGP→OSPF: redistribute bgp <AS> subnets [route-map FILTER] in router ospf',
                        'subnets: obbligatorio per redistribuire prefissi con maschera variabile in OSPF',
                        'Problema loop: A redistribuisce BGP→OSPF; B vede la rotta OSPF e la ri-annuncia in BGP',
                        'Soluzione route tagging: set tag su redistribute BGP→OSPF; deny tag in redistribute OSPF→BGP',
                        'Tag è un marcatore 32-bit che viaggia con la rotta e può essere testato con match tag',
                    ],
                    'key': 'Loop redistribuzione BGP↔OSPF: prevenire con route tag. Marcare BGP→OSPF, bloccare OSPF→BGP.',
                },
            ],
        },
    ],
    'config_section': {
        'slides': [
            {
                'type': 'config',
                'title': 'Redistribuzione con Tag Anti-Loop',
                'lines': [
                    '! OSPF: accetta BGP ma marca, blocca tagged OSPF→BGP',
                    'route-map BGP-TO-OSPF permit 10',
                    ' set tag 65000               ! marca rotte BGP in OSPF',
                    'route-map OSPF-TO-BGP deny 10',
                    ' match tag 65000             ! blocca rotte già da BGP',
                    'route-map OSPF-TO-BGP permit 20',
                    '!',
                    'router ospf 1',
                    ' redistribute bgp 65000 subnets route-map BGP-TO-OSPF',
                    'router bgp 65000',
                    ' redistribute ospf 1 route-map OSPF-TO-BGP',
                ],
                'device': 'CORE', 'hl': 3,
            },
        ],
    },
    'trouble': [
        ('Loop di redistribuzione', 'Manca anti-loop con tag — rotte BGP→OSPF→BGP si amplificano'),
        ('Rotte non redistribuite in OSPF', 'Manca subnets nella redistribuzione OSPF — solo classful redistribuito senza'),
        ('Prefix-list blocca tutto', 'Implicit deny finale — verificare il deny/permit dell\'ultima clausola'),
        ('Route-map non matcha', 'match ip address usa ACL o prefix-list — verificare che il nome corrisponda'),
    ],
    'exam_tips': [
        'redistribute bgp X subnets: subnets obbligatorio per redistribuire prefissi VLSM in OSPF',
        'Loop BGP↔OSPF: set tag in redistribuzione + match tag in deny per prevenire rientro',
        'prefix-list ge/le: filtra range di maschere — più preciso ed efficiente di ACL estesa',
        'Route-map implicit deny: senza permit finale, tutto viene bloccato',
    ],
    'exam_qa': [
        ('Perché subnets è necessario in redistribute bgp into OSPF?',
         'Senza subnets, OSPF redistribuisce solo rotte classful (legacy). Con subnets: tutto incluso VLSM.'),
        ('Come prevenire loop in redistribuzione bidirezionale?',
         'Marcare rotte BGP→OSPF con set tag. Bloccare in OSPF→BGP con match tag + deny.'),
    ],
    'summary': {
        'labels': ['Prefix-List ge/le', 'Route-Map IF/THEN', 'Anti-Loop Tag'],
        'bodies': [
            'prefix-list filtra per network + lunghezza maschera. ge/le per range. Più efficiente di ACL.',
            'Route-map: match (condizione) + set (azione). Implicit deny finale.',
            'Loop BGP↔OSPF: set tag in BGP→OSPF. match tag + deny in OSPF→BGP.',
        ],
    },
}

# ─── MOD-09 ───────────────────────────────────────────────────────────────────
MODULES['MOD-09'] = {
    'title': 'PBR & Route Manipulation Avanzata',
    'area': 'AREA 3 — ROUTE MANIPULATION', 'hours': '2h', 'codes': '1.2 · 1.6 · 3.2.d',
    'agenda': [
        'Policy-Based Routing (PBR): routing per sorgente',
        'ip route-map e set ip next-hop verify-availability',
        'Administrative Distance: preferenza protocolli',
        'Floating static route: backup condizionale',
        'IP SLA + Object Tracking: rotte condizionali',
    ],
    'topology': {
        'title': 'Topologia MOD-09 — PBR Guest vs Produzione',
        'caption': 'LAN-A: rete prod (10.10.0.0) via WAN-A; rete guest (10.99.0.0) forzata via WAN-B con PBR su CORE.',
        'nodes': [
            ('ISP-A\nAS 100', 0.08, 0.3),
            ('WAN-A\nAS 65001', 0.25, 0.3),
            ('CORE\nPBR qui', 0.5, 0.5),
            ('WAN-B\nAS 65002', 0.75, 0.3),
            ('ISP-B\nAS 200', 0.92, 0.3),
            ('LAN-A\nprod+guest', 0.3, 0.75),
        ],
    },
    'sections': [
        {
            'section': 'Policy-Based Routing',
            'subtitle': 'Routing basato su policy anziché destinazione',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'PBR — Routing per Sorgente (e non solo)',
                    'points': [
                        'PBR bypassa la routing table normale per il traffico matchato',
                        'Match: ip address (ACL che seleziona il traffico) — per IP sorgente, destinazione, DSCP...',
                        'set ip next-hop: forza il next-hop per il traffico matchato',
                        'set ip next-hop verify-availability: verifica che il next-hop sia raggiungibile via IP SLA',
                        'Applicato INBOUND sull\'interfaccia verso la sorgente: ip policy route-map NAME',
                        'Uso tipico: traffico guest→ISP-B, traffico prod→ISP-A su stesso router',
                    ],
                    'key': 'PBR = routing basato su attributi arbitrari del pacchetto — applicato inbound sull\'ingresso traffico.',
                },
            ],
        },
        {
            'section': 'Administrative Distance e IP SLA',
            'subtitle': 'Preferenza protocolli e rotte condizionali',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Administrative Distance — Preferenza Protocolli',
                    'points': [
                        'AD: valore 0-255, più basso = più fidato — usato per scegliere tra protocolli diversi',
                        'AD default: Connected=0, Static=1, EIGRP summary=5, OSPF=110, IS-IS=115, RIP=120',
                        'Floating static route: static con AD > protocollo dinamico (es. AD=200 vs OSPF=110)',
                        'Quando OSPF è attivo → OSPF installa la rotta. OSPF cade → static emerge',
                        'Modifica AD: ip route 0.0.0.0 0.0.0.0 10.0.0.1 200 (AD 200)',
                        'Utile per backup con failover automatico senza dover monitorare protocolli',
                    ],
                    'key': 'Floating static: AD più alto del protocollo dinamico — appare nel RIB solo se il protocollo cade.',
                },
                {
                    'type': 'teoria',
                    'title': 'IP SLA + Object Tracking',
                    'points': [
                        'IP SLA: sonda proattiva (ICMP echo, UDP jitter, HTTP...) verso un target',
                        'Risultato: UP se la sonda ha successo, DOWN se fallisce (configurable thresholds)',
                        'Object Tracking: traccia stato di un SLA, interfaccia o route',
                        'Rotta condizionale: ip route 0.0.0.0 ... track 1 — installata solo se track 1 = UP',
                        'Combinazione con PBR: set ip next-hop verify-availability usa track object implicito',
                        'Frequenza sonda: ip sla schedule — configurare probe frequency e react threshold',
                    ],
                    'key': 'IP SLA + track: rotta installata condizionalmente. PBR verify-availability: next-hop testato live.',
                },
            ],
        },
    ],
    'config_section': {
        'slides': [
            {
                'type': 'config',
                'title': 'PBR Guest Traffic',
                'lines': [
                    'ip access-list extended GUEST-TRAFFIC',
                    ' permit ip 10.99.0.0 0.0.0.255 any',
                    '!',
                    'route-map PBR-GUEST permit 10',
                    ' match ip address GUEST-TRAFFIC',
                    ' set ip next-hop verify-availability 10.0.23.1 1 track 10',
                    '!',
                    '! IP SLA per verificare WAN-B',
                    'ip sla 10',
                    ' icmp-echo 10.0.23.1',
                    'ip sla schedule 10 life forever start-time now',
                    'track 10 ip sla 10 reachability',
                ],
                'device': 'CORE', 'hl': 5,
            },
        ],
    },
    'trouble': [
        ('PBR non ha effetto', 'ip policy route-map non applicato inbound sull\'interfaccia corretta'),
        ('next-hop verify-availability fallisce', 'Track non configurato o SLA non schedulato — verify-availability richiede track object'),
        ('Floating static non emerge', 'AD della static NON è maggiore del protocollo dinamico — verificare valore AD'),
        ('IP SLA sempre DOWN', 'Target non raggiungibile o frequenza probe troppo bassa — verificare routing e sla schedule'),
    ],
    'exam_tips': [
        'PBR applicato INBOUND: ip policy route-map sull\'interfaccia di ingresso del traffico',
        'Floating static: AD > protocollo dinamico. OSPF=110 → static AD=150 per backup',
        'IP SLA: proattivo — rileva il fallimento PRIMA che OSPF/BGP convergano',
        'verify-availability nel PBR usa implicitamente un track object per testare il next-hop',
    ],
    'exam_qa': [
        ('Come forzare traffico da una subnet specifica verso un next-hop diverso dalla routing table?',
         'PBR: ACL seleziona la subnet + route-map con set ip next-hop + ip policy inbound sull\'interfaccia.'),
        ('Quando una floating static route viene installata nel RIB?',
         'Solo quando il protocollo con AD inferiore non ha più la rotta — floating static emerge automaticamente.'),
    ],
    'summary': {
        'labels': ['PBR Inbound', 'Floating Static', 'IP SLA + Track'],
        'bodies': [
            'PBR matcha ACL + set next-hop. Applicato inbound. Bypassa routing table per traffico selezionato.',
            'AD floating > protocollo: backup automatico. OSPF cade → static emerge. Zero configurazione manuale.',
            'IP SLA proba il next-hop. track object condiziona rotte e PBR. Failover pre-routing.',
        ],
    },
}


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def generate_module(mod_id):
    content = MODULES.get(mod_id)
    if not content:
        print(f"  [SKIP] {mod_id} — contenuto non ancora definito")
        return 0
    out_dir = mod_id
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, f"{mod_id}_slide.pptx")
    prs = make_deck(mod_id, content)
    prs.save(out_path)
    n = len(prs.slides)
    print(f"  [OK] {out_path}  ({n} slide)")
    return n


def main():
    target = sys.argv[1] if len(sys.argv) > 1 else None
    all_ids = [f"MOD-{i:02d}" for i in range(1, 36)]
    to_gen = [target] if target else all_ids

    total = 0
    for mod_id in to_gen:
        total += generate_module(mod_id)

    print(f"\n[DONE] Totale slide generate: {total}")


# ─── MOD-10 ───────────────────────────────────────────────────────────────────
MODULES['MOD-10'] = {
    'title': 'MPLS LDP & Fondamenta',
    'area': 'AREA 4 — MPLS', 'hours': '2h', 'codes': '2.1',
    'agenda': [
        'MPLS: cos\'è e perché esiste',
        'FEC, Label, LIB e LFIB',
        'LDP: distribuzione e sessioni',
        'PHP: Penultimate Hop Popping',
        'Configurazione e verifica LDP su IOS',
    ],
    'topology': {
        'title': 'Topologia MOD-10 — Backbone MPLS',
        'caption': 'PE1-P1-P2-PE2: backbone MPLS AS 65000. OSPF area 0 pre-configurato. LDP da abilitare su PE1 e P1.',
        'nodes': [
            ('PE1\n1.1.1.1\nLDP da conf', 0.12, 0.5),
            ('P1\n3.3.3.3\nLDP da conf', 0.38, 0.5),
            ('P2\n4.4.4.4\nLDP pre-conf', 0.62, 0.5),
            ('PE2\n2.2.2.2\nLDP pre-conf', 0.88, 0.5),
        ],
    },
    'sections': [
        {
            'section': 'MPLS — Fondamenta',
            'subtitle': 'Multiprotocol Label Switching: il perché',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'MPLS — Cos\'è e Perché Esiste',
                    'points': [
                        'MPLS: instradamento basato su label 20-bit anziché lookup IP — più veloce su hardware legacy',
                        'Label aggiunta tra header L2 e header IP (shim header) — non modifica IP',
                        'Motivazione originale: velocità (ora superata da ASIC IP). Uso attuale: VPN L3/L2, TE, QoS',
                        'FEC (Forwarding Equivalence Class): gruppo di pacchetti trattati allo stesso modo',
                        'Label Stack: MPLS supporta stack di label (LIFO) — es. VPN + TE tunnel',
                        'Bit speciali: S=1 (bottom of stack), TTL (loop prevention), TC (traffic class)',
                    ],
                    'key': 'MPLS oggi: non velocità, ma VPN (L3VPN, L2VPN), Traffic Engineering e gestione QoS.',
                },
                {
                    'type': 'teoria',
                    'title': 'LIB, LFIB e PHP',
                    'points': [
                        'LIB (Label Information Base): database di tutte le label imparate via LDP',
                        'LFIB (Label Forwarding Information Base): tabella di forwarding label → azione (swap/pop/push)',
                        'Azioni label: PUSH (aggiunge label), SWAP (sostituisce), POP (rimuove)',
                        'PHP (Penultimate Hop Popping): il penultimo router rimuove la label — il PE riceve IP puro',
                        'Implicit null (label 3): LDP annuncia label 3 per i propri prefissi → PHP attivato',
                        'Explicit null (label 0): PHP disabilitato — utile per preservare QoS marking nella label',
                    ],
                    'key': 'PHP: l\'ultimo hop MPLS (penultimate) fa pop — il PE riceve IP direttamente, non label.',
                },
            ],
        },
        {
            'section': 'LDP — Label Distribution Protocol',
            'subtitle': 'Come i router si scambiano le label',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'LDP — Sessioni e Funzionamento',
                    'points': [
                        'LDP: protocollo per distribuire binding label↔FEC tra router MPLS',
                        'Discovery: Hello multicast UDP 224.0.0.2:646 — scopre i vicini LDP',
                        'Sessione: TCP 646 tra Router-ID (loopback) — stabile, non dipende dal link',
                        'Ogni router genera label locale per ogni prefisso nel RIB → annuncia via LDP',
                        'Liberal label retention: conserva label da tutti i peer (anche non-best-path)',
                        'Router-ID LDP: deve essere lo stesso usato da OSPF — configurare esplicitamente',
                    ],
                    'key': 'LDP usa loopback per la sessione TCP — stessa regola di BGP: update-source Lo0.',
                },
            ],
        },
    ],
    'config_section': {
        'slides': [
            {
                'type': 'config',
                'title': 'LDP — Abilitazione su IOS',
                'lines': [
                    'mpls ldp router-id Loopback0 force  ! Router-ID LDP = loopback',
                    '!',
                    'interface Ethernet0/0.13',
                    ' mpls ip                            ! abilita LDP sull\'interfaccia',
                    '!',
                    'interface Ethernet0/0.34',
                    ' mpls ip',
                ],
                'device': 'P1', 'hl': 0,
            },
            {
                'type': 'verifica',
                'title': 'Verifica LDP e MPLS Forwarding',
                'lines': [
                    'P1# show mpls ldp neighbor',
                    'Peer LDP Ident: 1.1.1.1:0; Local LDP Ident 3.3.3.3:0',
                    '    TCP connection: 1.1.1.1.42316 - 3.3.3.3.646',
                    '    State: Oper; Msgs sent/rcvd: 52/51',
                    '',
                    'P1# show mpls forwarding-table',
                    'Local  Out   Prefix         Bytes Label  Outgoing',
                    '16     Pop   1.1.1.1/32     0      Et0/0.13',
                    '17     18    2.2.2.2/32     0      Et0/0.34',
                ],
                'hl': 6,
            },
        ],
    },
    'trouble': [
        ('Sessione LDP non si forma', 'Router-ID LDP diverso da OSPF — mpls ldp router-id Loopback0 force'),
        ('LFIB vuota dopo mpls ip', 'OSPF non ha rotte verso il peer loopback — verificare reachabilità loopback'),
        ('PHP non avviene', 'PHP richiede implicit-null (label 3) — verificare show mpls ldp bindings'),
        ('Sessione LDP flap', 'Instabilità OSPF — il loopback LDP-ID deve essere in OSPF passivo e stabile'),
    ],
    'exam_tips': [
        'mpls ldp router-id Loopback0 force: obbligatorio per stabilità sessione LDP',
        'PHP: label 3 (implicit-null) sul penultimo hop — il PE riceve IP puro',
        'LFIB: show mpls forwarding-table — Pop = PHP in atto, Swap = transito, Aggregate = destinazione locale',
        'LDP usa TCP 646 tra loopback — stessa regola di BGP: loopback deve essere raggiungibile',
    ],
    'exam_qa': [
        ('Cos\'è il Penultimate Hop Popping (PHP)?',
         'Il penultimo router MPLS rimuove la label prima di consegnare al PE — il PE fa lookup IP normale.'),
        ('Differenza tra LIB e LFIB?',
         'LIB: database completo label imparate. LFIB: solo label usate per il forwarding (best-path).'),
    ],
    'summary': {
        'labels': ['FEC + Label Stack', 'LDP via Loopback', 'PHP = Label 3'],
        'bodies': [
            'FEC: gruppo pacchetti con stesso trattamento. Label 20-bit tra L2 e IP. Stack per VPN+TE.',
            'LDP scopre peer via multicast, sessione TCP tra loopback. Router-ID = loopback esplicito.',
            'PHP: penultimo router fa pop (label 3). PE riceve IP puro — lookup normale.',
        ],
    },
}

# ─── MOD-11 ───────────────────────────────────────────────────────────────────
MODULES['MOD-11'] = {
    'title': 'MPLS L3VPN',
    'area': 'AREA 4 — MPLS', 'hours': '2h', 'codes': '2.2',
    'agenda': [
        'VRF: Virtual Routing and Forwarding',
        'Route Distinguisher (RD) e Route Target (RT)',
        'MP-BGP VPNv4: trasporto delle rotte VPN',
        'Label stack: VPN label + transport label',
        'Peering PE-CE: eBGP nel VRF',
    ],
    'topology': {
        'title': 'Topologia MOD-11 — MPLS L3VPN',
        'caption': 'PE1↔PE2: MP-BGP VPNv4 iBGP. CE1/CE2 (AS 65001/65002) in VRF CUST_A. RD 65000:100.',
        'nodes': [
            ('CE1\nAS 65001', 0.08, 0.5),
            ('PE1\n1.1.1.1\nVRF CUST_A', 0.28, 0.5),
            ('P1/P2\nCore MPLS', 0.55, 0.5),
            ('PE2\n2.2.2.2\nVRF CUST_A', 0.78, 0.5),
            ('CE2\nAS 65002', 0.95, 0.5),
        ],
    },
    'sections': [
        {
            'section': 'VRF e MP-BGP VPNv4',
            'subtitle': 'Separazione di routing e trasporto etichette',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'VRF — Virtual Routing and Forwarding',
                    'points': [
                        'VRF: istanza di routing separata — routing table, CEF table, forwarding indipendente',
                        'Permette overlapping address space: due clienti possono usare lo stesso range IP',
                        'vrf forwarding su interfaccia: assegna l\'interfaccia a un VRF specifico',
                        'Global table ≠ VRF table: le rotte non si "vedono" a vicenda senza leaking esplicito',
                        'VRF-Lite: VRF senza MPLS — solo per separazione locale (es. tenant su un router)',
                        'VRF con MPLS: aggiunge RD + RT + MP-BGP per trasporto inter-PE',
                    ],
                    'key': 'VRF = routing table virtuale. Overlapping IP OK. VRF-Lite = separazione locale, no MPLS.',
                },
                {
                    'type': 'teoria',
                    'title': 'RD, RT e MP-BGP VPNv4',
                    'points': [
                        'RD (Route Distinguisher): 64-bit prefisso alle rotte VPN per renderle uniche in BGP',
                        'RT (Route Target): community BGP extended — definisce import/export tra VRF',
                        'VPNv4 prefix = RD + IPv4 (96 bit totali) — trasportato via MP-BGP address-family vpnv4',
                        'Label VPN: assegnata dal PE egress — indica il VRF di destinazione all\'uscita',
                        'Label stack: [transport label (LDP)] + [VPN label (MP-BGP)] — 2 label in transit',
                        'PE1 → PE2: push VPN label + push transport label. P-router: swap transport. PE2: pop entrambe',
                    ],
                    'key': 'RD rende le rotte uniche. RT decide dove importarle. VPN label indica il VRF al PE egress.',
                },
            ],
        },
    ],
    'config_section': {
        'slides': [
            {
                'type': 'config',
                'title': 'L3VPN — VRF + MP-BGP + CE-PE',
                'lines': [
                    'vrf definition CUST_A',
                    ' rd 65000:100',
                    ' route-target export 65000:100',
                    ' route-target import 65000:100',
                    ' address-family ipv4',
                    '!',
                    'interface Ethernet0/0.11',
                    ' vrf forwarding CUST_A',
                    ' ip address 192.168.1.2 255.255.255.252',
                    '!',
                    'router bgp 65000',
                    ' address-family vpnv4',
                    '  neighbor 2.2.2.2 activate',
                    '  neighbor 2.2.2.2 send-community extended',
                ],
                'device': 'PE1', 'hl': 1,
            },
        ],
    },
    'trouble': [
        ('VRF non importa le rotte dal PE remoto', 'RT import/export non corrispondono tra PE1 e PE2'),
        ('CE non raggiunge CE remoto', 'MP-BGP vpnv4 non attivato o send-community extended mancante'),
        ('Label stack errato', 'LDP non attivo sulla backbone — PHP non funziona correttamente'),
        ('Overlapping IP non funziona', 'vrf forwarding non applicato all\'interfaccia CE-PE'),
    ],
    'exam_tips': [
        'RD rende il prefisso VPN unico in BGP. RT decide import/export tra VRF su PE diversi',
        'send-community extended obbligatorio per trasmettere RT con i prefix VPNv4',
        'Label stack L3VPN: [transport LDP] + [VPN label] — 2 etichette in transito',
        'VRF-Lite: VRF senza MPLS per separazione locale. RD/RT non necessari.',
    ],
    'exam_qa': [
        ('Differenza tra RD e RT?',
         'RD rende univoco il prefix VPN in BGP (solo identificatore). RT (extended community) controlla dove importare/esportare.'),
        ('Quante label ci sono in un pacchetto L3VPN in transito su un P-router?',
         'Due: transport label (LDP, swappata dal P) + VPN label (MP-BGP, non toccata dal P).'),
    ],
    'summary': {
        'labels': ['VRF = Routing Table', 'RD Unico + RT Policy', 'Label Stack 2 livelli'],
        'bodies': [
            'VRF: istanza routing separata. Overlapping IP. vrf forwarding sull\'interfaccia CE-PE.',
            'RD = unicità in BGP. RT (extended community) = import/export policy tra VRF.',
            'Transport label (LDP) + VPN label (MP-BGP). P-router tocca solo transport.',
        ],
    },
}

# ─── MOD-12 ───────────────────────────────────────────────────────────────────
MODULES['MOD-12'] = {
    'title': 'MPLS L2VPN (xconnect / AToM)',
    'area': 'AREA 4 — MPLS', 'hours': '1.5h', 'codes': '2.2',
    'agenda': [
        'L2VPN vs L3VPN: quando usare ciascuno',
        'Pseudowire e Attachment Circuit',
        'xconnect: configurazione AToM su IOS',
        'Verifica e troubleshooting pseudowire',
    ],
    'topology': {
        'title': 'Topologia MOD-12 — MPLS L2VPN AToM',
        'caption': 'CE1 (VLAN 101) ↔ PE1 xconnect ↔ PE2 xconnect ↔ CE2 (VLAN 202). VC-ID 101. Backbone LDP pre-conf.',
        'nodes': [
            ('CE1\n172.16.0.1', 0.08, 0.5),
            ('PE1\n1.1.1.1\nxconnect', 0.3, 0.5),
            ('P1-P2\nCore', 0.55, 0.5),
            ('PE2\n2.2.2.2\nxconnect', 0.78, 0.5),
            ('CE2\n172.16.0.2', 0.95, 0.5),
        ],
    },
    'sections': [
        {
            'section': 'L2VPN AToM',
            'subtitle': 'Trasportare L2 su backbone MPLS',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'L2VPN vs L3VPN',
                    'points': [
                        'L3VPN: PE fa routing — i CE sono in VRF, traffico IP instradato dal provider',
                        'L2VPN: PE fa bridging/switching — il CE "vede" un link L2 diretto all\'altro CE',
                        'Pseudowire: tunnel virtuale L2 tra due PE — trasparente al traffico CE',
                        'AToM (Any Transport over MPLS): standard Cisco per L2VPN su MPLS',
                        'Attachment Circuit (AC): il link fisico/VLAN tra CE e PE — ingresso del pseudowire',
                        'VC-ID: identificatore del circuito virtuale — deve essere uguale su entrambi i PE',
                    ],
                    'key': 'L2VPN: il provider trasporta frame L2 — i CE pensano di essere sullo stesso segmento.',
                },
                {
                    'type': 'config',
                    'title': 'xconnect — Configurazione AToM',
                    'lines': [
                        'interface Ethernet0/0.101',
                        ' encapsulation dot1Q 101     ! attachment circuit',
                        ' xconnect 2.2.2.2 101 encapsulation mpls',
                        '                 ^peer-Lo0 ^VC-ID',
                        '! Su PE2 (simmetrico):',
                        'interface Ethernet0/0.202',
                        ' encapsulation dot1Q 202     ! VLAN diversa, VC-ID uguale',
                        ' xconnect 1.1.1.1 101 encapsulation mpls',
                    ],
                    'device': 'PE1', 'hl': 2,
                },
            ],
        },
    ],
    'config_section': {
        'slides': [
            {
                'type': 'verifica',
                'title': 'Verifica Pseudowire AToM',
                'lines': [
                    'PE1# show xconnect all',
                    'Legend: XC ST=Xconnect State  S1=Segment1 State  S2=Segment2 State',
                    'ST  Segment-1         S1  Segment-2                    S2',
                    'UP  Et0/0.101:101     UP  MPLS 2.2.2.2:101            UP',
                    '',
                    'PE1# show mpls l2transport vc',
                    'Local intf  St  Dest address  VC id  S/R  Pkts',
                    'Et0/0.101   UP  2.2.2.2       101    UP   1024',
                ],
                'hl': 3,
            },
        ],
    },
    'trouble': [
        ('Pseudowire DOWN', 'VC-ID non corrispondente tra PE1 e PE2 — deve essere identico'),
        ('xconnect non accettato', 'LDP non attivo sulla backbone o PE loopback non raggiungibile'),
        ('CE1 non pinga CE2', 'IP stesso subnet ma su VLAN diverse — controllare ARP e tag VLAN CE'),
        ('Pseudowire UP ma traffico 0', 'Interfaccia CE shutdown o VLAN non permessa sul trunk verso CE'),
    ],
    'exam_tips': [
        'VC-ID deve essere identico su entrambi i PE — è l\'identificatore del circuito virtuale',
        'xconnect usa LDP backbone per trasportare frame L2 — LDP deve essere attivo',
        'L2VPN: PE non fa routing — trasparente al traffico CE. CE vede link L2 diretto',
        'VLAN CE lato 1 ≠ VLAN CE lato 2 — il PE fa retagging tramite xconnect',
    ],
    'exam_qa': [
        ('Cos\'è un Attachment Circuit in L2VPN?',
         'Il link fisico/logico (VLAN) tra il CE e il PE — è l\'ingresso del pseudowire.'),
        ('Perché i VC-ID devono corrispondere su entrambi i PE?',
         'Identificano il circuito virtuale end-to-end — PE1 cerca il peer con lo stesso VC-ID.'),
    ],
    'summary': {
        'labels': ['Pseudowire L2', 'xconnect VC-ID', 'AToM su MPLS'],
        'bodies': [
            'L2VPN: tunnel virtuale L2 tra PE. CE vede link diretto. PE non fa routing.',
            'xconnect: 3 parametri — peer loopback, VC-ID (uguale su entrambi), encapsulation mpls.',
            'AToM trasporta Ethernet/VLAN su backbone MPLS. VLAN CE può essere diversa da lato a lato.',
        ],
    },
}

# ─── MOD-13 ───────────────────────────────────────────────────────────────────
MODULES['MOD-13'] = {
    'title': 'EtherChannel LACP',
    'area': 'Layer 2 Technologies', 'hours': '1h', 'codes': '3.1.b',
    'agenda': [
        'EtherChannel: aggregazione di link fisici',
        'LACP vs PAgP vs Static: quando usare quale',
        'Requisiti di compatibilità per il bundle',
        'Load balancing: algoritmi e configurazione',
        'Verifica e troubleshooting EtherChannel',
    ],
    'topology': {
        'title': 'Topologia MOD-13 — EtherChannel LACP',
        'caption': 'SW1↔SW2: Po1 LACP (e0/2+e0/3). R1 uplink: VLAN 100 (SW1) + VLAN 200 (SW2). HSRP su SVI.',
        'nodes': [
            ('R1\nRouter', 0.5, 0.15),
            ('SW1\nHSRP Active V10', 0.25, 0.5),
            ('SW2\nHSRP Active V20', 0.75, 0.5),
            ('PC1\nVLAN 10', 0.15, 0.85),
            ('PC2\nVLAN 20', 0.85, 0.85),
        ],
    },
    'sections': [
        {
            'section': 'EtherChannel LACP',
            'subtitle': 'Aggregazione link per banda e ridondanza',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'EtherChannel — Fondamenta',
                    'points': [
                        'EtherChannel: aggrega 2-8 link fisici in un singolo bundle logico (Port-Channel)',
                        'Vantaggi: banda aggregata + ridondanza automatica (link failure non causa down)',
                        'STP vede il Port-Channel come singola porta — evita loop senza bloccare link',
                        'LACP (802.3ad): standard IEEE — negoziazione automatica, consigliato',
                        'PAgP (Port Aggregation Protocol): proprietario Cisco — legacy',
                        'Static (on): nessuna negoziazione — rischioso (misconfiguration non rilevata)',
                    ],
                    'key': 'LACP = 802.3ad standard. STP vede un solo link logico — usa Po1 anziché e0/2+e0/3.',
                },
                {
                    'type': 'teoria',
                    'title': 'LACP — Requisiti e Load Balancing',
                    'points': [
                        'Requisiti bundle: stessa velocità, duplex, VLAN allowed, native VLAN, modalità trunk/access',
                        'LACP mode: active/active (negozia) oppure active/passive (uno negozia, uno aspetta)',
                        'Algoritmo load balancing: src-mac, dst-mac, src-dst-mac, src-ip, dst-ip, src-dst-ip',
                        'Hash deterministico: stesso flusso usa sempre lo stesso link fisico — no reordering',
                        'Scegliere il metodo in base al traffico: src-dst-ip per traffico IP generico',
                        'port-channel load-balance src-dst-ip: configurato globalmente sullo switch',
                    ],
                    'key': 'LACP active/active: entrambi i lati negoziano. Stessi parametri L2 su tutti i link del bundle.',
                },
            ],
        },
    ],
    'config_section': {
        'slides': [
            {
                'type': 'config',
                'title': 'EtherChannel LACP — Configurazione',
                'lines': [
                    'interface range Ethernet0/2 - 3',
                    ' switchport trunk encapsulation dot1q',
                    ' switchport mode trunk',
                    ' switchport trunk allowed vlan 10,20',
                    ' channel-group 1 mode active   ! LACP active',
                    '!',
                    'interface Port-channel 1',
                    ' switchport trunk encapsulation dot1q',
                    ' switchport mode trunk',
                    ' switchport trunk allowed vlan 10,20',
                    '!',
                    'port-channel load-balance src-dst-ip',
                ],
                'device': 'SW1', 'hl': 4,
            },
            {
                'type': 'verifica',
                'title': 'Verifica EtherChannel',
                'lines': [
                    'SW1# show etherchannel summary',
                    'Flags: D-down  P-bundled  s-suspended  I-stand-alone',
                    '       H-Hot-standby  U-in-use  f-failed-alloc',
                    'Number of channel-groups in use: 1',
                    'Group  Port-channel  Protocol   Ports',
                    '1      Po1(SU)       LACP       Et0/2(P) Et0/3(P)',
                    '',
                    'SW1# show lacp neighbor',
                    'Channel group 1 neighbors — Partner ports: Et0/2, Et0/3',
                ],
                'hl': 5,
            },
        ],
    },
    'trouble': [
        ('Port in stato "s" (suspended)', 'Mismatch VLAN allowed o native VLAN tra i membri del bundle'),
        ('Port-Channel non si forma', 'Mode incompatibile: active↔on non funziona — usare active↔active o active↔passive'),
        ('Un link non si aggiunge al bundle', 'Velocità o duplex diverso — tutti i link devono essere identici'),
        ('Traffic sbilanciato su un solo link', 'Algoritmo load-balance non adatto al traffico — cambiare con port-channel load-balance'),
    ],
    'exam_tips': [
        'LACP (802.3ad) = standard. PAgP = Cisco. Static (on) = rischioso senza negoziazione',
        'active/active o active/passive per LACP. passive/passive = non funziona',
        'Tutti i link del bundle: stessa velocità, duplex, VLAN allowed, native VLAN, modo trunk',
        'STP vede Po1 come singola porta — eliminazione loop automatica',
    ],
    'exam_qa': [
        ('LACP mode active vs passive?',
         'Active: invia LACP PDU e negozia. Passive: risponde solo se riceve LACP. Entrambi passive = no bundle.'),
        ('Perché EtherChannel non causa loop STP?',
         'STP vede Port-Channel come singola interfaccia logica — non ci sono path multipli da considerare.'),
    ],
    'summary': {
        'labels': ['LACP 802.3ad', 'Bundle Requirements', 'STP + Load Balance'],
        'bodies': [
            'LACP standard IEEE. active/active consigliato. Static rischioso senza feedback.',
            'Tutti i link uguali: velocità, duplex, VLAN, native, mode. Uno sbagliato = suspended.',
            'STP vede Po1 come un link. Load balance: src-dst-ip per traffico generico.',
        ],
    },
}

# ─── MOD-14 ───────────────────────────────────────────────────────────────────
MODULES['MOD-14'] = {
    'title': 'Spanning Tree Protocol',
    'area': 'Layer 2 Technologies', 'hours': '1h', 'codes': '3.1.c',
    'agenda': [
        'STP 802.1D: root election e port roles',
        'RSTP 802.1w: convergenza rapida',
        'MST 802.1s: istanze multiple',
        'PortFast e BPDU Guard: protezione edge',
        'Root Guard: protezione root bridge',
    ],
    'topology': {
        'title': 'Topologia MOD-14 — STP con Po1',
        'caption': 'SW1: Root VLAN10 (prio 4096), Secondary VLAN20 (8192). SW2: Root VLAN20. Po1 LACP da MOD-13.',
        'nodes': [
            ('R1\nUpstream', 0.5, 0.1),
            ('SW1\nRoot VLAN10\nprio 4096', 0.25, 0.5),
            ('SW2\nRoot VLAN20\nprio 4096', 0.75, 0.5),
            ('PC1\nVLAN10\nPortFast', 0.15, 0.85),
            ('PC2\nVLAN20\nPortFast', 0.85, 0.85),
        ],
    },
    'sections': [
        {
            'section': 'STP — Fondamenta',
            'subtitle': '802.1D, RSTP e MST',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'STP 802.1D — Root Election e Port Roles',
                    'points': [
                        'Root Bridge: switch con Bridge Priority più bassa (default 32768 + VLAN ID)',
                        'Tie-break: MAC address più bassa. Consigliato: configurare priority esplicitamente',
                        'Port roles: Root Port (verso root), Designated Port (sul segmento), Non-designated (bloccata)',
                        'Stati: Blocking → Listening → Learning → Forwarding (30s totali con delay 15s+15s)',
                        'RSTP 802.1w: convergenza in secondi anziché minuti — port roles Alternate e Backup',
                        'MST 802.1s: mappa VLAN a istanze STP — riduce numero di istanze attive',
                    ],
                    'key': 'RSTP = 802.1w. Stessa struttura root election, convergenza molto più rapida. Oggi lo standard.',
                },
                {
                    'type': 'teoria',
                    'title': 'PortFast, BPDU Guard e Root Guard',
                    'points': [
                        'PortFast: porta va direttamente in Forwarding (bypassa L+L) — solo su porte host/access',
                        'BPDU Guard: porta PortFast riceve BPDU → err-disabled — previene switch non autorizzati',
                        'Root Guard: porta non può diventare Root Port — protegge il root bridge da cambio',
                        'BPDU Filter: sopprime invio/ricezione BPDU — usare con cautela (non sicuro come Guard)',
                        'Loop Guard: porta in Forwarding smette di ricevere BPDU → blocking (non forward)',
                        'Err-disabled recovery: errdisable recovery cause bpduguard interval 300',
                    ],
                    'key': 'PortFast + BPDU Guard: standard su tutte le porte host. Root Guard: su porte trunk verso aggregazione.',
                },
            ],
        },
    ],
    'config_section': {
        'slides': [
            {
                'type': 'config',
                'title': 'STP Root + PortFast + Guards',
                'lines': [
                    '! Root primario VLAN 10, secondario VLAN 20',
                    'spanning-tree vlan 10 priority 4096',
                    'spanning-tree vlan 20 priority 8192',
                    '!',
                    '! Porte host: PortFast + BPDU Guard',
                    'interface Ethernet1/0',
                    ' spanning-tree portfast',
                    ' spanning-tree bpduguard enable',
                    '!',
                    '! Porta trunk verso switch aggregazione: Root Guard',
                    'interface Port-channel1',
                    ' spanning-tree guard root',
                ],
                'device': 'SW1', 'hl': 1,
            },
        ],
    },
    'trouble': [
        ('Porta in err-disabled dopo PortFast', 'BPDU Guard attivato — switch connesso invia BPDU. shutdown/no shutdown + rimuovere switch'),
        ('Root bridge non è quello desiderato', 'Priority non configurata — switch con MAC più bassa vince. Configurare priority esplicitamente'),
        ('MST region mismatch', 'Configurazione MST deve essere identica su tutti gli switch della regione (name, revision, mapping)'),
        ('Convergenza STP lenta', 'Usare RSTP (spanning-tree mode rapid-pvst) invece di 802.1D legacy'),
    ],
    'exam_tips': [
        'PortFast solo su porte host/access — mai su trunk o uplink verso switch',
        'BPDU Guard mette la porta in err-disabled se riceve BPDU — switch non autorizzato rilevato',
        'Root Guard: la porta non può diventare Root Port — protegge il root bridge da cambio non desiderato',
        'RSTP 802.1w: Alternate Port = backup del Root Port. Backup Port = backup Designated sullo stesso segmento',
    ],
    'exam_qa': [
        ('Differenza tra BPDU Guard e Root Guard?',
         'BPDU Guard: porta in err-disabled se riceve BPDU. Root Guard: porta bloccata se riceverebbe BPDU superiore.'),
        ('Perché PortFast non va usato su trunk?',
         'PortFast bypassa STP — su un trunk potrebbe forwardare frame prima che STP abbia rilevato loop.'),
    ],
    'summary': {
        'labels': ['RSTP Standard', 'PortFast + BPDU Guard', 'Root Guard'],
        'bodies': [
            'RSTP 802.1w: convergenza in secondi. Stessa root election. Alternate/Backup port per failover rapido.',
            'PortFast su host: salta Listening/Learning. BPDU Guard: err-disabled se riceve BPDU.',
            'Root Guard su uplink aggregazione: impedisce cambio root bridge non desiderato.',
        ],
    },
}

# ─── MOD-15 ───────────────────────────────────────────────────────────────────
MODULES['MOD-15'] = {
    'title': 'FHRP — HSRP, VRRP & GLBP',
    'area': 'Layer 3 Technologies', 'hours': '2h', 'codes': '3.4.c',
    'agenda': [
        'FHRP: il problema del default gateway',
        'HSRPv2: Virtual IP, priority e preempt',
        'IP SLA + Object Tracking per failover',
        'VRRP 802.1X: differenze da HSRP',
        'GLBP: load balancing attivo-attivo',
    ],
    'topology': {
        'title': 'Topologia MOD-15 — HSRP Dual-Group',
        'caption': 'SW1: HSRP Active VLAN10 (prio 110), Standby VLAN20. SW2: Active VLAN20. VIP: 10.10.10.1/20.1.',
        'nodes': [
            ('R1\nUpstream\n1.1.1.1', 0.5, 0.12),
            ('SW1\nActive VLAN10\nprio 110', 0.25, 0.5),
            ('SW2\nActive VLAN20\nprio 110', 0.75, 0.5),
            ('PC1\nGW 10.10.10.1', 0.15, 0.85),
            ('PC2\nGW 10.10.20.1', 0.85, 0.85),
        ],
    },
    'sections': [
        {
            'section': 'HSRP e FHRP',
            'subtitle': 'Gateway virtuale ridondante',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'FHRP — Il Problema del Default Gateway',
                    'points': [
                        'Problema: se il default gateway fisico cade, tutti gli host su quel segmento perdono la connettività',
                        'FHRP (First Hop Redundancy Protocol): presenta un VIP condiviso da due o più router/switch',
                        'I client configurano il VIP come gateway — non sanno quale device è attivo',
                        'Failover: il dispositivo Standby diventa Active — il VIP "si sposta" senza riconfigurazione client',
                        'HSRP: Cisco proprietario. VRRP: IEEE 802.1X standard. GLBP: Cisco con load-balancing',
                        'MAC virtuale HSRP: 0000.0c07.acXX (XX = gruppo hex). VRRP: 0000.5e00.01XX',
                    ],
                    'key': 'VIP = indirizzo virtuale condiviso. Client puntano al VIP — trasparente al failover.',
                },
                {
                    'type': 'teoria',
                    'title': 'HSRPv2 — Priority, Preempt e Tracking',
                    'points': [
                        'Priority: più alto = Active (default 100). Tie-break: IP più alto vince',
                        'Preempt: il router rientrato riacquista il ruolo Active se ha priority più alta',
                        'Senza preempt: il router rientrato rimane Standby anche con priority più alta',
                        'IP SLA ICMP → Object Track → decremento priority: failover automatico se uplink cade',
                        'Esempio: priority 110 con decrement 20 → se track DOWN: priority 90 → Standby diventa Active',
                        'HSRPv2: group 0-4095, IPv6, autenticazione MD5 — upgrade da v1',
                    ],
                    'key': 'Preempt obbligatorio per riacquisire il ruolo Active dopo recovery. Senza: il vecchio Standby rimane Active.',
                },
                {
                    'type': 'teoria',
                    'title': 'VRRP vs HSRP vs GLBP',
                    'points': [
                        'VRRP (RFC 5798): standard IEEE — Master/Backup (non Active/Standby). Priority default 100',
                        'VRRP: il Master può essere il proprietario del VIP (IP reale = VIP) — differenza da HSRP',
                        'GLBP (Gateway Load Balancing Protocol): Cisco — un AVG assegna MAC virtuali a più AVF',
                        'GLBP: tutti i gateway sono attivi contemporaneamente — load balancing reale',
                        'GLBP algoritmi: round-robin, weighted, host-dependent',
                        'GLBP usa 4 MAC virtuali distinte per distribuire il traffico',
                    ],
                    'key': 'GLBP unico nel permettere load balancing attivo-attivo sul gateway — tutti i router forward traffico.',
                },
            ],
        },
    ],
    'config_section': {
        'slides': [
            {
                'type': 'config',
                'title': 'HSRPv2 con IP SLA Tracking',
                'lines': [
                    'ip sla 1',
                    ' icmp-echo 1.1.1.1 source-interface Vlan10',
                    'ip sla schedule 1 life forever start-time now',
                    'track 1 ip sla 1 reachability',
                    '!',
                    'interface Vlan10',
                    ' standby version 2',
                    ' standby 10 ip 10.10.10.1         ! VIP',
                    ' standby 10 priority 110',
                    ' standby 10 preempt',
                    ' standby 10 track 1 decrement 20   ! 110-20=90 < 100 → failover',
                ],
                'device': 'SW1', 'hl': 10,
            },
        ],
    },
    'trouble': [
        ('HSRP non fa failover', 'Manca preempt o decrement non sufficiente — priority deve scendere sotto quella dello Standby'),
        ('Entrambi i router Active', 'Problema di comunicazione Hello — verificare multicast 224.0.0.2 e VLAN'),
        ('VIP non risponde dopo failover', 'ARP non aggiornato — HSRP manda gratuitous ARP automaticamente. Verificare invio'),
        ('Track non cambia stato', 'IP SLA non schedulato o target non raggiungibile con il source corretto'),
    ],
    'exam_tips': [
        'preempt obbligatorio per riacquisire Active dopo recovery — senza preempt: Standby rimane Active',
        'HSRP: 0000.0c07.acXX. VRRP: 0000.5e00.01XX. GLBP: 0007.b400.XXYY',
        'GLBP = unico load-balancing FHRP — tutti i router inoltrano traffico contemporaneamente',
        'VRRP: il Master può essere il proprietario del VIP — l\'IP reale di un router può essere il VIP',
    ],
    'exam_qa': [
        ('Differenza tra HSRP e VRRP?',
         'HSRP: Cisco, Active/Standby, il VIP è sempre virtuale. VRRP: IEEE standard, Master/Backup, il Master può avere VIP = IP reale.'),
        ('Perché preempt è importante in HSRP?',
         'Senza preempt, lo Standby che è diventato Active rimane tale anche quando il router originale torna — loss of determinism.'),
    ],
    'summary': {
        'labels': ['VIP + MAC Virtuale', 'Priority + Preempt', 'GLBP Load Balance'],
        'bodies': [
            'VIP condiviso tra Active/Standby. Client puntano al VIP. MAC virtuale per ARP.',
            'Priority alta = Active. Preempt = riacquisto automatico. Track + decrement = failover automatico.',
            'GLBP: tutti i router forward simultaneamente. AVG assegna MAC virtuali. Round-robin default.',
        ],
    },
}

# ─── MOD-16 ───────────────────────────────────────────────────────────────────
MODULES['MOD-16'] = {
    'title': 'IP SLA & SPAN',
    'area': 'Network Assurance', 'hours': '2h', 'codes': '4.3 · 4.4',
    'agenda': [
        'IP SLA: monitoraggio proattivo della rete',
        'Object Tracking: azioni condizionali su SLA',
        'Local SPAN: cattura traffico su switch locale',
        'RSPAN: cattura traffico cross-switch',
        'ERSPAN: cattura traffico su reti IP remote',
    ],
    'topology': {
        'title': 'Topologia MOD-16 — IP SLA e SPAN',
        'caption': 'SPAN locale su SW1: sorgente e1/0 (PC1) → dest e1/1. RSPAN: SW1 e1/0 → VLAN 999 → SW2 e1/1.',
        'nodes': [
            ('R1\nTarget SLA\n1.1.1.1', 0.5, 0.1),
            ('SW1\nSPA Local\nSRC: e1/0', 0.25, 0.5),
            ('SW2\nRSPAN dest\ne1/1', 0.75, 0.5),
            ('PC1\nSorgente\ntraffic', 0.1, 0.85),
            ('Sniffer\nSPAN dest', 0.4, 0.85),
        ],
    },
    'sections': [
        {
            'section': 'IP SLA & SPAN',
            'subtitle': 'Monitoraggio proattivo e analisi traffico',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'IP SLA — Monitoraggio Proattivo',
                    'points': [
                        'IP SLA: genera traffico sintetico per misurare la disponibilità di un percorso',
                        'Tipi: icmp-echo (latency/availability), udp-jitter (voce/video), http, tcp-connect',
                        'Frequenza: probe ogni N secondi — rilevazione rapida rispetto a routing protocol',
                        'Threshold: configurare reaction per trigger alert (SNMP trap, syslog)',
                        'Object track: collega SLA a decisioni di routing (rotta condizionale, HSRP decrement)',
                        'ip sla schedule: obbligatorio per avviare la sonda — senza schedule il SLA non gira',
                    ],
                    'key': 'IP SLA = monitoraggio proattivo. Rileva failure prima che OSPF/BGP convergano.',
                },
                {
                    'type': 'teoria',
                    'title': 'SPAN, RSPAN e ERSPAN',
                    'points': [
                        'SPAN (Switched Port ANalyzer): copia traffico da source port/VLAN a destination port locale',
                        'Destinazione SPAN: porta collegata a sniffer/IDS — non invia traffico normale',
                        'RSPAN (Remote SPAN): usa una VLAN dedicata per trasportare il traffico mirror cross-switch',
                        'VLAN RSPAN: non inviare traffico utente — solo RSPAN traffic (no-learn, no-mac)',
                        'ERSPAN (Encapsulated RSPAN): incapsula in GRE — trasporta il mirror su rete IP routable',
                        'Limitazioni SPAN: può ridurre performance dello switch se alta percentuale del traffico',
                    ],
                    'key': 'SPAN=locale, RSPAN=cross-switch via VLAN dedicata, ERSPAN=routable via GRE. Scegliere in base alla distanza.',
                },
            ],
        },
    ],
    'config_section': {
        'slides': [
            {
                'type': 'config',
                'title': 'SPAN e RSPAN — Configurazione',
                'lines': [
                    '! Local SPAN su SW1',
                    'monitor session 1 source interface Et1/0 both',
                    'monitor session 1 destination interface Et1/1',
                    '!',
                    '! RSPAN — SW1 (sorgente)',
                    'vlan 999',
                    ' remote-span',
                    'monitor session 2 source interface Et1/0 both',
                    'monitor session 2 destination remote vlan 999',
                    '!',
                    '! RSPAN — SW2 (destinazione)',
                    'monitor session 2 source remote vlan 999',
                    'monitor session 2 destination interface Et1/1',
                ],
                'device': 'SW1', 'hl': 1,
            },
        ],
    },
    'trouble': [
        ('SPAN non cattura traffico', 'Sessione non configurata o porta dest in access — SPAN dest non deve ricevere traffico normale'),
        ('RSPAN VLAN non transita', 'VLAN 999 non permessa sul trunk SW1↔SW2 o remote-span non configurato sulla VLAN'),
        ('IP SLA sempre DOWN', 'ip sla schedule non configurato — senza schedule la sonda non gira'),
        ('SPAN degrada performance switch', 'Troppo traffico da specchiare — ridurre sorgenti o usare ERSPAN con sampling'),
    ],
    'exam_tips': [
        'SPAN: locale al singolo switch. RSPAN: cross-switch via VLAN dedicata. ERSPAN: IP routable via GRE',
        'RSPAN VLAN: no-learn, no-mac — non deve trasportare traffico normale utente',
        'IP SLA: ip sla schedule obbligatorio — senza schedule la sonda non parte',
        'SPAN destination port: non invia traffico normale — non connettere host normali',
    ],
    'exam_qa': [
        ('Quando usare RSPAN invece di SPAN?',
         'Quando sorgente e destinazione sono su switch diversi — RSPAN usa una VLAN dedicata per trasportare il mirror.'),
        ('IP SLA vs ICMP ping manuale?',
         'IP SLA: continuo, schedulato, integrato con object tracking. Ping manuale: one-shot, non automatizzabile.'),
    ],
    'summary': {
        'labels': ['IP SLA Proattivo', 'SPAN Locale', 'RSPAN / ERSPAN'],
        'bodies': [
            'IP SLA genera traffico sintetico continuo. Rileva failure rapido. Integrato con track e routing.',
            'SPAN: copia traffico porta/VLAN → dest locale. Non impatta routing. Dest non forward normale.',
            'RSPAN: VLAN dedicata cross-switch. ERSPAN: GRE su IP routable. Entrambi per analisi remota.',
        ],
    },
}

# ─── MOD-17 ───────────────────────────────────────────────────────────────────
MODULES['MOD-17'] = {
    'title': 'VRF-Lite & GRE Tunneling',
    'area': 'AREA 7 — OVERLAY & VPN', 'hours': '2h', 'codes': '4.1 · 4.2',
    'agenda': [
        'VRF-Lite: separazione routing senza MPLS',
        'GRE: Generic Routing Encapsulation',
        'Tunnel source, destination e keepalive',
        'Route statiche nel VRF via tunnel GRE',
        'Troubleshooting VRF e tunnel GRE',
    ],
    'topology': {
        'title': 'Topologia MOD-17 — VRF-Lite + GRE',
        'caption': 'HUB-SP1-SP2 su underlay ISP. VRF CUST-A/CUST-B separati. GRE tunnel P2P HUB↔SP1 e HUB↔SP2.',
        'nodes': [
            ('ISP\n192.0.2.253', 0.5, 0.15),
            ('HUB\n192.0.2.254\nVRF CUST-A/B', 0.25, 0.5),
            ('SP1\n198.51.100.254\nVRF CUST-A/B', 0.65, 0.3),
            ('SP2\n203.0.113.254\nVRF CUST-A/B', 0.65, 0.7),
        ],
    },
    'sections': [
        {
            'section': 'VRF-Lite e GRE',
            'subtitle': 'Segmentazione routing e tunnel overlay',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'VRF-Lite — Segmentazione Senza MPLS',
                    'points': [
                        'VRF-Lite: VRF locali senza MP-BGP o MPLS — solo separazione routing su singolo device',
                        'Ogni VRF ha routing table indipendente: CUST-A non vede le rotte di CUST-B',
                        'vrf definition: definisce il VRF con address-family ipv4',
                        'vrf forwarding: assegna l\'interfaccia al VRF — rimuove automaticamente l\'IP',
                        'Route nel VRF: ip route vrf CUST-A 10.1.2.0 255.255.255.0 172.16.101.2',
                        'ping nel VRF: ping vrf CUST-A 10.1.2.1 — senza vrf usa la global table',
                    ],
                    'key': 'VRF-Lite = separazione routing locale. vrf forwarding rimuove l\'IP — riconfigurare dopo.',
                },
                {
                    'type': 'teoria',
                    'title': 'GRE — Generic Routing Encapsulation',
                    'points': [
                        'GRE (RFC 2784): incapsula qualsiasi protocollo L3 in IP — tunnel IP-in-IP',
                        'Header GRE: 4 byte minimo tra outer IP e inner IP — Protocol 47',
                        'Tunnel source: interfaccia locale (tipicamente loopback per stabilità)',
                        'Tunnel destination: IP remoto del tunnel endpoint',
                        'tunnel mode gre ip (default): GRE su IPv4. tunnel mode gre ipv6: GRE su IPv6',
                        'Keepalive: verifica che il tunnel remoto sia attivo — senza, tunnel è sempre UP',
                    ],
                    'key': 'GRE Protocol 47. Source stabile (loopback). Keepalive per rilevare tunnel down.',
                },
            ],
        },
    ],
    'config_section': {
        'slides': [
            {
                'type': 'config',
                'title': 'VRF + GRE Tunnel',
                'lines': [
                    'vrf definition CUST-A',
                    ' address-family ipv4',
                    '!',
                    'interface Loopback1',
                    ' vrf forwarding CUST-A',
                    ' ip address 10.1.1.1 255.255.255.255',
                    '!',
                    'interface Tunnel101',
                    ' vrf forwarding CUST-A',
                    ' ip address 172.16.101.1 255.255.255.252',
                    ' tunnel source Loopback0          ! global table',
                    ' tunnel destination 198.51.100.254 ! SP1 loopback',
                    ' keepalive 10 3',
                ],
                'device': 'HUB', 'hl': 8,
            },
        ],
    },
    'trouble': [
        ('Tunnel DOWN dopo vrf forwarding', 'vrf forwarding rimuove IP — riconfigurare ip address dopo il comando'),
        ('Ping VRF non funziona', 'Manca "ping vrf NOME" — senza vrf usa global table'),
        ('GRE tunnel UP ma traffico non passa', 'Route statiche VRF mancanti o recursive routing (tunnel source in VRF sbagliato)'),
        ('Keepalive non funziona', 'Peer non ha keepalive configurato — keepalive unidirezionale non rileva entrambi i failure'),
    ],
    'exam_tips': [
        'vrf forwarding rimuove l\'IP configurato — sempre riapplicare ip address dopo',
        'ping vrf NOME: obbligatorio per testare connettività in un VRF specifico',
        'GRE tunnel source: usare loopback per stabilità — non interfaccia fisica (potrebbe andare down)',
        'Recursive routing GRE: tunnel destination deve essere raggiungibile senza il tunnel stesso',
    ],
    'exam_qa': [
        ('Cosa succede dopo vrf forwarding su un\'interfaccia con IP già configurato?',
         'L\'IP viene rimosso automaticamente — va riapplicato con ip address dopo il comando vrf forwarding.'),
        ('Differenza tra VRF-Lite e MPLS L3VPN?',
         'VRF-Lite: solo separazione locale, nessuna propagazione inter-PE. L3VPN: VRF propagati via MP-BGP tra PE.'),
    ],
    'summary': {
        'labels': ['VRF-Lite Locale', 'GRE Protocol 47', 'Route + Tunnel'],
        'bodies': [
            'VRF-Lite: routing table separata senza MPLS. vrf forwarding rimuove IP. ping vrf per test.',
            'GRE: IP-in-IP, Protocol 47. Source loopback. Keepalive per rilevare tunnel failure.',
            'Route statiche nel VRF via tunnel. Evitare recursive routing.',
        ],
    },
}

# ─── MOD-18 ───────────────────────────────────────────────────────────────────
MODULES['MOD-18'] = {
    'title': 'IPSec IKEv2 & VTI',
    'area': 'AREA 7 — OVERLAY & VPN', 'hours': '2h', 'codes': '4.4 · 4.5',
    'agenda': [
        'IPSec: suite crittografica e componenti',
        'ESP vs AH: header e protezioni',
        'IKEv2: flusso IKE_SA_INIT → IKE_AUTH',
        'VTI (Virtual Tunnel Interface) con tunnel protection',
        'Verifica SA e contatori IPSec',
    ],
    'topology': {
        'title': 'Topologia MOD-18 — IPSec IKEv2 su GRE',
        'caption': 'HUB↔SP1: GRE + IPSec (IKEv2 + AES-256). HUB↔SP2: reference pre-configurato. ISP underlay pubblico.',
        'nodes': [
            ('ISP\nUnderlay', 0.5, 0.15),
            ('HUB\n192.0.2.254\nIKEv2 initiator', 0.25, 0.5),
            ('SP1\n198.51.100.254\nIKEv2 responder', 0.75, 0.35),
            ('SP2\n203.0.113.254\nPre-conf ref', 0.75, 0.7),
        ],
    },
    'sections': [
        {
            'section': 'IPSec — Suite e Componenti',
            'subtitle': 'Cifratura, integrità e autenticazione peer',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'IPSec — ESP, AH e Modalità',
                    'points': [
                        'ESP (Encapsulating Security Payload): Protocol 50 — cifratura + integrità + autenticazione',
                        'AH (Authentication Header): Protocol 51 — solo integrità, nessuna cifratura',
                        'SPI (Security Parameter Index): identifica la SA (Security Association) — 32-bit nel header',
                        'ICV (Integrity Check Value): HMAC del payload — verifica integrità',
                        'Tunnel mode: incapsula intero IP originale — usato tra gateway (VPN)',
                        'Transport mode: protegge solo il payload L4 — usato end-to-end tra host',
                    ],
                    'key': 'ESP = cifratura + integrità. Tunnel mode per VPN tra gateway. Transport mode per host-to-host.',
                },
                {
                    'type': 'teoria',
                    'title': 'IKEv2 — Flusso di Negoziazione',
                    'points': [
                        'IKEv2: negozia la SA IPSec in due fasi — semplificato rispetto a IKEv1',
                        'IKE_SA_INIT: negozia algoritmi IKE, scambia chiavi Diffie-Hellman (DH group)',
                        'IKE_AUTH: autentica i peer (PSK o certificati), stabilisce la Child SA (IPSec SA)',
                        'CREATE_CHILD_SA: aggiunge SA aggiuntive o esegue rekey',
                        'Vantaggio IKEv2: meno messaggi di scambio, PFS, MOBIKE, EAP integrato',
                        'SA lifetime: rikey automatico prima della scadenza — no disruption del traffico',
                    ],
                    'key': 'IKEv2: IKE_SA_INIT (Diffie-Hellman) + IKE_AUTH (autenticazione + IPSec SA). Più efficiente di IKEv1.',
                },
                {
                    'type': 'teoria',
                    'title': 'VTI — tunnel protection ipsec profile',
                    'points': [
                        'tunnel protection ipsec profile: associa un profilo IPSec a un tunnel GRE — VTI',
                        'VTI vs crypto map: VTI è compatibile con mGRE/DMVPN, crypto map no',
                        'Crypto map: legacy, richiede ACL per selezionare traffico, non scalabile',
                        'IPSec profile: definisce transform-set + DH group — applicato al tunnel',
                        'GRE + IPSec: GRE fornisce il tunnel, IPSec cifra tutto il contenuto GRE',
                        'Verifica: show crypto ikev2 sa + show crypto ipsec sa + contatori pkts encaps/decaps',
                    ],
                    'key': 'tunnel protection ipsec profile = VTI moderno. Compatibile con DMVPN/mGRE. Crypto map è legacy.',
                },
            ],
        },
    ],
    'config_section': {
        'slides': [
            {
                'type': 'config',
                'title': 'IKEv2 + IPSec Profile + VTI',
                'lines': [
                    'crypto ikev2 proposal PROP-1',
                    ' encryption aes-cbc-256',
                    ' integrity sha256',
                    ' group 14',
                    'crypto ikev2 policy POL-1',
                    ' proposal PROP-1',
                    'crypto ikev2 keyring KR-1',
                    ' peer SP1',
                    '  address 198.51.100.254',
                    '  pre-shared-key local K3yS3cr3t',
                    '  pre-shared-key remote K3yS3cr3t',
                    'crypto ipsec transform-set TS esp-aes 256 esp-sha256-hmac',
                    ' mode tunnel',
                    'crypto ipsec profile IPSEC-PROF',
                    ' set transform-set TS',
                    'interface Tunnel101',
                    ' tunnel protection ipsec profile IPSEC-PROF',
                ],
                'device': 'HUB', 'hl': 15,
            },
            {
                'type': 'verifica',
                'title': 'Verifica IKEv2 e IPSec SA',
                'lines': [
                    'HUB# show crypto ikev2 sa',
                    'Tunnel-id  Local              Remote             Status',
                    '1          192.0.2.254/4500   198.51.100.254/4500  READY',
                    '',
                    'HUB# show crypto ipsec sa | incl pkts',
                    '    #pkts encaps: 1024, #pkts encrypt: 1024',
                    '    #pkts decaps: 998,  #pkts decrypt: 998',
                ],
                'hl': 2,
            },
        ],
    },
    'trouble': [
        ('IKEv2 SA non si forma', 'Proposal non corrispondenti tra peer — encryption/integrity/group diversi'),
        ('Phase 1 UP ma pkts encaps = 0', 'Transform-set mismatch (Phase 2 fallisce) — show crypto ipsec sa'),
        ('Tunnel UP ma no traffico cifrato', 'tunnel protection non applicato al tunnel o profilo errato'),
        ('Rekey failure', 'DH group diverso per PFS — verificare group nel profilo IKEv2 e IPSec'),
    ],
    'exam_tips': [
        'ESP Protocol 50 (non TCP/UDP). SPI identifica la SA. ICV verifica integrità',
        'IKEv2: IKE_SA_INIT + IKE_AUTH → Phase 1. CREATE_CHILD_SA → Phase 2 (IPSec SA)',
        'tunnel protection ipsec profile: obbligatorio per DMVPN — crypto map non compatibile con mGRE',
        'pkts encaps/decaps = 0 dopo SA UP: transform-set mismatch tra peer',
    ],
    'exam_qa': [
        ('Differenza tra tunnel protection e crypto map?',
         'tunnel protection: per VTI/DMVPN, compatibile mGRE, no ACL. Crypto map: legacy, richiede ACL, non scalabile.'),
        ('Cosa identifica l\'SPI in IPSec?',
         'Security Parameter Index: identifica la Security Association (SA) — indica quale SA usare per il pacchetto.'),
    ],
    'summary': {
        'labels': ['ESP Protocol 50', 'IKEv2 Due Fasi', 'VTI tunnel protection'],
        'bodies': [
            'ESP = cifratura + integrità. AH = solo integrità. Tunnel mode per VPN gateway-to-gateway.',
            'IKE_SA_INIT: DH + algoritmi. IKE_AUTH: autenticazione + IPSec SA. Più efficiente di IKEv1.',
            'tunnel protection ipsec profile: associa IPSec al GRE tunnel. Compatibile DMVPN.',
        ],
    },
}

# ─── MOD-19 ───────────────────────────────────────────────────────────────────
MODULES['MOD-19'] = {
    'title': 'DMVPN Phase 1, 2 & 3',
    'area': 'AREA 7 — OVERLAY & VPN', 'hours': '3h', 'codes': '4.6 · 4.7 · 4.8',
    'agenda': [
        'DMVPN: mGRE, NHRP e IPSec',
        'NHRP: NHS, NHC e registrazione',
        'Phase 1: hub-and-spoke classico',
        'Phase 2: spoke-to-spoke diretti',
        'Phase 3: shortcut routing avanzato',
    ],
    'topology': {
        'title': 'Topologia MOD-19 — DMVPN Cloud',
        'caption': 'HUB: NHS + mGRE. SP1/SP2: NHC. Phase 2/3: shortcut diretto SP1↔SP2 senza passare per HUB.',
        'nodes': [
            ('HUB\nNHS\nTu110/210', 0.5, 0.2),
            ('SP1\nNHC\n.110.11', 0.2, 0.7),
            ('SP2\nNHC\n.110.12', 0.8, 0.7),
            ('ISP\nUnderlay', 0.5, 0.85),
        ],
    },
    'sections': [
        {
            'section': 'DMVPN — Architettura',
            'subtitle': 'mGRE + NHRP + IPSec = DMVPN',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'DMVPN — I Tre Componenti',
                    'points': [
                        'mGRE (multipoint GRE): un tunnel fisico supporta connessioni dinamiche a N peer',
                        'NHRP (Next Hop Resolution Protocol): mappa IP tunnel → NBMA (IP pubblico)',
                        'IPSec: cifra il traffico GRE — tunnel protection ipsec profile',
                        'Hub = NHS (Next Hop Server): registra e risponde alle richieste NHRP',
                        'Spoke = NHC (Next Hop Client): si registra all\'Hub, risolve IP tunnel degli altri spoke',
                        'Vantaggio: spoke non deve conoscere a priori gli altri spoke — risoluzione dinamica',
                    ],
                    'key': 'DMVPN = mGRE (1 tunnel N peer) + NHRP (risolve IP tunnel) + IPSec (cifratura). Dinamico.',
                },
                {
                    'type': 'teoria',
                    'title': 'DMVPN Phase 1, 2 e 3',
                    'points': [
                        'Phase 1: tutto il traffico passa per l\'Hub — Hub invia il traffico agli spoke',
                        'Phase 2: spoke risolve via NHRP l\'IP dell\'altro spoke → shortcut diretto',
                        'Phase 3: Hub manda NHRP Redirect → spoke installa rotta host via shortcut',
                        'HUB Phase 2: no split-horizon EIGRP + no next-hop-self → spoke vede IP spoke remoto',
                        'HUB Phase 3: ip nhrp redirect + ip summary-address EIGRP sull\'Hub',
                        'Phase 3: rotte più specifiche installate dynamicamente — scalabilità massima',
                    ],
                    'key': 'Phase 2: no split-horizon + no next-hop-self sull\'Hub. Phase 3: nhrp redirect per shortcut.',
                },
            ],
        },
    ],
    'config_section': {
        'slides': [
            {
                'type': 'config',
                'title': 'DMVPN Phase 2 — Hub Config',
                'lines': [
                    'interface Tunnel110',
                    ' vrf forwarding CUST-A',
                    ' ip address 172.16.110.1 255.255.255.0',
                    ' tunnel source Loopback0',
                    ' tunnel mode gre multipoint       ! mGRE',
                    ' ip nhrp network-id 110',
                    ' ip nhrp map multicast dynamic    ! spoke registra multicast',
                    ' tunnel protection ipsec profile IPSEC-PROF',
                    '!',
                    'router eigrp DMVPN',
                    ' address-family ipv4 vrf CUST-A',
                    '  af-interface Tunnel110',
                    '   no split-horizon',
                    '   no next-hop-self',
                ],
                'device': 'HUB', 'hl': 12,
            },
        ],
    },
    'trouble': [
        ('Spoke non si registra all\'NHS', 'NHRP NHS IP sbagliato o NBMA map mancante — ip nhrp nhs + ip nhrp map'),
        ('Phase 2 shortcut non funziona', 'split-horizon o next-hop-self non rimosso sull\'Hub — EIGRP non propaga correttamente'),
        ('IPSec non negozia tra spoke', 'Spoke non hanno keyring per peer-to-peer — usare wildcard in keyring'),
        ('show dmvpn mostra NHRP INCOMPLETE', 'NHRP resolution fallita — verificare che lo spoke sia registrato con show ip nhrp'),
    ],
    'exam_tips': [
        'no split-horizon + no next-hop-self sull\'Hub: entrambi critici per DMVPN Phase 2',
        'mGRE: tunnel mode gre multipoint — un tunnel fisico gestisce N spoke dinamicamente',
        'NHRP: NHS=Hub registra tutti gli spoke. NHC=Spoke risolve IP degli altri spoke',
        'tunnel protection ipsec profile: obbligatorio. Crypto map non compatibile con mGRE',
    ],
    'exam_qa': [
        ('Perché no split-horizon è necessario sull\'Hub in DMVPN Phase 2?',
         'Split-horizon blocca EIGRP dal ri-annunciare rotte spoke ad altri spoke. Senza: spoke non vede le rotte degli altri.'),
        ('Differenza tra DMVPN Phase 2 e Phase 3?',
         'Phase 2: shortcut IP spoke. Phase 3: nhrp redirect + rotte host dinamiche — più scalabile, shortcut più efficiente.'),
    ],
    'summary': {
        'labels': ['mGRE Dinamico', 'NHRP NHS/NHC', 'Phase 2: no NH-self'],
        'bodies': [
            'mGRE: un tunnel, N spoke dinamici. NHRP risolve IP tunnel → NBMA. IPSec cifra tutto.',
            'NHS=Hub: registra spoke. NHC=Spoke: si registra + risolve peer. Risoluzione on-demand.',
            'Phase 2: no split-horizon + no next-hop-self sull\'Hub. Phase 3: nhrp redirect per shortcut.',
        ],
    },
}


MODULES['MOD-20'] = {
    'title': 'LISP & VXLAN — Architettura SD-Access',
    'area': 'AREA 4 — OVERLAY & SD-ACCESS',
    'hours': '1.5',
    'codes': '4.9, 4.10',
    'agenda': [
        'Il problema della mobilità IP',
        'LISP — EID, RLOC e flusso Map-Request',
        'VXLAN — VNI, VTEP e incapsulamento',
        'LISP + VXLAN in SD-Access',
        'Exam Tips & Summary',
    ],
    'topology': {
        'title': 'Architettura LISP + VXLAN — SD-Access',
        'nodes': [('Host A\n(EID)', 0.1, 0.5), ('ITR (xTR-A)', 0.3, 0.5), ('LISP MS/MR', 0.5, 0.2), ('ETR (xTR-B)', 0.7, 0.5), ('Host B\n(EID)', 0.9, 0.5)],
        'caption': 'Separazione EID (identità) vs RLOC (locazione) — flusso Map-Request/Reply',
    },
    'sections': [
        {
            'section': 'PARTE 1 — LISP: Locator/ID Separation',
            'subtitle': 'Separazione identità/locazione',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Il Problema della Mobilità IP',
                    'points': [
                        'In IP tradizionale un indirizzo ha due ruoli: identità (chi sei) e locazione (dove sei)',
                        'La mobilità è costosa: ogni spostamento richiede cambio IP o aggiornamenti routing massivi',
                        'LISP separa i due piani: EID (Endpoint Identifier) identifica chi sei — fisso',
                        'RLOC (Routing Locator) identifica dove sei — cambia con la posizione nella rete',
                        'Analogia: EID = numero di telefono (fisso), RLOC = cella GSM corrente (variabile)',
                    ],
                    'key': 'LISP = separazione identità/locazione — host mantiene EID anche cambiando rete',
                },
                {
                    'type': 'teoria',
                    'title': 'Componenti LISP',
                    'points': [
                        'EID: indirizzo logico del dispositivo — invariato con lo spostamento',
                        'RLOC: indirizzo del router di bordo che raggiunge l\'EID — varia con posizione',
                        'Map-Server (MS): riceve registrazioni EID→RLOC dagli xTR, popola il database',
                        'Map-Resolver (MR): risponde alle Map-Request degli ITR con il RLOC corretto',
                        'ITR: riceve traffico, invia Map-Request, incapsula in UDP/4341 verso RLOC',
                        'ETR: riceve pacchetti LISP incapsulati, decapsula, consegna all\'EID locale',
                    ],
                    'key': 'xTR = router che svolge sia ITR che ETR — ruolo tipico in SD-Access',
                },
                {
                    'type': 'teoria',
                    'title': 'Flusso Map-Request/Reply',
                    'points': [
                        '1. Host A (EID: 10.1.1.10) vuole raggiungere Host B (EID: 10.2.1.10)',
                        '2. ITR non ha mappatura EID→RLOC in cache: invia Map-Request al MR',
                        '3. MR risponde con Map-Reply: "RLOC = 203.0.113.254 (xTR di Host B)"',
                        '4. ITR incapsula: [outer IP: ITR→RLOC][UDP 4341][EID payload originale]',
                        '5. ETR (203.0.113.254) riceve, decapsula, consegna a Host B',
                        '6. ITR salva mappatura in cache (TTL configurabile) — no Map-Request ripetute',
                    ],
                    'key': 'Porta LISP = UDP 4341. Cache ITR evita Map-Request per ogni pacchetto.',
                },
            ],
        },
        {
            'section': 'PARTE 2 — VXLAN: Virtual eXtensible LAN',
            'subtitle': 'Segmentazione estesa su rete IP',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Il Limite di 802.1Q e la Soluzione VXLAN',
                    'points': [
                        '802.1Q VLAN tag: 12 bit → max 4.096 segmenti — insufficiente per datacenter multi-tenant',
                        'VXLAN usa VNI (VXLAN Network Identifier): 24 bit → ~16.777.216 segmenti',
                        'VXLAN incapsula l\'intero frame L2 originale in pacchetto UDP/IP',
                        'VTEP (VXLAN Tunnel Endpoint): dispositivo che incapsula/decapsula — switch, hypervisor, router',
                        'Underlay: rete IP routable tra VTEP. Overlay: segmenti L2 virtuali (VNI) su L3',
                    ],
                    'key': 'VNI 24-bit = ~16 milioni di segmenti vs 4.096 VLAN 802.1Q',
                },
                {
                    'type': 'teoria',
                    'title': 'Struttura Pacchetto VXLAN',
                    'points': [
                        'Stack: Outer IP | UDP (dport 4789) | VXLAN Header | Inner Ethernet Frame',
                        'Outer IP Header: VTEP-A → VTEP-B (underlay — qualsiasi rete IP routable)',
                        'UDP Header: porta destinazione 4789 (IANA), porta sorgente = hash del flusso',
                        'VXLAN Header: 8 byte — VNI 24-bit + flag I + reserved',
                        'Inner Ethernet Frame: frame L2 originale dell\'host — MAC src/dst inclusi',
                        'Overhead: ~50 byte per frame vs 4 byte tag 802.1Q',
                    ],
                    'key': 'Porta VXLAN = UDP 4789. L\'intero frame Ethernet è preservato nell\'overlay.',
                },
                {
                    'type': 'teoria',
                    'title': 'BUM Traffic e VTEP Discovery',
                    'points': [
                        'BUM = Broadcast/Unknown unicast/Multicast — traffico da inondare a tutti i VTEP del VNI',
                        'Opzione 1 — Multicast underlay: ogni VNI mappa a un gruppo multicast IP nell\'underlay',
                        'Opzione 2 — Ingress Replication: VTEP mittente invia copia unicast a ciascun VTEP',
                        'Multicast: scalabile, richiede multicast routing nell\'underlay',
                        'Ingress Replication: semplice, no multicast, overhead proporzionale a N VTEP',
                        'Cisco SD-Access usa ingress replication — no multicast richiesto nel campus',
                    ],
                    'key': 'SD-Access: ingress replication per BUM — semplicità di deployment campus',
                },
            ],
        },
        {
            'section': 'PARTE 3 — LISP + VXLAN in SD-Access',
            'subtitle': 'Cisco Campus Fabric Architecture',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'SD-Access: Separazione dei Piani',
                    'points': [
                        'Control Plane = LISP: gestisce mobilità endpoint, mappa EID→RLOC',
                        'Data Plane = VXLAN: trasporta frame L2 su L3 con segmentazione VNI',
                        'Policy Plane = SGT (Security Group Tag): policy per gruppo, indipendente dall\'IP',
                        'Fabric Edge (xTR + VTEP): router di accesso — incapsula/decapsula + registra EID',
                        'Fabric Border: connette la fabric a reti esterne, ha ruolo MS/MR',
                        'Catalyst Center: orchestrazione centralizzata — provisioning automatico dei nodi',
                    ],
                    'key': 'SD-Access: LISP=control · VXLAN=data · SGT=policy — tripletta da memorizzare',
                },
                {
                    'type': 'teoria',
                    'title': 'Confronto SD-Access vs SD-WAN',
                    'points': [
                        'SD-Access (campus): Fabric Edge = switch Catalyst, gestito da Catalyst Center',
                        'SD-WAN (WAN): vEdge/cEdge, gestito da Cisco vManage',
                        'Analogia controllo: LISP (SD-Access) vs OMP (SD-WAN)',
                        'Analogia locazione: RLOC (SD-Access) vs TLOC (SD-WAN)',
                        'Entrambi separano control plane da data plane — stesso principio, ambiti diversi',
                        'Domanda ENCOR: "quale è il control plane SD-Access?" = LISP (Fabric Control Plane node)',
                    ],
                    'key': 'NOTA: LISP e VXLAN non configurabili su IOU — questo modulo è teoria pura',
                },
            ],
        },
    ],
    'trouble': [
        ('Map-Request senza risposta', 'MR non raggiungibile o xTR non registrato — verificare ip lisp database e connettività al MS'),
        ('VTEP non si trova', 'Underlay routing non converge — ping tra VTEP IP (RLOC) deve funzionare prima di VXLAN'),
        ('Host non si sposta correttamente', 'Cache ITR scaduta o non aggiornata — clear lisp cache / attendere TTL'),
        ('BUM flooding eccessivo', 'VNI non mappato a gruppo multicast o lista ingress replication incompleta'),
    ],
    'exam_tips': [
        'EID = chi sei (fisso), RLOC = dove sei (variabile) — separazione fondamentale di LISP',
        'MS riceve registrazioni EID→RLOC, MR risponde alle Map-Request — ruoli distinti',
        'VNI = 24 bit = ~16 milioni di segmenti (vs 4.096 VLAN 802.1Q con 12 bit)',
        'VTEP = dispositivo che incapsula/decapsula frame Ethernet in UDP/4789',
        'SD-Access: LISP=control plane, VXLAN=data plane, SGT=policy plane',
    ],
    'exam_qa': [
        ('In LISP, quale componente riceve le registrazioni EID→RLOC dagli xTR?',
         'Map-Server (MS). Il Map-Resolver (MR) risponde alle Map-Request — sono ruoli separati anche se spesso sullo stesso nodo.'),
        ('Quale porta UDP usa VXLAN per il trasporto?',
         'UDP 4789 (IANA assigned). LISP usa UDP 4341.'),
        ('In SD-Access, cosa succede all\'EID di un host che si sposta tra due Fabric Edge?',
         'L\'EID rimane invariato — solo il RLOC cambia nel database LISP (nuovo xTR). Nessuna riconfigurazione sull\'host.'),
    ],
    'summary': {
        'labels': ['LISP: EID vs RLOC', 'VXLAN: 24-bit VNI', 'SD-Access: 3 Piani'],
        'bodies': [
            'LISP separa identità (EID fisso) da locazione (RLOC variabile). MS registra, MR risponde. UDP 4341.',
            'VXLAN: 16M segmenti VNI vs 4096 VLAN. VTEP incapsula frame L2 in UDP/4789. Underlay IP qualsiasi.',
            'SD-Access: LISP=control, VXLAN=data, SGT=policy. Fabric Edge=xTR+VTEP. Catalyst Center orchestra.',
        ],
    },
}


MODULES['MOD-21'] = {
    'title': 'SD-WAN — Architettura e Componenti',
    'area': 'AREA 8 — SD-WAN & SD-ACCESS', 'hours': '1.5', 'codes': '4.7, 4.8',
    'agenda': ['Limiti WAN tradizionale', 'Componenti SD-WAN Cisco', 'OMP e TLOC', 'ZTP e Onboarding', 'Exam Tips'],
    'topology': {
        'title': 'Architettura SD-WAN — Piano di Controllo e Dati',
        'nodes': [('vManage\n(gestione)', 0.2, 0.2), ('vBond\n(orchestrator)', 0.5, 0.2), ('vSmart\n(controller)', 0.8, 0.2), ('cEdge\nSite A', 0.2, 0.7), ('cEdge\nSite B', 0.8, 0.7)],
        'caption': 'vManage=NMS · vBond=autenticazione · vSmart=OMP · cEdge=data plane',
    },
    'sections': [
        {
            'section': 'Componenti SD-WAN Cisco',
            'subtitle': 'Architettura a 4 piani separati',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'I Quattro Componenti SD-WAN',
                    'points': [
                        'vManage: NMS centralizzato — configura, monitora e fa troubleshooting di tutti i device',
                        'vBond: orchestrator — primo punto di contatto, autentica i device con certificati',
                        'vSmart: controller del control plane — distribuisce policy e routing via OMP',
                        'vEdge / cEdge: router WAN — forwarding del traffico (data plane), aderisce alle policy',
                        'Separazione: management plane (vManage), control plane (vSmart), data plane (cEdge)',
                        'Tutti i componenti comunicano su canale TLS/DTLS cifrato — nessun testo in chiaro',
                    ],
                    'key': 'vBond = punto di ingresso obbligatorio. Senza vBond, nessun device si connette alla fabric.',
                },
                {
                    'type': 'teoria',
                    'title': 'OMP — Overlay Management Protocol',
                    'points': [
                        'OMP: protocollo proprietario Cisco, sostituisce BGP/OSPF nel control plane SD-WAN',
                        'Trasporta: rotte (prefissi), policy, servizi tra vSmart e cEdge',
                        'TLOC (Transport Locator): identifica un link WAN fisico — tupla (sistema, colore, incapsulamento)',
                        'Colori TLOC: categorizzano il tipo di link (mpls, biz-internet, lte, private1, ecc.)',
                        'vSmart seleziona i migliori TLOC per ogni prefisso e li propaga ai cEdge',
                        'Analogia: OMP = BGP del SD-WAN. TLOC = RLOC di LISP applicato alla WAN.',
                    ],
                    'key': 'OMP separa il control plane SD-WAN dal data plane — cEdge non parla direttamente tra loro',
                },
            ],
        },
        {
            'section': 'ZTP, Policy e Data Plane',
            'subtitle': 'Onboarding automatico e forwarding cifrato',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Zero Touch Provisioning (ZTP)',
                    'points': [
                        'ZTP: il cEdge si connette a vBond automaticamente al primo avvio — nessun pre-config manuale',
                        'Requisito: certificato del device firmato dalla CA Cisco (o privata) preinstallato in fabbrica',
                        'Flusso: cEdge → DHCP → ottiene IP + indirizzo vBond → TLS handshake → autenticato',
                        'vBond informa vManage e vSmart del nuovo device — scarica template di configurazione',
                        'Vantaggioso: uffici remoti collegati senza intervento tecnico locale (truck roll zero)',
                        'Prerequisito: connettività internet sul link bootstrap prima dell\'attivazione SD-WAN',
                    ],
                    'key': 'ZTP = deploy senza tecnico in loco. Il certificato preinstallato garantisce identità sicura.',
                },
                {
                    'type': 'teoria',
                    'title': 'Data Plane: IPSec tra cEdge',
                    'points': [
                        'cEdge stabilisce tunnel IPSec direttamente con altri cEdge (peer-to-peer)',
                        'I TLOC scambiati via OMP forniscono gli IP pubblici dei peer — no vSmart nel data path',
                        'Ogni coppie di TLOC forma un tunnel IPSec indipendente — multi-path nativo',
                        'Algoritmo: AES-256-GCM per default — FIPS compliant',
                        'BFD: rileva guasti sul tunnel in secondi — rerouting automatico su TLOC alternativo',
                        'Policy applicata al data plane: centralized policy (vSmart) o localized policy (cEdge)',
                    ],
                    'key': 'IPSec diretti tra cEdge: vSmart controlla QUALE tunnel usare, non il traffico stesso',
                },
            ],
        },
    ],
    'trouble': [
        ('cEdge non si connette a vBond', 'DNS risolve ztp.viptela.com? Certificato valido? Firewall blocca UDP/12346?'),
        ('OMP session down', 'vSmart raggiungibile? Mismatch site-id o org-name? Verificare show sdwan omp summary'),
        ('Tunnel IPSec non forma', 'BFD failure? NAT simmetrico sul path? Colore TLOC non annunciato da vSmart?'),
        ('Policy non applicata', 'Policy attiva su vManage? cEdge ha ricevuto policy via OMP? show sdwan policy from-vsmart'),
    ],
    'exam_tips': [
        'vBond = orchestrator: autentica i device. vSmart = controller: distribuisce OMP. vManage = NMS.',
        'OMP trasporta rotte, TLOC e policy tra vSmart e cEdge — non BGP/OSPF nel control plane SD-WAN',
        'TLOC = (system-ip, colore, incapsulamento) — identifica univocamente un link WAN fisico',
        'ZTP: device si auto-registra via vBond — certificato preinstallato in fabbrica è il prerequisito',
        'IPSec tra cEdge è peer-to-peer — vSmart non è nel data path, solo nel control path',
    ],
    'exam_qa': [
        ('In Cisco SD-WAN, quale componente autentica i device al primo avvio?', 'vBond (orchestrator). vSmart distribuisce policy, vManage gestisce la configurazione.'),
        ('Cosa identifica un TLOC in SD-WAN?', 'La tupla (system-ip, colore, incapsulamento) — identifica un link WAN fisico specifico su un cEdge.'),
    ],
    'summary': {
        'labels': ['4 Componenti', 'OMP + TLOC', 'ZTP + IPSec'],
        'bodies': [
            'vManage=NMS, vBond=auth, vSmart=controller OMP, cEdge=data plane. TLS/DTLS cifra tutto.',
            'OMP: rotte+policy tra vSmart/cEdge. TLOC=(system-ip, colore, encap) identifica link WAN.',
            'ZTP: auto-onboarding via certificato. IPSec peer-to-peer tra cEdge — vSmart non nel data path.',
        ],
    },
}

MODULES['MOD-22'] = {
    'title': 'SD-Access — Campus Fabric Cisco',
    'area': 'AREA 8 — SD-WAN & SD-ACCESS', 'hours': '1.5', 'codes': '4.9, 4.10',
    'agenda': ['Limiti del campus tradizionale', 'Fabric: Edge, Border, Control', 'Catalyst Center e provisioning', 'SGT e macro-segmentazione', 'Exam Tips'],
    'topology': {
        'title': 'SD-Access Campus Fabric — Piani e Nodi',
        'nodes': [('Catalyst Center\n(orchestrazione)', 0.5, 0.1), ('Fabric Border\nMS/MR', 0.5, 0.4), ('Fabric Edge A\nxTR+VTEP', 0.2, 0.7), ('Fabric Edge B\nxTR+VTEP', 0.8, 0.7)],
        'caption': 'Control: LISP · Data: VXLAN · Policy: SGT · Orchestrazione: Catalyst Center',
    },
    'sections': [
        {
            'section': 'Architettura Campus Fabric',
            'subtitle': 'Nodi, ruoli e piani di comunicazione',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'I Nodi della Fabric SD-Access',
                    'points': [
                        'Fabric Edge: switch di accesso — connette endpoint, incapsula in VXLAN, registra EID in LISP',
                        'Fabric Border: gateway tra fabric e rete esterna (internet, WAN, datacenter)',
                        'Fabric Control Plane: nodo con LISP Map-Server + Map-Resolver — gestisce EID→RLOC',
                        'Intermediate Node: switch di distribuzione/core — forwarding IP underlay, no ruolo fabric',
                        'Catalyst Center: orchestrazione e provisioning automatico di tutti i nodi sopra',
                        'ISE (Identity Services Engine): fornisce autenticazione 802.1X e assegna SGT agli endpoint',
                    ],
                    'key': 'Fabric Edge = xTR + VTEP. Fa tutto: encap VXLAN, registra LISP, applica SGT.',
                },
                {
                    'type': 'teoria',
                    'title': 'SGT — Security Group Tag',
                    'points': [
                        'SGT: tag a 16 bit assegnato a ogni endpoint dall\'ISE al momento dell\'autenticazione 802.1X',
                        'Identifica il gruppo di sicurezza (es. SGT 10 = "Dipendenti", SGT 20 = "Guest")',
                        'Policy: "gruppo X può parlare con gruppo Y" — indipendente dall\'IP o dalla VLAN',
                        'Vantaggio: policy invariante con lo spostamento dell\'host — segue il tag, non l\'IP',
                        'Macro-segmentazione: VN (Virtual Network) separa domini completi (es. Employees vs IoT)',
                        'Micro-segmentazione: SGT separa gruppi all\'interno dello stesso VN',
                    ],
                    'key': 'SGT = policy per identità, non per posizione. Invariante al movimento host.',
                },
                {
                    'type': 'teoria',
                    'title': 'Catalyst Center — Provisioning e DNA',
                    'points': [
                        'Catalyst Center (ex DNA Center): controller centralizzato per il campus Cisco',
                        'Funzioni: provisioning zero-touch, template di configurazione, assurance (AI/ML)',
                        'Provisioning: definire un sito, aggiungere device, assegnare ruolo fabric — automazione completa',
                        'DNA Assurance: telemetria streaming da tutti i device — rileva anomalie prima che impattino',
                        'Intent-based networking: l\'operatore definisce l\'OBIETTIVO, Catalyst Center lo traduce in config',
                        'API northbound: REST/RESTCONF per integrazione con sistemi terzi (ITSM, SIEM)',
                    ],
                    'key': 'Intent-based: definisci cosa vuoi, non come configurarlo — Catalyst Center traduce.',
                },
            ],
        },
    ],
    'trouble': [
        ('Host non autenticato in fabric', 'ISE policy? 802.1X non configurato sul Fabric Edge? Verificare show authentication sessions'),
        ('SGT non propagato', 'SXP abilitato? ISE-Fabric Edge trust? Verificare show cts role-based sgt-map'),
        ('LISP registration fallisce', 'Fabric Control Plane raggiungibile? show lisp service ipv4 | include map-server'),
        ('Catalyst Center non vede device', 'CDP/LLDP abilitato? Credenziali SNMP/SSH corrette? Device nella subnet gestita?'),
    ],
    'exam_tips': [
        'Fabric Edge = xTR (LISP) + VTEP (VXLAN) — incapsula e registra EID all\'ingresso della fabric',
        'SGT: assegnato da ISE al login 802.1X — policy per gruppo, indipendente da IP e VLAN',
        'Macro-segmentazione = VN (Virtual Network). Micro-segmentazione = SGT dentro il VN.',
        'Catalyst Center orchestra tutto: provisioning automatico, assurance AI, API REST northbound',
        'Fabric Control Plane = LISP MS/MR. Non è un router — è un database di mapping EID→RLOC.',
    ],
    'exam_qa': [
        ('Quale componente SD-Access assegna il SGT a un endpoint?', 'ISE (Identity Services Engine), al momento dell\'autenticazione 802.1X. Il Fabric Edge applica il tag al traffico.'),
        ('Differenza tra macro e micro-segmentazione in SD-Access?', 'Macro: VN (Virtual Network) separa domini completi. Micro: SGT separa gruppi all\'interno dello stesso VN.'),
    ],
    'summary': {
        'labels': ['Nodi Fabric', 'SGT: Policy per ID', 'Catalyst Center'],
        'bodies': [
            'Edge=xTR+VTEP, Border=gateway, Control=LISP MS/MR. Intermediate node=solo underlay.',
            'SGT 16-bit da ISE 802.1X. Policy gruppo-a-gruppo indipendente da IP. VN=macro, SGT=micro.',
            'Intent-based: provisioning zero-touch, DNA Assurance AI/ML, API REST northbound.',
        ],
    },
}

MODULES['MOD-23'] = {
    'title': 'Wireless RF & Fondamenta 802.11',
    'area': 'AREA 9 — WIRELESS', 'hours': '2', 'codes': '2.1.a · 2.1.b · 2.1.c',
    'agenda': ['Fondamenta RF: frequenze e canali', 'Standard 802.11 a/b/g/n/ac/ax', 'BSS, ESS e roaming', 'CSMA/CA e hidden node', 'Exam Tips'],
    'topology': {
        'title': 'Architettura Wireless — BSS ed ESS',
        'nodes': [('AP-1\nBSSID-A', 0.25, 0.4), ('AP-2\nBSSID-B', 0.75, 0.4), ('Client A', 0.1, 0.7), ('Client B', 0.4, 0.7), ('DS\n(Rete cablata)', 0.5, 0.15)],
        'caption': 'BSS = AP + client associati · ESS = più BSS con stesso SSID · DS = distribution system',
    },
    'sections': [
        {
            'section': 'Radio Frequency Fundamentals',
            'subtitle': 'Frequenze, canali e propagazione',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Bande di Frequenza Wi-Fi',
                    'points': [
                        '2.4 GHz: range maggiore, penetrazione pareti ottima — ma solo 3 canali non sovrapposti (1, 6, 11)',
                        '5 GHz: range minore, più interferenze, ma 25+ canali non sovrapposti — throughput maggiore',
                        '6 GHz (Wi-Fi 6E): banda nuova, 59 canali da 20 MHz, solo dispositivi moderni',
                        'RSSI (Received Signal Strength Indicator): misura in dBm — più vicino a 0 = più forte',
                        '-70 dBm: soglia minima per dati, -80 dBm: voce marginalizzata, -90 dBm: inutilizzabile',
                        'Path loss: il segnale decade con il quadrato della distanza + ostacoli fisici (pareti, metallo)',
                    ],
                    'key': '2.4 GHz: range migliore. 5 GHz: throughput migliore. Canali non sovrapposti = no co-channel interference.',
                },
                {
                    'type': 'teoria',
                    'title': 'Standard 802.11 — Evoluzione',
                    'points': [
                        '802.11b (1999): 2.4 GHz, 11 Mbps — prima diffusione di massa',
                        '802.11a (1999): 5 GHz, 54 Mbps — OFDM, meno diffuso per costo',
                        '802.11g (2003): 2.4 GHz, 54 Mbps, OFDM — retrocompatibile con b',
                        '802.11n (2009): 2.4+5 GHz, fino a 600 Mbps — introduce MIMO (più antenne)',
                        '802.11ac (2013): 5 GHz, fino a 3.5 Gbps — MU-MIMO, beamforming, canali 80/160 MHz',
                        '802.11ax / Wi-Fi 6 (2019): 2.4+5+6 GHz — OFDMA, TWT, BSS Coloring, fino a 9.6 Gbps',
                    ],
                    'key': 'MIMO = più flussi spaziali. OFDMA (ax) = divide canale tra più client simultaneamente.',
                },
            ],
        },
        {
            'section': 'BSS, ESS, CSMA/CA',
            'subtitle': 'Associazione e accesso al mezzo',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'BSS, SSID, BSSID ed ESS',
                    'points': [
                        'BSS (Basic Service Set): un AP + tutti i client ad esso associati',
                        'BSSID: indirizzo MAC dell\'AP — identifica univocamente un BSS',
                        'SSID: nome della rete wireless (stringa, fino a 32 caratteri)',
                        'ESSID: stesso SSID su più AP — forma l\'ESS (Extended Service Set)',
                        'DS (Distribution System): rete cablata che connette gli AP tra loro',
                        'Roaming: client si sposta da AP-1 a AP-2 mantenendo stesso SSID — reassociation',
                    ],
                    'key': 'BSS = una cella AP. ESS = più celle stessa rete. BSSID=MAC AP, SSID=nome rete.',
                },
                {
                    'type': 'teoria',
                    'title': 'CSMA/CA e Hidden Node Problem',
                    'points': [
                        'CSMA/CA: Carrier Sense Multiple Access / Collision Avoidance — ascolta prima di trasmettere',
                        'Wireless non può rilevare collisioni (half-duplex) — usa backoff casuale per evitarle',
                        'DCF (Distributed Coordination Function): DIFS + random backoff + ACK obbligatorio',
                        'Hidden node: Client A e Client B non si sentono tra loro — collisioni sull\'AP invisibili',
                        'Soluzione hidden node: RTS/CTS (Request to Send / Clear to Send) — riserva il mezzo',
                        'RTS/CTS: overhead aggiuntivo — usato solo quando hidden node è un problema reale',
                    ],
                    'key': 'CSMA/CA ≠ CSMA/CD. Wireless usa AVOIDANCE (evita) non DETECTION (rileva).',
                },
            ],
        },
    ],
    'trouble': [
        ('Throughput basso su 2.4 GHz', 'Co-channel interference? Canali sovrapposti (solo 1/6/11 non sovrapposti). Analisi spettro.'),
        ('Client non si associa', 'SSID nascosto? Tipo di sicurezza mismatch (WPA2 vs WPA3)? RSSI troppo basso (-80 dBm)?'),
        ('Roaming lento (sticky client)', 'Client rimane su AP distante. 802.11r (fast BSS transition) non supportato? BSS Max Idle?'),
        ('Hidden node problem', 'Attivare RTS/CTS: dot11 rts threshold 500 sull\'AP — soglia in byte'),
    ],
    'exam_tips': [
        'Canali non sovrapposti 2.4 GHz: 1, 6, 11 — tutti gli altri si sovrappongono',
        '802.11ax (Wi-Fi 6): OFDMA divide il canale tra più client — differenza chiave rispetto ad ac',
        'BSSID = MAC del radio dell\'AP. SSID = nome rete. Un AP può avere più SSID (diversi BSSID)',
        'CSMA/CA: listen before talk + random backoff + ACK obbligatorio — no collision detection',
        'Hidden node → soluzione = RTS/CTS. Aumenta overhead ma elimina collisioni nascoste.',
    ],
    'exam_qa': [
        ('Quanti canali non sovrapposti ha la banda 2.4 GHz?', '3 canali non sovrapposti: 1, 6 e 11. Tutti gli altri si sovrappongono e causano co-channel interference.'),
        ('Qual è la differenza principale tra 802.11ac e 802.11ax?', '802.11ax introduce OFDMA: divide il canale tra più client simultaneamente. 802.11ac usa OFDM con un client per volta per slot temporale.'),
    ],
    'summary': {
        'labels': ['2.4 vs 5 GHz', '802.11ax: OFDMA', 'CSMA/CA + RTS/CTS'],
        'bodies': [
            '2.4 GHz: range > penetrazione, 3 canali non sovrapposti. 5 GHz: throughput > canali. 6 GHz: Wi-Fi 6E.',
            'ax/Wi-Fi 6: OFDMA (multi-client), BSS Coloring, TWT. MIMO=flussi spaziali. MU-MIMO=multi-utente.',
            'CSMA/CA: DIFS+backoff+ACK. Hidden node→RTS/CTS. BSS=1 AP. ESS=più AP stesso SSID.',
        ],
    },
}

MODULES['MOD-24'] = {
    'title': 'Wireless Deployment & WLC',
    'area': 'AREA 9 — WIRELESS', 'hours': '2', 'codes': '2.1.d · 2.1.e · 2.1.f',
    'agenda': ['AP autonomo vs lightweight', 'CAPWAP: control e data', 'WLC: architettura e funzioni', 'FlexConnect e roaming', 'Exam Tips'],
    'topology': {
        'title': 'Architettura WLC — CAPWAP Tunnel',
        'nodes': [('WLC\n(controller)', 0.5, 0.2), ('AP-1\nLightweight', 0.2, 0.6), ('AP-2\nLightweight', 0.8, 0.6), ('Client', 0.2, 0.9)],
        'caption': 'CAPWAP control tunnel: WLC-AP (gestione) · CAPWAP data tunnel: traffico client centralizzato',
    },
    'sections': [
        {
            'section': 'AP Autonomo vs Lightweight',
            'subtitle': 'Modelli di deployment wireless',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Autonomous vs Lightweight AP',
                    'points': [
                        'AP Autonomo: configurazione completa locale (SSID, sicurezza, VLAN) — standalone',
                        'Gestione autonomo: individuale (SSH/console) — non scalabile oltre 10-15 AP',
                        'AP Lightweight: configurazione minima locale — delega la gestione al WLC via CAPWAP',
                        'WLC centralizza: RF management, sicurezza, roaming, QoS, policy per tutti gli AP',
                        'Conversione: AP Cisco può essere convertito tra autonomo e lightweight con boot da TFTP',
                        'Modello cloud: Cisco Meraki — AP lightweight con WLC cloud-based (Meraki Dashboard)',
                    ],
                    'key': 'Lightweight AP = "thin" — fa solo RF e trasmissione. WLC fa tutto il resto.',
                },
                {
                    'type': 'teoria',
                    'title': 'CAPWAP — Control and Provisioning of Wireless APs',
                    'points': [
                        'CAPWAP: protocollo standard (RFC 5415) per comunicazione WLC-AP',
                        'CAPWAP Control: UDP 5246 — configurazione, statistiche, autenticazione AP',
                        'CAPWAP Data: UDP 5247 — traffico client incapsulato e inviato al WLC',
                        'Discovery: AP invia CAPWAP Discovery Request broadcast → WLC risponde → join',
                        'Split MAC: funzioni MAC divise tra AP (real-time: beacon, probe, ack) e WLC (auth, assoc)',
                        'Sicurezza: tunnel DTLS opzionale per cifrare CAPWAP data — richiede licenza su alcuni modelli',
                    ],
                    'key': 'CAPWAP Control=UDP 5246 · Data=UDP 5247. Split MAC: AP gestisce solo il real-time.',
                },
            ],
        },
        {
            'section': 'WLC, FlexConnect e Roaming',
            'subtitle': 'Centralizzazione e mobilità avanzata',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'WLC — Funzioni Principali',
                    'points': [
                        'RF Management: seleziona canale e potenza per ogni AP automaticamente (RRM)',
                        'Client Authentication: 802.1X/EAP, PSK, WebAuth — centralizzato per tutti gli AP',
                        'Roaming: gestisce il passaggio di client tra AP — L2 roaming (stesso subnet) e L3 (anchor)',
                        'QoS: marking e queuing centralizzati — policy uniforme su tutta la rete wireless',
                        'WLAN Profile: SSID + sicurezza + VLAN mappata — applicato a gruppi di AP (AP Group)',
                        'High Availability: WLC HA pairs (1:1) con SSO (Stateful Switchover) — zero downtime client',
                    ],
                    'key': 'RRM: Radio Resource Management — WLC ottimizza RF automaticamente (canale+potenza).',
                },
                {
                    'type': 'teoria',
                    'title': 'FlexConnect — AP in Uffici Remoti',
                    'points': [
                        'Problema: WAN lenta/down → AP in modalità locale perde connessione WLC → client giù',
                        'FlexConnect: AP può commutare il traffico localmente anche senza WLC (standalone mode)',
                        'Local Switching: traffico client esce localmente sull\'AP — non transita nel CAPWAP data',
                        'Local Auth: AP autentica i client in locale (cache credenziali) se WLC non raggiungibile',
                        'Connected mode: AP è connesso a WLC — WLC gestisce normalmente',
                        'Standalone mode: WLC non raggiungibile — AP continua a servire client con config cache',
                    ],
                    'key': 'FlexConnect = AP resiliente su WAN. Local switching riduce traffico verso sede centrale.',
                },
            ],
        },
    ],
    'trouble': [
        ('AP non fa join al WLC', 'CAPWAP Discovery fallisce? AP in subnet diversa da WLC: configurare option 43 DHCP o DNS cisco-capwap-controller'),
        ('Client autentica ma non ottiene IP', 'VLAN mappata nel WLAN profile? DHCP server raggiungibile dall\'AP/WLC? Trunk sulla porta AP?'),
        ('Roaming lento', '802.11r (FT) abilitato sul WLAN profile? PMK caching attivo? Stessa mobility group?'),
        ('FlexConnect standalone: client disconnesso', 'WLAN non in FlexConnect local switching? AP non ha cached credentials per local auth?'),
    ],
    'exam_tips': [
        'CAPWAP Control=UDP 5246, Data=UDP 5247 — memorizzare le porte per l\'esame',
        'Split MAC: AP gestisce beacon/probe/ACK (real-time). WLC gestisce autenticazione/associazione.',
        'FlexConnect: local switching = traffico non passa per WLC. Standalone = WLC non raggiungibile.',
        'RRM: Radio Resource Management — WLC sceglie canale e potenza per ogni AP automaticamente',
        'AP Discovery: CAPWAP broadcast → DHCP option 43 → DNS cisco-capwap-controller → WLC configurato',
    ],
    'exam_qa': [
        ('Quali porte UDP usa CAPWAP?', 'CAPWAP Control = UDP 5246. CAPWAP Data = UDP 5247.'),
        ('Cosa succede a un AP FlexConnect quando la WAN verso il WLC va down?', 'Entra in standalone mode: continua a servire i client con configurazione in cache e autenticazione locale.'),
    ],
    'summary': {
        'labels': ['Lightweight + CAPWAP', 'WLC: RF + Auth', 'FlexConnect'],
        'bodies': [
            'AP Lightweight: solo RF real-time. CAPWAP: Control 5246/Data 5247. Split MAC divide funzioni.',
            'WLC: RRM (RF auto), 802.1X centralizzato, roaming L2/L3, QoS, HA SSO.',
            'FlexConnect: local switching (no CAPWAP data) + standalone mode (WAN down, client serviti).',
        ],
    },
}

MODULES['MOD-25'] = {
    'title': 'Wireless Security — WPA2, WPA3 e 802.1X',
    'area': 'AREA 9 — WIRELESS', 'hours': '1.5', 'codes': '2.1.g · 2.1.h',
    'agenda': ['WEP e la sua vulnerabilità', 'WPA, WPA2, WPA3', '802.1X EAP — autenticazione enterprise', 'PMF e protezione management frame', 'Exam Tips'],
    'topology': {
        'title': 'Wireless Security — Flusso 802.1X',
        'nodes': [('Client\n(Supplicant)', 0.1, 0.5), ('AP\n(Authenticator)', 0.4, 0.5), ('RADIUS\n(Auth Server)', 0.75, 0.5), ('WLC', 0.75, 0.2)],
        'caption': 'EAP over wireless (EAPOL) tra Client-AP · RADIUS tra AP/WLC e Auth Server',
    },
    'sections': [
        {
            'section': 'Evoluzione della Sicurezza Wireless',
            'subtitle': 'Da WEP a WPA3',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'WEP, WPA, WPA2 — Confronto',
                    'points': [
                        'WEP (1997): RC4 con IV 24-bit fisso → statisticamente forzabile in minuti. DEPRECATO.',
                        'WPA (2003): TKIP (Temporal Key Integrity Protocol) — IV esteso, MIC — patch su hardware WEP',
                        'WPA2 (2004): CCMP (AES-128 in CTR/CBC-MAC) — cifratura forte. Standard attuale.',
                        'WPA3 (2018): SAE (Simultaneous Authentication of Equals) — sostituisce PSK, forward secrecy',
                        'WPA3-Enterprise: 192-bit suite crittografica per ambienti high-security (governo, finanza)',
                        'CCMP = Counter Mode Cipher Block Chaining Message Authentication Code Protocol',
                    ],
                    'key': 'WPA2/CCMP = standard minimo. WPA3/SAE = no PSK catturabili, forward secrecy nativa.',
                },
                {
                    'type': 'teoria',
                    'title': '802.1X EAP — Autenticazione Enterprise',
                    'points': [
                        '802.1X: framework di autenticazione basato su porta — Supplicant / Authenticator / Auth Server',
                        'Supplicant: client wireless che chiede l\'accesso alla rete',
                        'Authenticator: AP o WLC — "portinaia" che blocca il traffico fino all\'autenticazione',
                        'Auth Server: RADIUS (porta 1812/1813) — verifica credenziali e autorizza',
                        'EAP (Extensible Authentication Protocol): trasporta le credenziali su 802.1X',
                        'Tipi EAP: PEAP (username/password in TLS tunnel), EAP-TLS (certificati mutual), EAP-FAST',
                    ],
                    'key': '802.1X = 3 ruoli: Supplicant (client), Authenticator (AP/WLC), Auth Server (RADIUS)',
                },
                {
                    'type': 'teoria',
                    'title': 'PMF e Protezione Frame Management',
                    'points': [
                        'Problema: frame management 802.11 (Deauth, Disassoc, Beacon) non cifrati in WPA2',
                        'Attacco: Deauth spoofato → client disconnesso da rete — denial of service wireless',
                        'PMF (Protected Management Frames): 802.11w — cifra e autentica frame management',
                        'PMF Optional: client con e senza PMF coesistono — compatibilità massima',
                        'PMF Required: solo client con PMF ammessi — sicurezza massima, no backward compat',
                        'WPA3 impone PMF obbligatorio — prerequisito per certificazione Wi-Fi Alliance WPA3',
                    ],
                    'key': 'PMF (802.11w) = cifra management frames. WPA3 lo richiede obbligatorio.',
                },
            ],
        },
    ],
    'trouble': [
        ('Client non si connette con WPA2-Enterprise', 'Certificato CA del RADIUS non fidato dal client? EAP type mismatch? RADIUS secret errata?'),
        ('Deauth attack in corso', 'PMF non abilitato? Attivare PMF Required o Optional. Analisi spettro per trovare sorgente.'),
        ('WPA3 SAE fallisce con client vecchi', 'Client non supportano SAE? Configurare WPA2/WPA3 in modalità transitional (mixed)'),
        ('RADIUS timeout', 'WLC raggiunge RADIUS su UDP 1812? Shared secret identica? NAS-IP configurato correttamente?'),
    ],
    'exam_tips': [
        'WEP = RC4+IV24bit → ROTTO. WPA = TKIP (patch). WPA2 = CCMP/AES (standard). WPA3 = SAE.',
        '802.1X: Supplicant=client, Authenticator=AP/WLC, Auth Server=RADIUS (UDP 1812/1813)',
        'PEAP: credenziali in tunnel TLS (no certificato client). EAP-TLS: certificati mutual (client+server).',
        'PMF (802.11w): protegge frame management da deauth spoofing. WPA3 lo richiede sempre.',
        'SAE (WPA3): no PSK trasmesso — Diffie-Hellman locale, forward secrecy — resistente a offline attack.',
    ],
    'exam_qa': [
        ('Qual è il protocollo di cifratura usato da WPA2?', 'CCMP (Counter Mode Cipher Block Chaining MAC Protocol) basato su AES-128.'),
        ('In 802.1X, quale dispositivo blocca il traffico fino al completamento dell\'autenticazione?', 'L\'Authenticator — di solito l\'AP o il WLC. Apre la "porta logica" solo dopo RADIUS Access-Accept.'),
    ],
    'summary': {
        'labels': ['WPA2/WPA3', '802.1X EAP', 'PMF 802.11w'],
        'bodies': [
            'WEP=rotto. WPA=TKIP(patch). WPA2=CCMP/AES. WPA3=SAE(forward secrecy, no PSK catturabile).',
            '802.1X: Supplicant(client)+Authenticator(AP)+RADIUS. PEAP=user/pass in TLS. EAP-TLS=cert mutual.',
            'PMF: cifra Deauth/Disassoc — blocca deauth spoofing. WPA3 impone PMF obbligatorio.',
        ],
    },
}

MODULES['MOD-26'] = {
    'title': 'QoS MQC — Classification, Marking e Queuing',
    'area': 'AREA 10 — QoS', 'hours': '2', 'codes': '1.5.a · 1.5.b',
    'agenda': ['Perché QoS: congestione e priorità', 'MQC: class-map e policy-map', 'DSCP e IP Precedence', 'CBWFQ e LLQ', 'Policing vs Shaping'],
    'topology': {
        'title': 'QoS MQC — Topologia Lab',
        'nodes': [('R1\n(trust edge)', 0.15, 0.5), ('R2\n(core)', 0.5, 0.5), ('R3\n(egress)', 0.85, 0.5)],
        'caption': 'Marking al bordo (R1) · Queuing al core (R2) · Shaping/Policing in uscita (R3)',
    },
    'sections': [
        {
            'section': 'MQC: Classificazione e Marking',
            'subtitle': 'Identificare e marcare il traffico',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Il Modello MQC (Modular QoS CLI)',
                    'points': [
                        'MQC: framework IOS per QoS — separa classificazione (class-map) da azione (policy-map)',
                        'class-map: identifica il traffico — match protocol, match dscp, match access-group',
                        'policy-map: definisce l\'azione per ogni classe — set, bandwidth, priority, police',
                        'service-policy: applica la policy-map a un\'interfaccia (input o output)',
                        'Hierarchical: policy-map può contenere un\'altra policy-map (parent/child) per traffic shaping',
                        'Regola: se il traffico non matcha nessuna class-map → va nella class class-default',
                    ],
                    'key': 'class-map = CHI. policy-map = COSA fare. service-policy = DOVE applicare.',
                },
                {
                    'type': 'teoria',
                    'title': 'DSCP — Differentiated Services Code Point',
                    'points': [
                        'DSCP: 6 bit nel byte ToS dell\'header IP — 64 valori possibili (0-63)',
                        'Backward compatible con IP Precedence (primi 3 bit del ToS)',
                        'DSCP 46 (EF): Expedited Forwarding — voce, video real-time — latenza minima garantita',
                        'DSCP 34 (AF41): Assured Forwarding 4 class 1 drop — video conferenza',
                        'DSCP 32 (CS4): Class Selector 4 — traffico a media priorità',
                        'DSCP 0 (BE): Best Effort — tutto il traffico non marcato — nessuna garanzia',
                    ],
                    'key': 'DSCP EF=46: voce. AF41=34: video. CS0/BE=0: default. Marcare al bordo della rete.',
                },
                {
                    'type': 'config',
                    'title': 'MQC — class-map e policy-map',
                    'device': 'R1',
                    'lines': [
                        'class-map match-any VOCE',
                        ' match dscp ef              ! DSCP 46',
                        ' match protocol rtp',
                        '!',
                        'class-map match-any DATI',
                        ' match dscp af41',
                        '!',
                        'policy-map PM-QOS',
                        ' class VOCE',
                        '  priority 512              ! LLQ — 512 kbps garantiti',
                        ' class DATI',
                        '  bandwidth 256             ! CBWFQ — 256 kbps garantiti',
                        ' class class-default',
                        '  fair-queue                ! WFQ per il resto',
                        '!',
                        'interface Serial0/0',
                        ' service-policy output PM-QOS',
                    ],
                    'hl': 9,
                },
            ],
        },
        {
            'section': 'CBWFQ, LLQ, Policing e Shaping',
            'subtitle': 'Queuing e gestione della banda',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'CBWFQ e LLQ',
                    'points': [
                        'FIFO: prima arrivato, primo servito — nessuna priorità. Default sulle interfacce veloci.',
                        'WFQ: Weighted Fair Queuing — code multiple con peso proporzionale alla priorità IP Prec',
                        'CBWFQ (Class-Based WFQ): garantisce banda minima per classe (bandwidth kbps)',
                        'Problema CBWFQ per voce: garantisce banda ma non la latenza — jitter può essere alto',
                        'LLQ (Low Latency Queuing): aggiunge una coda priority a CBWFQ — servita PRIMA di tutto',
                        'LLQ = CBWFQ + priority queue: perfetto per voce (latenza < 150ms, jitter < 30ms)',
                    ],
                    'key': 'LLQ = priority + CBWFQ. La coda priority è svuotata PRIMA di qualsiasi altra coda.',
                },
                {
                    'type': 'teoria',
                    'title': 'Policing vs Shaping',
                    'points': [
                        'Policing: controlla la banda — traffico in eccesso SCARTATO (drop) o remarked',
                        'Policing è bidirezionale — tipicamente applicato in input sul router di bordo ISP',
                        'Shaping: controlla la banda — traffico in eccesso BUFFERIZZATO (ritardato), non droppato',
                        'Shaping introduce latenza e jitter — non adatto per voce/video real-time',
                        'Policing: reattivo (drop immediato), shaping: proattivo (ritardo + smoothing)',
                        'Uso tipico: policing in ingresso (ISP), shaping in uscita (edge verso WAN lenta)',
                    ],
                    'key': 'Policing = drop. Shaping = delay. Policing per ingresso ISP. Shaping per WAN lenta.',
                },
            ],
        },
    ],
    'trouble': [
        ('Voce non funziona con QoS attivo', 'LLQ configurato? priority keyword invece di bandwidth per la classe VOCE?'),
        ('policy-map non applicata', 'service-policy mancante sull\'interfaccia? input vs output corretto?'),
        ('Traffico va in class-default', 'class-map non matcha? match protocol richiede NBAR (ip nbar protocol-discovery)?'),
        ('Shaping non funziona', 'Hierarchical policy-map richiesta? shape average vs shape peak?'),
    ],
    'exam_tips': [
        'MQC: class-map (classificazione) → policy-map (azione) → service-policy (applicazione interfaccia)',
        'DSCP EF=46 (voce), AF41=34 (video conferenza), BE=0 (default non marcato)',
        'LLQ = priority + CBWFQ. Voce DEVE usare priority (LLQ), non bandwidth (CBWFQ)',
        'Policing: drop del surplus. Shaping: buffer del surplus. Policing ha meno latenza.',
        'class class-default: raccoglie tutto il traffico che non matcha nessuna class-map',
    ],
    'exam_qa': [
        ('Qual è la differenza tra policing e shaping?', 'Policing droppa il traffico in eccesso. Shaping lo bufferizza e ritarda. Policing introduce meno latenza.'),
        ('Perché usare LLQ invece di CBWFQ per il traffico voce?', 'CBWFQ garantisce banda ma non latenza — la coda può crescere. LLQ aggiunge una coda priority svuotata prima di tutte le altre, garantendo latenza minima.'),
    ],
    'summary': {
        'labels': ['MQC: 3 componenti', 'DSCP EF/AF/BE', 'LLQ vs Shaping'],
        'bodies': [
            'class-map=classifica, policy-map=azione, service-policy=applica. class-default=catch-all.',
            'DSCP EF=46 voce, AF41=34 video, BE=0 default. 6 bit nel ToS. Marcare al bordo.',
            'LLQ=priority+CBWFQ(latenza garantita per voce). Policing=drop eccesso. Shaping=buffer eccesso.',
        ],
    },
}

MODULES['MOD-27'] = {
    'title': 'NAT, PAT & NTP',
    'area': 'AREA 11 — IP SERVICES', 'hours': '2', 'codes': '3.4.a · 3.4.b',
    'agenda': ['NAT statico e dinamico', 'PAT / NAPT — overload', 'Verifica e troubleshooting NAT', 'NTP: stratum e sincronizzazione', 'Exam Tips'],
    'topology': {
        'title': 'Topologia NAT/PAT — Lab MOD-27',
        'nodes': [('R1\nNAT Inside', 0.15, 0.5), ('R2\nNAT Router', 0.45, 0.5), ('ISP\nR3', 0.75, 0.5), ('PC-A\n192.168.1.x', 0.15, 0.8)],
        'caption': 'R1=inside network · R2=NAT router (inside/outside) · R3=ISP (pubblico)',
    },
    'sections': [
        {
            'section': 'NAT Statico, Dinamico e PAT',
            'subtitle': 'Traduzione degli indirizzi IP',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'NAT — Tipi e Funzionamento',
                    'points': [
                        'NAT (Network Address Translation): traduce indirizzi IP privati in pubblici',
                        'NAT Statico: mappatura 1:1 fissa — un IP privato = un IP pubblico permanente',
                        'Uso NAT statico: server interni raggiungibili dall\'esterno (web, mail, VPN)',
                        'NAT Dinamico: pool di IP pubblici — assegnazione temporanea dal pool',
                        'PAT / NAT Overload: più IP privati → un IP pubblico — differenziati dalla porta sorgente',
                        'PAT: il più usato in reti aziendali e SOHO — tutta la LAN esce con un IP pubblico',
                    ],
                    'key': 'PAT = molti:1. Usa porta sorgente come discriminante. Nessuna connessione entrante senza mapping.',
                },
                {
                    'type': 'config',
                    'title': 'PAT Configuration — R2',
                    'device': 'R2',
                    'lines': [
                        'interface Ethernet0/0',
                        ' ip address 192.168.1.254 255.255.255.0',
                        ' ip nat inside          ! inside = rete privata',
                        '!',
                        'interface Ethernet0/1',
                        ' ip address 203.0.113.1 255.255.255.0',
                        ' ip nat outside         ! outside = rete pubblica',
                        '!',
                        'ip nat inside source list 10 interface Ethernet0/1 overload',
                        '!',
                        'access-list 10 permit 192.168.1.0 0.0.0.255',
                    ],
                    'hl': 8,
                },
                {
                    'type': 'verifica',
                    'title': 'Verifica NAT — show ip nat translations',
                    'lines': [
                        'R2# show ip nat translations',
                        'Pro Inside global   Inside local    Outside local   Outside global',
                        'tcp 203.0.113.1:1025  192.168.1.10:1025  8.8.8.8:53   8.8.8.8:53',
                        'tcp 203.0.113.1:1026  192.168.1.11:1025  8.8.4.4:443  8.8.4.4:443',
                        '',
                        'R2# show ip nat statistics',
                        'Total active translations: 2 (0 static, 2 dynamic; 2 extended)',
                        'Outside interfaces: Ethernet0/1',
                        'Inside interfaces: Ethernet0/0',
                    ],
                    'hl': 1,
                },
            ],
        },
        {
            'section': 'NTP — Network Time Protocol',
            'subtitle': 'Sincronizzazione oraria nella rete',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'NTP — Stratum e Sincronizzazione',
                    'points': [
                        'NTP: protocollo di sincronizzazione oraria (RFC 5905) — UDP porta 123',
                        'Stratum: misura la distanza dal reference clock. Stratum 0 = orologio atomico/GPS',
                        'Stratum 1: server collegato direttamente al reference (stratum 0)',
                        'Stratum 2: sincronizzato da stratum 1 — tipico in reti enterprise',
                        'Regola: un device è sempre uno stratum in più del suo server NTP',
                        'Importanza: log coerenti, AAA accounting, certificati TLS, BGP/OSPF logging — tutto richiede tempo esatto',
                    ],
                    'key': 'Stratum = distanza dal reference clock. Stratum 0 = fonte. Router = stratum 3-5 tipicamente.',
                },
                {
                    'type': 'config',
                    'title': 'NTP Configuration — Client e Server',
                    'device': 'R1',
                    'lines': [
                        '! Configurazione NTP client',
                        'ntp server 192.168.1.1 prefer   ! server primario',
                        'ntp server 192.168.1.2          ! server secondario',
                        'ntp update-calendar              ! sync hardware clock',
                        '!',
                        '! Configurazione NTP master (router come stratum locale)',
                        'ntp master 3                    ! diventa stratum 3',
                        '!',
                        '! Verifica',
                        'R1# show ntp status',
                        'Clock is synchronized, stratum 4, reference is 192.168.1.1',
                        'R1# show ntp associations',
                        ' address         ref clock    st  when  poll reach',
                        '*~192.168.1.1    .GPS.         1    45    64  377',
                    ],
                    'hl': 10,
                },
            ],
        },
    ],
    'trouble': [
        ('NAT non traduce', 'ip nat inside/outside configurati? ACL permette il subnet? overload keyword presente per PAT?'),
        ('NAT statico non funziona in ingresso', 'ip nat inside source static <private> <public> presente? Outside interface corretta?'),
        ('NTP non sincronizza', 'UDP 123 bloccato da ACL? ntp server raggiungibile (ping)? show ntp associations mostra * ?'),
        ('Translazioni NAT si esauriscono', 'ip nat translation timeout troppo lungo? clear ip nat translation * per svuotare la tabella'),
    ],
    'exam_tips': [
        'NAT inside source: traduce sorgente dei pacchetti uscenti (inside→outside)',
        'PAT usa ip nat inside source list ... interface ... overload — porta differenzia i flussi',
        'show ip nat translations: vede la tabella. show ip nat statistics: contatori hit/miss',
        'NTP stratum: il router è stratum+1 rispetto al suo server. ntp master N configura il router come source.',
        'ip nat outside: interfaccia verso internet. ip nat inside: interfaccia verso la LAN privata.',
    ],
    'exam_qa': [
        ('Qual è la differenza tra NAT dinamico e PAT?', 'NAT dinamico: pool di IP pubblici, mappatura temporanea 1:1. PAT: un solo IP pubblico condiviso da tutti, la porta sorgente differenzia i flussi.'),
        ('Cosa indica il simbolo * in show ntp associations?', 'Il * indica il server NTP selezionato come riferimento primario (synchronized peer).'),
    ],
    'summary': {
        'labels': ['NAT Statico/PAT', 'NAT Verifica', 'NTP Stratum'],
        'bodies': [
            'Statico=1:1 fisso. Dinamico=pool. PAT=overload, porta differenzia. inside/outside obbligatori.',
            'show ip nat translations (tabella) · show ip nat statistics (contatori) · debug ip nat.',
            'NTP UDP 123. Stratum=distanza da ref. ntp server X = client. ntp master N = source locale.',
        ],
    },
}

MODULES['MOD-28'] = {
    'title': 'Multicast — PIM-SM e Auto-RP',
    'area': 'AREA 11 — IP SERVICES', 'hours': '2', 'codes': '3.3 · 3.4.d',
    'agenda': ['Perché il multicast: one-to-many', 'IGMP: group membership', 'PIM-SM e Rendezvous Point', 'Auto-RP e BSR', 'Verifica e Troubleshooting'],
    'topology': {
        'title': 'Topologia Multicast — PIM-SM Lab',
        'nodes': [('Source\n10.0.1.1', 0.1, 0.5), ('R1\nPIM-SM', 0.3, 0.5), ('R2\nRP', 0.5, 0.3), ('R3\nPIM-SM', 0.7, 0.5), ('Receiver\n10.0.3.x', 0.9, 0.5)],
        'caption': 'R2=Rendezvous Point (RP) · Shared Tree (*,G) → Shortest Path Tree (S,G)',
    },
    'sections': [
        {
            'section': 'IGMP e PIM-SM',
            'subtitle': 'Group membership e distribuzione multicast',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'IGMP — Internet Group Management Protocol',
                    'points': [
                        'IGMP: protocollo L3 tra host e router diretto per gestire le iscrizioni ai gruppi multicast',
                        'IGMPv1: host invia Join. Router invia Query periodica per verificare interesse.',
                        'IGMPv2: aggiunge Leave Group — più rapida la cancellazione dal gruppo',
                        'IGMPv3: Source Specific Multicast (SSM) — host specifica quale sorgente vuole',
                        'IGMP Snooping: switch legge i messaggi IGMP e limita flooding multicast alla porta corretta',
                        'Senza IGMP snooping: il multicast si comporta come broadcast a livello L2',
                    ],
                    'key': 'IGMPv2 = standard più diffuso. IGMP Snooping = essenziale per non inondare la LAN.',
                },
                {
                    'type': 'teoria',
                    'title': 'PIM-SM — Protocol Independent Multicast Sparse Mode',
                    'points': [
                        'PIM-SM: modalità sparsa — traffico inviato solo quando richiesto da un receiver',
                        'RP (Rendezvous Point): punto di incontro tra sorgenti e ricevitori — ruolo centrale',
                        'Shared Tree (*,G): percorso RP → receiver — usato inizialmente per ogni gruppo',
                        'SPT Switchover (S,G): dopo la prima consegna, il router può passare al percorso più breve',
                        'PIM Register: il DR sorgente incapsula in unicast il primo pacchetto multicast verso il RP',
                        'Join/Prune: router invia PIM Join verso RP, PIM Prune quando non ci sono più receiver',
                    ],
                    'key': 'PIM-SM: (*,G)=shared tree via RP. (S,G)=shortest path tree dopo SPT switchover.',
                },
                {
                    'type': 'config',
                    'title': 'PIM-SM + Auto-RP Configuration',
                    'device': 'R2 (RP)',
                    'lines': [
                        'ip multicast-routing              ! abilita multicast globale',
                        '!',
                        'interface Loopback0',
                        ' ip address 2.2.2.2 255.255.255.255',
                        ' ip pim sparse-mode',
                        '!',
                        'interface Ethernet0/0',
                        ' ip pim sparse-mode              ! PIM su tutte le interfacce',
                        '!',
                        '! Auto-RP: R2 si annuncia come RP per tutti i gruppi',
                        'ip pim send-rp-announce Loopback0 scope 10',
                        'ip pim send-rp-discovery Loopback0 scope 10',
                        '!',
                        '! Su tutti i router: abilitare Auto-RP listener',
                        'ip pim autorp listener',
                    ],
                    'hl': 10,
                },
            ],
        },
        {
            'section': 'Verifica Multicast',
            'subtitle': 'Comandi di verifica PIM e IGMP',
            'slides': [
                {
                    'type': 'verifica',
                    'title': 'Verifica PIM-SM — show ip mroute',
                    'lines': [
                        'R3# show ip mroute',
                        '(*, 239.1.1.1), 00:05:12/stopped, RP 2.2.2.2, flags: S',
                        '  Incoming interface: Ethernet0/0, RPF nbr 10.0.23.2',
                        '  Outgoing interface list:',
                        '    Ethernet0/1, Forward/Sparse, 00:05:10',
                        '',
                        '(10.0.1.1, 239.1.1.1), 00:00:05/00:02:54, flags: T',
                        '  Incoming interface: Ethernet0/0, RPF nbr 10.0.23.2',
                        '  Outgoing interface list:',
                        '    Ethernet0/1, Forward/Sparse, 00:00:05',
                    ],
                    'hl': 1,
                },
            ],
        },
    ],
    'trouble': [
        ('Receiver non riceve traffico multicast', 'IGMP join inviato? PIM sparse-mode su tutte le interfacce? RP raggiungibile?'),
        ('RP non trovato (Auto-RP)', 'ip pim autorp listener su tutti i router? 224.0.1.39/40 non bloccato da ACL?'),
        ('show ip mroute vuoto', 'ip multicast-routing abilitato globalmente? PIM abilitato sulle interfacce?'),
        ('SPT switchover non avviene', 'ip pim spt-threshold infinity blocca il switchover — rimuovere per SPT automatico'),
    ],
    'exam_tips': [
        'PIM-SM: sparse mode — traffico solo su richiesta. RP = punto di incontro sorgente/receiver.',
        '(*,G) = shared tree verso RP. (S,G) = shortest path tree dopo SPT switchover.',
        'Auto-RP usa 224.0.1.39 (Mapping Agent) e 224.0.1.40 (Discovery). BSR è standard IETF alternativo.',
        'IGMP: host → router. PIM: router → router. IGMP snooping: switch impara le porte multicast.',
        'ip multicast-routing: comando globale obbligatorio. ip pim sparse-mode: per interfaccia.',
    ],
    'exam_qa': [
        ('Qual è la differenza tra (*,G) e (S,G) nella tabella multicast?', '(*,G) = shared tree: qualsiasi sorgente, percorso via RP. (S,G) = source-specific tree: percorso ottimale diretto dalla sorgente S al receiver.'),
        ('Cosa fa IGMP Snooping sugli switch?', 'Legge i messaggi IGMP e apprende su quale porta ci sono receiver. Invia il traffico multicast solo alle porte con receiver, invece di inondare tutta la VLAN.'),
    ],
    'summary': {
        'labels': ['IGMP + Snooping', 'PIM-SM: RP e Tree', 'Auto-RP'],
        'bodies': [
            'IGMPv2: Join/Leave tra host e router. IGMPv3: SSM. Snooping: switch impara porte multicast.',
            'PIM-SM: (*,G)=shared tree via RP. (S,G)=SPT switchover. PIM Register: primo pacchetto al RP.',
            'Auto-RP: 224.0.1.39 announce, 224.0.1.40 discovery. ip pim autorp listener su tutti.',
        ],
    },
}

MODULES['MOD-29'] = {
    'title': 'Network Assurance — SNMP, NetFlow e SPAN',
    'area': 'AREA 11 — IP SERVICES', 'hours': '2', 'codes': '4.1 · 4.2 · 4.3 · 4.4 · 4.5 · 4.6',
    'agenda': ['SNMP v2c/v3: monitoring', 'Syslog: livelli e best practice', 'NetFlow: traffic accounting', 'IP SLA: synthetic testing', 'SPAN/RSPAN: traffic mirroring'],
    'topology': {
        'title': 'Network Assurance — Architettura Monitoring',
        'nodes': [('Router/Switch\nDevice', 0.2, 0.5), ('Syslog\nServer', 0.6, 0.2), ('SNMP NMS\n(PRTG/Zabbix)', 0.6, 0.5), ('NetFlow\nCollector', 0.6, 0.8)],
        'caption': 'SNMP: polling/trap · Syslog: messaggi log · NetFlow: statistiche flusso',
    },
    'sections': [
        {
            'section': 'SNMP e Syslog',
            'subtitle': 'Monitoring e logging centralizzato',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'SNMP — Simple Network Management Protocol',
                    'points': [
                        'SNMP: protocollo di monitoring e gestione dispositivi di rete (UDP 161/162)',
                        'MIB (Management Information Base): database gerarchico di oggetti monitorabili (OID)',
                        'SNMPv2c: community string (testo in chiaro) — semplice ma non sicuro',
                        'SNMPv3: autenticazione (MD5/SHA) + cifratura (DES/AES) — standard enterprise',
                        'Polling: NMS chiede (GET) al device — pull model, a intervalli regolari',
                        'Trap/Inform: device notifica proattivamente l\'NMS (SNMP Trap=UDP 162) — push model',
                    ],
                    'key': 'SNMPv3 = unico sicuro. Trap: device → NMS, UDP 162. Inform: trap con ACK.',
                },
                {
                    'type': 'teoria',
                    'title': 'Syslog — Livelli di Severità',
                    'points': [
                        'Syslog: standard di logging (RFC 5424) — messaggi di testo a un server centralizzato',
                        'Porta: UDP 514 (default) o TCP 514 (reliable syslog)',
                        'Livello 0 (Emergency): sistema inutilizzabile — critico',
                        'Livello 1 (Alert): azione immediata richiesta',
                        'Livello 2 (Critical): condizione critica hardware/software',
                        'Livello 3 (Error): errori non fatali · 4 (Warning) · 5 (Notice) · 6 (Informational) · 7 (Debug)',
                    ],
                    'key': 'Mnemonica livelli: Every Awesome Cisco Engineer Will Need Icecream Daily (0-7)',
                },
            ],
        },
        {
            'section': 'NetFlow, IP SLA e SPAN',
            'subtitle': 'Traffic accounting e synthetic monitoring',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'NetFlow e IP SLA',
                    'points': [
                        'NetFlow: raccoglie statistiche sui flussi IP (src/dst IP, porta, protocollo, byte, pacchetti)',
                        'Flusso NetFlow: identificato da 7 tuple — sorgente, destinazione, protocollo, porte, ToS, interfaccia',
                        'Export: router invia record al collector via UDP 2055 (NetFlow v5/v9) o UDP 4739 (IPFIX)',
                        'Uso: capacity planning, security (anomaly detection), billing, troubleshooting',
                        'IP SLA: invia traffico sintetico per misurare latenza, jitter, packet loss su path specifici',
                        'IP SLA ICMP Echo: ping sintetico periodico. IP SLA UDP Jitter: misura jitter VoIP.',
                    ],
                    'key': 'NetFlow = osserva traffico reale. IP SLA = genera traffico sintetico per misurare il path.',
                },
                {
                    'type': 'teoria',
                    'title': 'SPAN e RSPAN — Traffic Mirroring',
                    'points': [
                        'SPAN (Switched Port Analyzer): copia il traffico di una porta verso una porta monitor',
                        'Uso SPAN: collegare un analizzatore (Wireshark, IDS) per vedere il traffico in tempo reale',
                        'SPAN locale: porta sorgente e porta destinazione sullo stesso switch',
                        'RSPAN (Remote SPAN): traffico mirrorato trasportato su una VLAN dedicata tra switch remoti',
                        'ERSPAN: incapsula il traffico in GRE — trasporto su rete IP (non solo L2)',
                        'Limitazione: SPAN non cattura errori L1/L2 (FCS error) — per quello serve tap hardware',
                    ],
                    'key': 'SPAN: mirror locale. RSPAN: mirror cross-switch via VLAN dedicata. ERSPAN: mirror via IP/GRE.',
                },
            ],
        },
    ],
    'trouble': [
        ('SNMP trap non arriva al NMS', 'snmp-server host configurato? UDP 162 non bloccato da ACL? Community string corretta?'),
        ('NetFlow record non arrivano al collector', 'ip flow export destination configurato? UDP 2055 aperto? ip flow ingress sull\'interfaccia?'),
        ('IP SLA non parte', 'ip sla schedule configurato? ip sla N: operation definita correttamente?'),
        ('SPAN non cattura traffico', 'monitor session definita? Porta destinazione non in trunk/access normale?'),
    ],
    'exam_tips': [
        'SNMPv3: autenticazione + cifratura. SNMPv2c: solo community string (no encryption). v1: obsoleto.',
        'Syslog livelli 0-7: Emergency, Alert, Critical, Error, Warning, Notice, Info, Debug',
        'NetFlow: 7-tuple flusso. Export UDP 2055 (v5/v9) o 4739 (IPFIX). Collector separato dal router.',
        'IP SLA: traffico sintetico per misurare latenza/jitter. Usato con tracking per floating static route.',
        'SPAN: locale. RSPAN: cross-switch via VLAN. ERSPAN: over IP/GRE. SPAN porta destinazione = solo receive.',
    ],
    'exam_qa': [
        ('Differenza tra SNMP Trap e SNMP Inform?', 'Trap: UDP non affidabile, nessun ACK. Inform: il NMS invia un ACK — più affidabile ma richiede più risorse sul device.'),
        ('Cosa misura IP SLA UDP Jitter?', 'Misura latenza, jitter (variazione latenza) e packet loss su un path UDP — ideale per monitorare la qualità di path VoIP.'),
    ],
    'summary': {
        'labels': ['SNMP v2c/v3', 'NetFlow + IPFIX', 'SPAN/RSPAN/IP SLA'],
        'bodies': [
            'SNMPv3: auth+encrypt. Trap=UDP162 (no ACK). Inform=ACK. Polling=GET. MIB=OID database.',
            'NetFlow: 7-tuple, export UDP 2055. IPFIX=RFC NetFlow. Uso: security, capacity, billing.',
            'IP SLA: synthetic test (ping/jitter). SPAN: mirror locale. RSPAN: cross-switch. ERSPAN: GRE/IP.',
        ],
    },
}

MODULES['MOD-30'] = {
    'title': 'Device Security & AAA',
    'area': 'AREA 12 — SECURITY', 'hours': '2', 'codes': '5.1.a · 5.1.b',
    'agenda': ['AAA: Authentication, Authorization, Accounting', 'RADIUS vs TACACS+', 'Hardening device: SSH, banners, privilege levels', 'Local fallback e best practice', 'Exam Tips'],
    'topology': {
        'title': 'AAA Architecture — RADIUS/TACACS+',
        'nodes': [('Admin\n(Client)', 0.1, 0.5), ('Router/Switch\n(NAS)', 0.4, 0.5), ('RADIUS\nServer', 0.75, 0.3), ('TACACS+\nServer', 0.75, 0.7)],
        'caption': 'NAS = Network Access Server · RADIUS UDP 1812/1813 · TACACS+ TCP 49',
    },
    'sections': [
        {
            'section': 'AAA — Tre Funzioni',
            'subtitle': 'Authentication, Authorization, Accounting',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'AAA: Concetti Fondamentali',
                    'points': [
                        'Authentication: verifica l\'identità — chi sei? Username/password, certificato, OTP',
                        'Authorization: determina i permessi — cosa puoi fare? Comandi, livelli di privilegio',
                        'Accounting: registra le azioni — cosa hai fatto? Log di sessione, comandi eseguiti',
                        'NAS (Network Access Server): il device (router/switch) che delega AAA al server',
                        'RADIUS: UDP 1812 (auth) / 1813 (accounting) — combina autenticazione e autorizzazione',
                        'TACACS+: TCP 49 — separa le tre funzioni AAA, cifra l\'intero payload',
                    ],
                    'key': 'TACACS+: TCP, cifra tutto, separa AAA — preferito per gestione dispositivi di rete.',
                },
                {
                    'type': 'teoria',
                    'title': 'RADIUS vs TACACS+ — Confronto',
                    'points': [
                        'RADIUS: UDP 1812/1813. Cifra solo la password. Combina authn e authz in un\'unica risposta.',
                        'RADIUS: ideale per accesso utenti end (VPN, wireless 802.1X, accesso internet)',
                        'TACACS+: TCP 49. Cifra l\'intero pacchetto. Separa Authentication / Authorization / Accounting.',
                        'TACACS+: ideale per gestione dispositivi (CLI router/switch) — command authorization',
                        'Command authorization: TACACS+ può autorizzare/bloccare singoli comandi IOS',
                        'Cisco raccomanda TACACS+ per device management, RADIUS per network access',
                    ],
                    'key': 'RADIUS=UDP/network access. TACACS+=TCP/device management. TACACS+ cifra tutto.',
                },
                {
                    'type': 'config',
                    'title': 'AAA Configuration — TACACS+',
                    'device': 'R1',
                    'lines': [
                        'aaa new-model                     ! abilita AAA globalmente',
                        '!',
                        'tacacs server TACACS-PRIMARY',
                        ' address ipv4 192.168.122.10',
                        ' key CiscoTACACS123',
                        '!',
                        'aaa group server tacacs+ TACACS-GRP',
                        ' server name TACACS-PRIMARY',
                        '!',
                        'aaa authentication login default group TACACS-GRP local',
                        'aaa authorization exec default group TACACS-GRP local',
                        'aaa accounting exec default start-stop group TACACS-GRP',
                        '!',
                        'line vty 0 4',
                        ' login authentication default',
                        ' transport input ssh',
                    ],
                    'hl': 9,
                },
            ],
        },
        {
            'section': 'Hardening e SSH',
            'subtitle': 'Best practice per la sicurezza del device',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Device Hardening — Best Practice',
                    'points': [
                        'Disabilitare servizi non necessari: no ip http server, no cdp run (dove non necessario)',
                        'Privilege levels: 0 (no access), 1 (user EXEC), 15 (privileged) — livelli custom 2-14',
                        'Enable secret vs enable password: secret usa MD5/SCRYPT, password è testo in chiaro',
                        'service password-encryption: offusca (non cifra) password in chiaro nel config',
                        'Banner motd/login/exec: messaggio legale obbligatorio prima del login',
                        'SSH: versione 2 obbligatoria. ip ssh version 2, crypto key generate rsa modulus 2048',
                    ],
                    'key': 'enable secret > enable password. SSH v2 obbligatorio. Banner = avviso legale.',
                },
            ],
        },
    ],
    'trouble': [
        ('AAA fallisce: nessun accesso al device', 'Sempre configurare fallback locale: aaa authentication login default group TACACS local'),
        ('TACACS+ non risponde', 'TCP 49 raggiungibile? Shared key identica? test aaa group TACACS username password legacy'),
        ('Command authorization troppo restrittiva', 'TACACS+ nega comandi legittimi? Verificare shell:priv-lvl nel profilo TACACS'),
        ('SSH non funziona', 'crypto key generate rsa modulus 2048? ip ssh version 2? ip domain-name configurato?'),
    ],
    'exam_tips': [
        'RADIUS: UDP 1812 (auth) / 1813 (acct). TACACS+: TCP 49. TACACS+ cifra tutto il pacchetto.',
        'aaa new-model: comando obbligatorio per abilitare AAA. Senza: il device usa autenticazione locale.',
        'TACACS+: separa Authentication / Authorization / Accounting — RADIUS li combina.',
        'Fallback locale: sempre aggiungere local alla fine della lista AAA — evita lockout totale.',
        'enable secret: MD5 hash (o SCRYPT). enable password: testo in chiaro — NON usare mai in produzione.',
    ],
    'exam_qa': [
        ('Perché TACACS+ è preferito a RADIUS per la gestione dei dispositivi di rete?', 'TACACS+ separa Authentication/Authorization/Accounting, cifra l\'intero pacchetto e supporta command authorization. RADIUS combina authn e authz e cifra solo la password.'),
        ('Cosa succede se il server TACACS+ è irraggiungibile e non c\'è fallback locale?', 'Nessun accesso al device. Per questo il fallback local è obbligatorio: aaa authentication login default group TACACS+ local'),
    ],
    'summary': {
        'labels': ['AAA: 3 funzioni', 'RADIUS vs TACACS+', 'Hardening + SSH'],
        'bodies': [
            'Authentication=chi sei, Authorization=cosa puoi, Accounting=cosa hai fatto. aaa new-model abilita.',
            'RADIUS UDP 1812/1813 (network access). TACACS+ TCP 49 (device mgmt, cifra tutto, separa AAA).',
            'enable secret > password. SSH v2 + RSA 2048. Banner legale. Fallback local sempre configurato.',
        ],
    },
}

MODULES['MOD-31'] = {
    'title': 'ACL & Control Plane Policing',
    'area': 'AREA 12 — SECURITY', 'hours': '2', 'codes': '5.2.a · 5.2.b',
    'agenda': ['ACL standard ed estese', 'ACL named e best practice', 'Posizionamento ACL', 'CoPP: protezione del control plane', 'Verifica e Troubleshooting'],
    'topology': {
        'title': 'ACL e CoPP — Topologia Lab',
        'nodes': [('PC-A\n10.1.1.10', 0.1, 0.5), ('R1\nACL in/out', 0.4, 0.5), ('R2\nCoPP', 0.7, 0.5), ('Server\n10.2.1.1', 0.95, 0.5)],
        'caption': 'ACL filtrano traffico transit · CoPP protegge CPU del router da flooding',
    },
    'sections': [
        {
            'section': 'ACL Standard ed Estese',
            'subtitle': 'Filtraggio del traffico IP',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'ACL — Concetti Fondamentali',
                    'points': [
                        'ACL: lista ordinata di regole permit/deny — valutate in ordine, prima match = decisione',
                        'Implicit deny all: ogni ACL termina con "deny any" implicito — non visibile nel config',
                        'ACL Standard (1-99, 1300-1999): filtra solo sorgente IP — granularità bassa',
                        'ACL Extended (100-199, 2000-2699): filtra src/dst IP, protocollo, porta — granularità alta',
                        'Named ACL: nome descrittivo, modificabile senza riscrivere tutta la lista',
                        'Wildcard mask: inverso della subnet mask — 0=bit deve matchare, 1=bit ignorato',
                    ],
                    'key': 'Posizionamento: ACL Extended → vicino alla sorgente. ACL Standard → vicino alla destinazione.',
                },
                {
                    'type': 'config',
                    'title': 'ACL Extended Named — Configurazione',
                    'device': 'R1',
                    'lines': [
                        'ip access-list extended FILTER-WEB',
                        ' 10 permit tcp 10.1.1.0 0.0.0.255 host 10.2.1.1 eq 443',
                        ' 20 permit tcp 10.1.1.0 0.0.0.255 host 10.2.1.1 eq 80',
                        ' 30 deny   tcp any host 10.2.1.1 eq 22  ! blocca SSH da fuori',
                        ' 40 permit ip any any               ! permetti il resto',
                        '!',
                        'interface Ethernet0/0',
                        ' ip access-group FILTER-WEB in    ! applica in ingresso',
                        '!',
                        '! Verifica',
                        'R1# show ip access-lists FILTER-WEB',
                        'Extended IP access list FILTER-WEB',
                        '    10 permit tcp 10.1.1.0/24 host 10.2.1.1 eq 443 (45 matches)',
                        '    20 permit tcp 10.1.1.0/24 host 10.2.1.1 eq 80 (12 matches)',
                    ],
                    'hl': 7,
                },
            ],
        },
        {
            'section': 'CoPP — Control Plane Policing',
            'subtitle': 'Protezione della CPU del router',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'CoPP — Perché e Come',
                    'points': [
                        'Problema: traffico destinato alla CPU del router (ICMP, BGP, OSPF, SSH, SNMP) può saturarla',
                        'Attacco: flooding di pacchetti verso la CPU → CPU overload → router non gestisce più il routing',
                        'CoPP: applica rate-limiting al traffico diretto al control plane usando MQC',
                        'Protegge: OSPF Hello, BGP keepalive, SSH, SNMP, ICMP — tutto il traffico verso CPU',
                        'Policy: priorità alta per protocolli critici (OSPF, BGP), rate-limit per il resto',
                        'Implementazione: service-policy input applicata all\'interfaccia "control-plane"',
                    ],
                    'key': 'CoPP = MQC applicato al control plane. Protegge la CPU da flooding intenzionale o accidentale.',
                },
                {
                    'type': 'config',
                    'title': 'CoPP — Configurazione Base',
                    'device': 'R2',
                    'lines': [
                        'class-map match-any ROUTING-PROTOCOLS',
                        ' match protocol ospf',
                        ' match protocol bgp',
                        '!',
                        'class-map match-any MANAGEMENT',
                        ' match protocol ssh',
                        ' match protocol snmp',
                        '!',
                        'policy-map COPP-POLICY',
                        ' class ROUTING-PROTOCOLS',
                        '  police rate 64000 bps       ! 64 kbps per OSPF/BGP',
                        ' class MANAGEMENT',
                        '  police rate 32000 bps       ! 32 kbps per SSH/SNMP',
                        ' class class-default',
                        '  police rate 8000 bps        ! 8 kbps per il resto',
                        '!',
                        'control-plane',
                        ' service-policy input COPP-POLICY',
                    ],
                    'hl': 16,
                },
            ],
        },
    ],
    'trouble': [
        ('ACL non filtra correttamente', 'Ordine delle regole? Prima match vince. show ip access-lists mostra i contatori hit.'),
        ('ACL blocca traffico legittimo', 'Implicit deny at end? Aggiungere permit ip any any se necessario. Verificare wildcard mask.'),
        ('CoPP droppa OSPF/BGP', 'Rate limit troppo basso? Aumentare il rate per ROUTING-PROTOCOLS. show policy-map control-plane'),
        ('ACL non applicata', 'ip access-group applicato sull\'interfaccia? in vs out corretto rispetto al flusso del traffico?'),
    ],
    'exam_tips': [
        'ACL Standard: numeri 1-99, 1300-1999 — filtra solo IP sorgente',
        'ACL Extended: numeri 100-199, 2000-2699 — filtra src+dst+protocollo+porta',
        'Posizionamento: Extended vicino alla sorgente, Standard vicino alla destinazione',
        'Implicit deny all: ogni ACL ha "deny any" finale implicito — non visibile nel running config',
        'CoPP: MQC + control-plane { service-policy input }. Protegge CPU da DoS.',
    ],
    'exam_qa': [
        ('Perché le ACL Extended devono essere posizionate vicino alla sorgente?', 'Perché filtrano all\'inizio del percorso, evitando che il traffico da bloccare attraversi inutilmente la rete prima di essere scartato.'),
        ('Cosa protegge CoPP e come si implementa?', 'CoPP protegge la CPU del router da flooding di traffico control plane. Si implementa con MQC applicato all\'interfaccia logica "control-plane" con service-policy input.'),
    ],
    'summary': {
        'labels': ['ACL Standard/Extended', 'Posizionamento ACL', 'CoPP'],
        'bodies': [
            'Standard=src IP (1-99). Extended=src+dst+proto+porta (100-199). Named=modificabile. Implicit deny.',
            'Extended: vicino alla sorgente (più efficiente). Standard: vicino alla destinazione.',
            'CoPP: MQC + control-plane service-policy input. Rate-limit traffico verso CPU. Protegge da DoS.',
        ],
    },
}

MODULES['MOD-32'] = {
    'title': 'EEM & Python Automation',
    'area': 'AREA 13 — AUTOMATION', 'hours': '2', 'codes': '6.1 · 6.2 · 6.6',
    'agenda': ['EEM: event-driven automation IOS', 'EEM Applet: sintassi e trigger', 'Python GuestShell su IOS-XE', 'Netmiko: SSH automation', 'Exam Tips'],
    'topology': {
        'title': 'EEM e Python — Architettura Automation',
        'nodes': [('R1\nEEM + GuestShell', 0.25, 0.5), ('Script Python\n(Netmiko)', 0.65, 0.3), ('Syslog/SNMP\nTrap', 0.65, 0.7)],
        'caption': 'EEM: automation locale al device · Python/Netmiko: automation da host esterno',
    },
    'sections': [
        {
            'section': 'EEM — Embedded Event Manager',
            'subtitle': 'Automation nativa in IOS',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'EEM — Concetti e Architettura',
                    'points': [
                        'EEM: framework di automation integrato in IOS — reagisce a eventi della rete',
                        'Evento: syslog pattern, CLI command, interfaccia down, timer, OIR, SNMP threshold',
                        'Azione: esegue CLI, invia syslog, invia trap SNMP, esegue script TCL/Python',
                        'Applet: policy EEM semplice definita in IOS — nessun file esterno richiesto',
                        'Script TCL/Python: policy complessa — file .tcl/.py caricato nel flash o da TFTP',
                        'Uso tipico: backup automatico config, remediation automatica, alerting avanzato',
                    ],
                    'key': 'EEM = automazione locale al router, zero dipendenze esterne. Reagisce agli eventi IOS.',
                },
                {
                    'type': 'config',
                    'title': 'EEM Applet — Backup Automatico',
                    'device': 'R1',
                    'lines': [
                        'event manager applet BACKUP-CONFIG',
                        ' event syslog pattern "Configured from"',
                        ' action 1.0 cli command "enable"',
                        ' action 2.0 cli command "copy running-config tftp:"',
                        ' action 3.0 cli command "192.168.122.1"',
                        ' action 4.0 cli command "r1-backup.cfg"',
                        ' action 5.0 syslog msg "Config backed up via EEM"',
                        '!',
                        '! Verifica',
                        'R1# show event manager policy registered',
                        'No.  Class     Type    Event Type     Trap  Time Registered',
                        '1    applet    system  syslog         Off   Mon May 18 10:00:00',
                        '  Name: BACKUP-CONFIG',
                    ],
                    'hl': 1,
                },
            ],
        },
        {
            'section': 'Python e Netmiko',
            'subtitle': 'SSH automation da host esterno',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Python per Network Automation',
                    'points': [
                        'Python: linguaggio di scripting dominante per network automation',
                        'Librerie chiave: Netmiko (SSH), NAPALM (multi-vendor), Nornir (framework), Requests (REST)',
                        'Netmiko: ConnectHandler — astrae la connessione SSH a diversi OS (IOS, IOS-XE, NX-OS)',
                        'send_command(): invia comando e restituisce output come stringa',
                        'send_config_set(): invia lista di comandi di configurazione',
                        'GuestShell (IOS-XE): container Linux integrato nel router — esegue Python localmente',
                    ],
                    'key': 'Netmiko: SSH multi-vendor in 5 righe. GuestShell: Python sul router senza server esterno.',
                },
                {
                    'type': 'config',
                    'title': 'Netmiko — Script SSH Base',
                    'device': 'Python',
                    'lines': [
                        'from netmiko import ConnectHandler',
                        '',
                        'device = {',
                        '    "device_type": "cisco_ios",',
                        '    "host": "10.0.0.1",',
                        '    "username": "admin",',
                        '    "password": "cisco",',
                        '}',
                        '',
                        'conn = ConnectHandler(**device)',
                        'output = conn.send_command("show ip route")',
                        'print(output)',
                        'conn.disconnect()',
                    ],
                    'hl': 10,
                },
            ],
        },
    ],
    'trouble': [
        ('EEM applet non si attiva', 'Pattern syslog corretto? show logging per vedere i messaggi. event manager run APPLET-NAME per test.'),
        ('Netmiko: timeout SSH', 'SSH raggiungibile? timeout= e banner_timeout= parametri ConnectHandler. Chiave RSA generata?'),
        ('GuestShell non parte', 'iox abilitato? virtual-service install name guestshell? guestshell enable?'),
        ('EEM action cli non esegue', 'Credenziali enable configurate? action 1.0 cli command "enable" prima dei comandi privilegiati?'),
    ],
    'exam_tips': [
        'EEM: event-driven automation locale in IOS. Applet = policy semplice. Script TCL/Python = complessa.',
        'EEM event types: syslog, timer, interface, cli, oir, snmp — variegato',
        'Netmiko: ConnectHandler + send_command() + send_config_set() — le 3 funzioni base',
        'GuestShell: container Linux su IOS-XE. guestshell enable → Python 3 disponibile localmente.',
        'NAPALM: multi-vendor API (get_facts, get_interfaces, get_bgp_neighbors) — abstraction layer.',
    ],
    'exam_qa': [
        ('Cosa fa EEM in Cisco IOS?', 'EEM (Embedded Event Manager) è un framework di automation che reagisce a eventi del sistema (syslog, timer, interfacce) ed esegue azioni (CLI, syslog, trap) senza dipendenze esterne.'),
        ('Qual è la differenza tra send_command e send_config_set in Netmiko?', 'send_command() esegue un comando in EXEC mode e restituisce l\'output. send_config_set() entra in config mode, invia una lista di comandi di configurazione e poi esce.'),
    ],
    'summary': {
        'labels': ['EEM Applet', 'Netmiko SSH', 'GuestShell'],
        'bodies': [
            'EEM: event-driven locale. Applet: policy IOS (syslog→action cli). Script: TCL/Python da flash.',
            'Netmiko: ConnectHandler(device_type, host, user, pass). send_command/send_config_set.',
            'GuestShell: container Linux su IOS-XE. iox + guestshell enable. Python 3 locale sul router.',
        ],
    },
}

MODULES['MOD-33'] = {
    'title': 'Netmiko & Nornir — Framework Automation',
    'area': 'AREA 13 — AUTOMATION', 'hours': '2', 'codes': '6.2 · 6.3',
    'agenda': ['Netmiko avanzato: TextFSM e parsing', 'Nornir: inventory e task', 'Esecuzione parallela con Nornir', 'Output e report', 'Exam Tips'],
    'topology': {
        'title': 'Nornir — Architettura Framework',
        'nodes': [('Inventory\n(hosts.yaml)', 0.15, 0.5), ('Nornir Core\n(runner)', 0.45, 0.5), ('Task\n(funzione Python)', 0.75, 0.3), ('Result\n(dict)', 0.75, 0.7)],
        'caption': 'Nornir: inventory + runner + tasks. Esecuzione parallela multi-device nativa.',
    },
    'sections': [
        {
            'section': 'Netmiko Avanzato e TextFSM',
            'subtitle': 'Parsing strutturato dell\'output CLI',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'TextFSM — Parsing Output CLI',
                    'points': [
                        'Problema: output CLI è testo non strutturato — difficile da elaborare programmaticamente',
                        'TextFSM: libreria Google che converte output CLI in liste di dizionari strutturati',
                        'Template: definisce pattern regex per estrarre campi dall\'output',
                        'Netmiko + TextFSM: use_textfsm=True in send_command() → ritorna lista di dizionari',
                        'NTC-Templates: repository pubblico di template TextFSM per comandi Cisco/Arista/Juniper',
                        'Alternativa: Genie Parser (NAPALM) — template prebuilt per IOS, IOS-XE, NX-OS',
                    ],
                    'key': 'TextFSM: testo non strutturato → dizionari Python. use_textfsm=True in Netmiko.',
                },
                {
                    'type': 'config',
                    'title': 'Netmiko + TextFSM — Esempio',
                    'device': 'Python',
                    'lines': [
                        'from netmiko import ConnectHandler',
                        '',
                        'device = {"device_type": "cisco_ios",',
                        '          "host": "10.0.0.1", "username": "admin", "password": "cisco"}',
                        '',
                        'conn = ConnectHandler(**device)',
                        '',
                        '# use_textfsm=True: output parsato automaticamente',
                        'routes = conn.send_command("show ip route",',
                        '                           use_textfsm=True)',
                        '',
                        '# routes è ora una lista di dizionari',
                        'for route in routes:',
                        '    print(f"{route[\'network\']}/{route[\'mask\']} via {route[\'nexthop_ip\']}")',
                    ],
                    'hl': 8,
                },
            ],
        },
        {
            'section': 'Nornir Framework',
            'subtitle': 'Automation multi-device parallela',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Nornir — Architettura e Vantaggi',
                    'points': [
                        'Nornir: framework Python puro per network automation — no DSL, tutto Python',
                        'Inventory: definisce host (IP, credenziali, OS, gruppi) in YAML o SimpleInventory',
                        'Runner: gestisce l\'esecuzione — SerialRunner (sequenziale) o ThreadedRunner (parallelo)',
                        'Task: funzione Python che esegue operazioni su un singolo host',
                        'Result: oggetto che raccoglie l\'output di ogni task per ogni host',
                        'Vantaggioso su Netmiko puro: parallelismo nativo, inventory centralizzato, report strutturato',
                    ],
                    'key': 'Nornir: inventory YAML + tasks Python + ThreadedRunner = automation parallela su 100+ device.',
                },
                {
                    'type': 'config',
                    'title': 'Nornir — Script Base',
                    'device': 'Python',
                    'lines': [
                        'from nornir import InitNornir',
                        'from nornir_netmiko.tasks import netmiko_send_command',
                        'from nornir_utils.plugins.functions import print_result',
                        '',
                        '# Inizializza Nornir con config file',
                        'nr = InitNornir(config_file="config.yaml")',
                        '',
                        '# Filtra host per gruppo',
                        'routers = nr.filter(groups=["routers"])',
                        '',
                        '# Esegui task su tutti i router in parallelo',
                        'result = routers.run(',
                        '    task=netmiko_send_command,',
                        '    command_string="show ip ospf neighbor")',
                        '',
                        'print_result(result)',
                    ],
                    'hl': 11,
                },
            ],
        },
    ],
    'trouble': [
        ('TextFSM non parsa correttamente', 'Template corretto per il device OS? ntc-templates installato? Testare con parse_output() standalone.'),
        ('Nornir: host non raggiungibile', 'Hostname/IP corretto in hosts.yaml? Credenziali? SSH abilitato? Timeout configurato?'),
        ('Nornir: task fallisce su alcuni host', 'Errori in AggregatedResult? print_result() mostra dettagli. on_failed=True per continuare nonostante errori.'),
        ('Parallelismo causa problemi', 'num_workers in config.yaml? Ridurre per debugging. SerialRunner per esecuzione sequenziale.'),
    ],
    'exam_tips': [
        'TextFSM: converte output CLI in struttura dati Python. use_textfsm=True in Netmiko.',
        'Nornir: inventory YAML, tasks Python, risultati strutturati — alternativa a Ansible più flessibile',
        'ThreadedRunner: esecuzione parallela multi-device — molto più veloce di loop sequenziale',
        'NTC-Templates: repository di template TextFSM per parsing multi-vendor',
        'Nornir non ha DSL (come Ansible) — tutto Python puro: più flessibile, curva apprendimento più alta',
    ],
    'exam_qa': [
        ('Che vantaggio offre Nornir rispetto a un semplice loop Netmiko?', 'Nornir fornisce inventory centralizzato, esecuzione parallela nativa (ThreadedRunner), gestione strutturata degli errori e result aggregati — senza scrivere il boilerplate di threading manualmente.'),
        ('Cosa fa TextFSM e perché è utile?', 'TextFSM converte output CLI non strutturato in liste di dizionari Python usando template regex. Questo permette di accedere ai dati programmaticamente (es. route["nexthop_ip"]) invece di fare parsing manuale delle stringhe.'),
    ],
    'summary': {
        'labels': ['TextFSM Parsing', 'Nornir Framework', 'ThreadedRunner'],
        'bodies': [
            'TextFSM: template regex → dizionari Python. use_textfsm=True in Netmiko. NTC-Templates: prebuilt.',
            'Nornir: inventory YAML + tasks Python + AggregatedResult. InitNornir + nr.run(task=...).',
            'ThreadedRunner: N device in parallelo. SerialRunner: sequenziale per debug. on_failed per resilienza.',
        ],
    },
}

MODULES['MOD-34'] = {
    'title': 'Ansible & Git per Network Automation',
    'area': 'AREA 13 — AUTOMATION', 'hours': '2', 'codes': '6.7',
    'agenda': ['Git: version control per le configurazioni', 'Ansible: architettura e concetti', 'Playbook IOS: ios_command e ios_config', 'Inventory e variables', 'Exam Tips'],
    'topology': {
        'title': 'Ansible — Control Node e Managed Nodes',
        'nodes': [('Control Node\n(Ansible)', 0.2, 0.5), ('R1\nManaged', 0.55, 0.3), ('R2\nManaged', 0.55, 0.5), ('R3\nManaged', 0.55, 0.7)],
        'caption': 'Ansible: agentless — SSH puro dal control node. No software sui managed node.',
    },
    'sections': [
        {
            'section': 'Git — Version Control',
            'subtitle': 'Controllo versione per configurazioni di rete',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Git — Concetti Fondamentali',
                    'points': [
                        'Git: sistema di version control distribuito — traccia le modifiche ai file nel tempo',
                        'Repository: cartella con history completa di tutte le modifiche',
                        'Commit: snapshot del repository in un momento — immutabile, identificato da hash SHA1',
                        'Branch: linea di sviluppo parallela — main/master per produzione, feature branch per sviluppo',
                        'Merge / Pull Request: integra le modifiche di un branch in un altro',
                        'Uso network: versioning delle configurazioni, rollback, peer review (pull request)',
                    ],
                    'key': 'Git per il networking: ogni cambio config = un commit. Rollback = git revert. Review = PR.',
                },
            ],
        },
        {
            'section': 'Ansible per IOS',
            'subtitle': 'Automation dichiarativa agentless',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'Ansible — Architettura e Concetti',
                    'points': [
                        'Ansible: framework di automation dichiarativa — descrive lo stato desiderato',
                        'Agentless: nessun software da installare sui managed node — usa SSH o API',
                        'Inventory: file che lista i managed node con IP, credenziali, variabili, gruppi',
                        'Playbook: YAML — descrive la sequenza di task da eseguire sui host',
                        'Task: chiama un modulo (ios_command, ios_config, uri) con parametri',
                        'Idempotente: eseguire un playbook più volte produce lo stesso risultato — safe to re-run',
                    ],
                    'key': 'Ansible = dichiarativo + idempotente + agentless. YAML playbook: legibile anche dai non-dev.',
                },
                {
                    'type': 'config',
                    'title': 'Playbook Ansible — ios_config',
                    'device': 'YAML',
                    'lines': [
                        '---',
                        '- name: Configure OSPF on routers',
                        '  hosts: routers',
                        '  gather_facts: false',
                        '  tasks:',
                        '    - name: Enable OSPF process',
                        '      cisco.ios.ios_config:',
                        '        lines:',
                        '          - router ospf 100',
                        '          - router-id {{ router_id }}',
                        '          - network 10.0.0.0 0.0.0.255 area 0',
                        '        save_when: changed',
                        '',
                        '    - name: Verify OSPF neighbors',
                        '      cisco.ios.ios_command:',
                        '        commands: show ip ospf neighbor',
                        '      register: ospf_output',
                    ],
                    'hl': 6,
                },
            ],
        },
    ],
    'trouble': [
        ('Ansible: SSH connection refused', 'SSH abilitato sul router? Credenziali corrette? ansible_connection=network_cli configurato?'),
        ('ios_config non applica', 'Modulo cisco.ios installato? ansible-galaxy collection install cisco.ios'),
        ('Playbook non idempotente', 'ios_config verifica lo stato prima di applicare. Check mode: ansible-playbook --check'),
        ('Git: merge conflict', 'Due branch modificano stessa riga. git diff per vedere il conflitto, editare manualmente, git add + commit.'),
    ],
    'exam_tips': [
        'Ansible: agentless (solo SSH). Dichiarativo (stato desiderato). Idempotente (safe to re-run).',
        'ios_command: esegue comandi show (read-only). ios_config: applica configurazione (write).',
        'Git: commit=snapshot, branch=linea parallela, merge=integrazione, revert=rollback sicuro.',
        'Ansible inventory: hosts.ini o YAML. Variabili: group_vars/ e host_vars/ per override.',
        'YAML: sensibile all\'indentazione. 2 spazi standard. Linting: yamllint prima di eseguire.',
    ],
    'exam_qa': [
        ('Cosa significa che Ansible è idempotente?', 'Eseguire lo stesso playbook più volte produce sempre lo stesso risultato finale. Se la configurazione desiderata è già applicata, Ansible non fa nulla (no change).'),
        ('Differenza tra ios_command e ios_config in Ansible?', 'ios_command: esegue comandi show/exec e restituisce output. ios_config: entra in configuration mode e applica comandi — verifica la differenza prima di applicare.'),
    ],
    'summary': {
        'labels': ['Git Version Control', 'Ansible Agentless', 'ios_config + ios_command'],
        'bodies': [
            'Git: commit=snapshot, branch=parallela, merge PR=review. Config network in Git = rollback facile.',
            'Ansible: YAML playbook, agentless SSH, dichiarativo, idempotente. No software sui router.',
            'ios_config: applica config (idempotente). ios_command: legge output. register: salva risultato.',
        ],
    },
}

MODULES['MOD-35'] = {
    'title': 'API & RESTCONF — Automation Programmatica',
    'area': 'AREA 13 — AUTOMATION', 'hours': '1.5', 'codes': '6.4 · 6.5',
    'agenda': ['REST API: concetti e metodi HTTP', 'RESTCONF — RFC 8040', 'YANG data models', 'Chiamate RESTCONF con Python requests', 'Exam Tips'],
    'topology': {
        'title': 'RESTCONF — Client e Device',
        'nodes': [('Client\nPython/curl', 0.15, 0.5), ('HTTPS\nRESTCONF', 0.45, 0.5), ('Router\nIOS-XE', 0.75, 0.5)],
        'caption': 'RESTCONF: HTTPS + JSON/XML · YANG: modello dati · RFC 8040 — no CLI, solo API',
    },
    'sections': [
        {
            'section': 'REST API e RESTCONF',
            'subtitle': 'Automazione via HTTP/HTTPS',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'REST API — Concetti Fondamentali',
                    'points': [
                        'REST (Representational State Transfer): architettura per API web — stateless, basata su HTTP',
                        'Metodi HTTP: GET (legge), POST (crea), PUT (sostituisce), PATCH (modifica), DELETE (cancella)',
                        'Response codes: 200 OK, 201 Created, 204 No Content, 400 Bad Request, 404 Not Found',
                        'Formato dati: JSON (leggero, human-readable) o XML (verboso, schema validation)',
                        'Headers: Content-Type: application/json, Accept: application/json, Authorization: Basic',
                        'Stateless: ogni richiesta contiene tutte le informazioni — no sessione server-side',
                    ],
                    'key': 'REST: GET=leggi, POST=crea, PUT=sostituisci, PATCH=modifica, DELETE=cancella.',
                },
                {
                    'type': 'teoria',
                    'title': 'RESTCONF e YANG',
                    'points': [
                        'RESTCONF (RFC 8040): standard IETF per configurare device di rete via REST su HTTPS',
                        'Base path: https://{device}/restconf/data/{yang-module}:{container}/{leaf}',
                        'YANG: linguaggio di modellazione dati (RFC 6020) — definisce struttura della configurazione',
                        'Namespace YANG: Cisco-IOS-XE-native (config Cisco) o ietf-interfaces (standard)',
                        'Content-Type RESTCONF: application/yang-data+json o application/yang-data+xml',
                        'NOTA: RESTCONF non disponibile su IOU — richiede IOS-XE (CSR1000v, Cat8000v)',
                    ],
                    'key': 'RESTCONF = REST + YANG. Configurazione strutturata via API HTTPS. No CLI.',
                },
                {
                    'type': 'config',
                    'title': 'RESTCONF — GET con Python requests',
                    'device': 'Python',
                    'lines': [
                        'import requests',
                        'import json',
                        '',
                        'url = "https://10.0.0.1/restconf/data/ietf-interfaces:interfaces"',
                        'headers = {',
                        '    "Accept": "application/yang-data+json",',
                        '    "Content-Type": "application/yang-data+json"',
                        '}',
                        '',
                        'response = requests.get(',
                        '    url, headers=headers,',
                        '    auth=("admin", "cisco"),',
                        '    verify=False)  # no SSL verify in lab',
                        '',
                        'data = response.json()',
                        'print(json.dumps(data, indent=2))',
                    ],
                    'hl': 9,
                },
            ],
        },
        {
            'section': 'NETCONF e Confronto Protocolli',
            'subtitle': 'Panoramica protocolli di automation',
            'slides': [
                {
                    'type': 'teoria',
                    'title': 'NETCONF vs RESTCONF vs SNMP',
                    'points': [
                        'SNMP: monitoring (pull + trap). MIB/OID non intuitivi. Standard ma limitato per config.',
                        'NETCONF (RFC 6241): SSH + XML + YANG — config/operazionale. Transazioni atomiche.',
                        'RESTCONF (RFC 8040): HTTP + JSON/XML + YANG — subset di NETCONF più semplice',
                        'RESTCONF: ideale per sviluppatori (stesse librerie HTTP del web). No lock/commit.',
                        'NETCONF: più potente (rollback, candidate config, lock) — preferito per automazione enterprise',
                        'gRPC/gNMI: streaming telemetria — sostituto di SNMP per monitoring ad alta frequenza',
                    ],
                    'key': 'ENCOR exam: RESTCONF=HTTP+YANG. NETCONF=SSH+XML+YANG. SNMP=monitoring legacy.',
                },
            ],
        },
    ],
    'trouble': [
        ('RESTCONF: 404 Not Found', 'URL corretto? YANG path valido? ip http secure-server abilitato su IOS-XE?'),
        ('RESTCONF: 401 Unauthorized', 'Credenziali Base64 corrette? aaa authorization exec default local configurato?'),
        ('RESTCONF non disponibile su IOU', 'IOU non supporta RESTCONF — usare CSR1000v o IOS-XE su GNS3/EVE-NG per i test'),
        ('SSL verify error', 'Certificato self-signed su lab: verify=False in requests.get() per lab. In produzione: importare CA.'),
    ],
    'exam_tips': [
        'RESTCONF: RFC 8040. Base path: /restconf/data/. Usa HTTPS (443). Dati YANG in JSON o XML.',
        'HTTP methods: GET=read, POST=create, PUT=replace, PATCH=update, DELETE=remove',
        'YANG: linguaggio di modellazione dati. Cisco-IOS-XE-native: config Cisco. ietf-interfaces: standard.',
        'Content-Type RESTCONF: application/yang-data+json (non application/json classico)',
        'NETCONF: SSH+XML, transazioni atomiche, candidate config. RESTCONF: HTTP+JSON, più semplice.',
    ],
    'exam_qa': [
        ('Qual è la differenza tra NETCONF e RESTCONF?', 'NETCONF (RFC 6241): SSH + XML, supporta transazioni atomiche, candidate config e lock. RESTCONF (RFC 8040): HTTP + JSON/XML, subset di NETCONF, più semplice ma senza transazioni.'),
        ('Quale metodo HTTP si usa per modificare parzialmente una configurazione RESTCONF?', 'PATCH — modifica solo i campi specificati senza sostituire l\'intera risorsa. PUT sostituisce completamente la risorsa.'),
    ],
    'summary': {
        'labels': ['REST: metodi HTTP', 'RESTCONF + YANG', 'NETCONF vs RESTCONF'],
        'bodies': [
            'GET/POST/PUT/PATCH/DELETE. 200=OK, 201=Created, 404=Not Found, 401=Unauthorized.',
            'RESTCONF: HTTPS+YANG. Path: /restconf/data/module:container. JSON: yang-data+json.',
            'NETCONF: SSH+XML, atomico, candidate. RESTCONF: HTTP+JSON, semplice. gNMI: telemetria.',
        ],
    },
}

if __name__ == "__main__":
    main()
